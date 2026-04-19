from __future__ import annotations

import base64
import difflib
import hashlib
import io
import ipaddress
import json
import math
import re
import shutil
import subprocess
import tempfile
import uuid
import zipfile
from collections import defaultdict
from collections.abc import Iterable
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

import yaml

from sediment.diagnostics import DiagnosticLogger
from sediment.insights import (
    build_cluster_key,
    cluster_state,
    compute_demand_score,
    compute_maturity_score,
    detect_intent,
    detect_query_language,
    infer_insight_kind,
    infer_recommended_action,
    insight_frontmatter,
    insight_title_from_cluster,
    is_ready_for_materialization,
    normalize_query_for_kb,
    normalize_subject,
    parse_insight,
    render_insight_markdown,
    slugify_filename,
    validate_insight_content,
)
from sediment.kb import (
    audit_kb,
    extract_wikilinks,
    index_config,
    inventory,
    resolve_kb_document_path,
    split_sections,
    split_frontmatter,
    validate_entry,
    validate_index,
)
from sediment.platform_store import PlatformStore, utc_now

LOGGER = DiagnosticLogger("platform_services")

FORMAL_STATUSES = {"fact", "inferred", "disputed"}
ALLOWED_MIME_TYPES = {
    "text/plain",
    "text/markdown",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/zip",
    "application/x-zip-compressed",
}
ZIP_MIME_TYPES = {"application/zip", "application/x-zip-compressed"}
SUPPORTED_UPLOAD_SUFFIXES = {".txt", ".md", ".docx", ".pptx"}
GRAPH_EVENT_BASE_ENERGY = {
    "ingest_created": 1.0,
    "ingest_updated": 0.72,
    "ask_reinforced": 0.78,
    "proposal_materialized": 0.92,
    "insight_promoted": 1.05,
    "insight_merged": 0.88,
}
GRAPH_SCENE_BUDGETS = {
    "home": {"events": 14, "nodes": 40, "edges": 80, "scene_mode": "portal-story"},
    "full": {"events": 24, "nodes": 70, "edges": 120, "scene_mode": "portal-immersive"},
    "admin": {"events": 28, "nodes": 90, "edges": 140, "scene_mode": "admin-governance"},
}
GRAPH_EVENT_PRIORITY = {
    "insight_promoted": 6,
    "insight_merged": 5,
    "proposal_materialized": 4,
    "ingest_created": 3,
    "ingest_updated": 2,
    "ask_reinforced": 1,
}
SUBMISSION_ANALYSIS_SCHEMA = json.dumps(
    {
        "type": "object",
        "additionalProperties": False,
        "properties": {
            "summary": {"type": "string"},
            "recommended_title": {"type": "string"},
            "recommended_type": {
                "type": "string",
                "enum": ["concept", "lesson", "feedback"],
            },
            "duplicate_risk": {
                "type": "string",
                "enum": ["low", "medium", "high"],
            },
            "committer_action": {
                "type": "string",
                "enum": ["manual_review", "ingest", "merge_or_link"],
            },
            "committer_note": {"type": "string"},
            "related_entries": {
                "type": "array",
                "items": {
                    "type": "object",
                    "additionalProperties": False,
                    "properties": {
                        "name": {"type": "string"},
                        "reason": {"type": "string"},
                    },
                    "required": ["name", "reason"],
                },
            },
        },
        "required": [
            "summary",
            "recommended_title",
            "recommended_type",
            "duplicate_risk",
            "committer_action",
            "committer_note",
            "related_entries",
        ],
    },
    ensure_ascii=False,
)


def infer_mime_type(filename: str) -> str | None:
    lower = str(filename or "").strip().lower()
    if lower.endswith(".md"):
        return "text/markdown"
    if lower.endswith(".txt"):
        return "text/plain"
    if lower.endswith(".docx"):
        return "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    if lower.endswith(".pptx"):
        return "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    if lower.endswith(".zip"):
        return "application/zip"
    return None


def ensure_platform_state(
    *,
    store: PlatformStore,
    state_dir: str | Path,
    uploads_dir: str | Path,
    workspaces_dir: str | Path,
) -> None:
    Path(state_dir).mkdir(parents=True, exist_ok=True)
    Path(uploads_dir).mkdir(parents=True, exist_ok=True)
    Path(workspaces_dir).mkdir(parents=True, exist_ok=True)
    store.init()


def build_submission_hash(*parts: str) -> str:
    digest = hashlib.sha256()
    for part in parts:
        digest.update(part.encode("utf-8"))
        digest.update(b"\n---\n")
    return digest.hexdigest()


def normalize_name(raw_value: str, *, fallback: str) -> str:
    cleaned = re.sub(r"\s+", " ", raw_value).strip()
    return cleaned or fallback


def parse_trusted_proxy_cidrs(raw_value: str | None) -> tuple[ipaddress._BaseNetwork, ...]:
    if not raw_value:
        return ()
    networks: list[ipaddress._BaseNetwork] = []
    for item in raw_value.split(","):
        token = item.strip()
        if not token:
            continue
        networks.append(ipaddress.ip_network(token, strict=False))
    return tuple(networks)


def detect_submitter_ip(
    headers: dict[str, str] | None,
    client_host: str | None,
    *,
    trust_proxy_headers: bool = False,
    trusted_proxy_cidrs: Iterable[ipaddress._BaseNetwork] = (),
) -> str:
    headers = headers or {}
    normalized_client = normalize_ip(client_host)
    if not trust_proxy_headers:
        return normalized_client

    trusted_networks = tuple(trusted_proxy_cidrs)
    if not trusted_networks:
        return normalized_client
    if not ip_in_networks(normalized_client, trusted_networks):
        return normalized_client

    forwarded = normalize_ip(headers.get("x-forwarded-for", "").split(",")[0].strip())
    real_ip = normalize_ip(headers.get("x-real-ip", "").strip())
    return forwarded or real_ip or normalized_client


def normalize_ip(raw_value: str | None) -> str:
    value = (raw_value or "").strip()
    if not value:
        return "unknown"
    candidate = value
    if value.startswith("[") and "]" in value:
        candidate = value[1 : value.index("]")]
    elif value.count(":") == 1 and "." in value:
        host, _, port = value.rpartition(":")
        if port.isdigit():
            candidate = host
    try:
        return str(ipaddress.ip_address(candidate))
    except ValueError:
        return value


def ip_in_networks(
    client_ip: str,
    networks: Iterable[ipaddress._BaseNetwork],
) -> bool:
    try:
        address = ipaddress.ip_address(client_ip)
    except ValueError:
        return False
    return any(address in network for network in networks)


def extract_upload_text(file_path: str | Path, mime_type: str) -> str:
    path = Path(file_path)
    if mime_type == "text/plain":
        return path.read_text(encoding="utf-8", errors="replace")
    if mime_type == "text/markdown":
        return path.read_text(encoding="utf-8", errors="replace")
    if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        from docx import Document

        doc = Document(path)
        paragraphs = [item.text.strip() for item in doc.paragraphs if item.text.strip()]
        return "\n\n".join(paragraphs).strip()
    if mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        from pptx import Presentation

        presentation = Presentation(path)
        chunks: list[str] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_lines: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = str(shape.text).strip()
                    if text:
                        slide_lines.append(text)
            if slide_lines:
                chunks.append(f"Slide {slide_index}\n" + "\n".join(slide_lines))
        return "\n\n".join(chunks).strip()
    raise ValueError(f"Unsupported mime type: {mime_type}")


def analyze_text_submission(
    *,
    kb_path: str | Path,
    title: str,
    content: str,
    submission_type: str,
) -> dict[str, Any]:
    cleaned_title = normalize_name(title, fallback="未命名提交")
    cleaned_content = content.strip()
    shortlist = search_kb(kb_path, f"{cleaned_title}\n{cleaned_content}", limit=6)
    fallback = _fallback_submission_analysis(
        title=cleaned_title,
        submission_type=submission_type,
        shortlist=shortlist,
    )
    LOGGER.info(
        "submission.analysis.start",
        "Starting submission analysis.",
        details={
            "title": cleaned_title,
            "submission_type": submission_type,
            "shortlist_count": len(shortlist),
        },
    )
    try:
        from sediment.llm_cli import build_cli_command, collect_output, parse_json_object
        from sediment.settings import load_settings

        settings = load_settings()
        timeout_seconds = min(
            max(15, int(settings["agent"]["exec_timeout_seconds"])),
            90,
        )
        prompt = _build_submission_analysis_prompt(
            title=cleaned_title,
            content=cleaned_content,
            submission_type=submission_type,
            shortlist=shortlist,
        )
        with tempfile.TemporaryDirectory(prefix="sediment-submit-analysis-") as temp_dir:
            temp_root = Path(temp_dir)
            prompt_file = temp_root / "prompt.txt"
            payload_file = temp_root / "payload.json"
            prompt_file.write_text(prompt, encoding="utf-8")
            payload_file.write_text(
                json.dumps(
                    {
                        "title": cleaned_title,
                        "submission_type": submission_type,
                        "content": cleaned_content,
                        "shortlist": shortlist,
                    },
                    ensure_ascii=False,
                    indent=2,
                ),
                encoding="utf-8",
            )
            invocation = build_cli_command(
                settings,
                prompt,
                prompt_file=prompt_file,
                payload_file=payload_file,
                cwd=Path(settings["workspace_root"]),
                extra_args=["--json-schema", SUBMISSION_ANALYSIS_SCHEMA],
            )
            result = subprocess.run(
                invocation.command,
                input=invocation.stdin_data,
                capture_output=True,
                text=True,
                timeout=timeout_seconds,
                check=False,
                cwd=str(Path(settings["workspace_root"])),
            )
        if result.returncode != 0:
            detail = (
                result.stderr.strip()
                or result.stdout.strip()
                or f"exit code {result.returncode}"
            )
            fallback["status"] = "degraded"
            fallback["warnings"] = [f"Agent analysis failed: {detail}"]
            LOGGER.warning(
                "submission.analysis.degraded",
                "Submission analysis degraded because the agent CLI failed.",
                details={
                    "returncode": result.returncode,
                    "stderr": result.stderr,
                    "stdout": result.stdout,
                },
            )
            return fallback
        raw_output = collect_output(invocation, stdout=result.stdout, stderr=result.stderr)
        parsed = parse_json_object(raw_output)
        normalized = _normalize_submission_analysis(
            parsed,
            title=cleaned_title,
            submission_type=submission_type,
            shortlist=shortlist,
        )
        normalized["status"] = "ok"
        LOGGER.info(
            "submission.analysis.completed",
            "Submission analysis completed.",
            details={
                "recommended_type": normalized.get("recommended_type"),
                "committer_action": normalized.get("committer_action"),
            },
        )
        return normalized
    except Exception as exc:  # noqa: BLE001
        fallback["status"] = "degraded"
        fallback["warnings"] = [f"Agent analysis unavailable: {exc}"]
        LOGGER.error(
            "submission.analysis.unavailable",
            "Submission analysis fell back because the agent runtime was unavailable.",
            error=exc,
            details={
                "title": cleaned_title,
                "submission_type": submission_type,
            },
        )
        return fallback


def prepare_document_submission(
    *,
    filename: str,
    mime_type: str,
    file_bytes: bytes,
    uploads: list[dict[str, Any]] | None,
    max_upload_bytes: int,
) -> dict[str, Any]:
    normalized_uploads, skipped = _expand_document_uploads(
        filename=filename,
        mime_type=mime_type,
        file_bytes=file_bytes,
        uploads=uploads,
        max_upload_bytes=max_upload_bytes,
    )
    if not normalized_uploads:
        raise ValueError("could not find a supported document to ingest")
    LOGGER.info(
        "document.prepare.start",
        "Preparing document submission payload.",
        details={
            "filename": filename,
            "mime_type": mime_type,
            "upload_count": len(normalized_uploads),
            "skipped_count": len(skipped),
        },
    )

    extracted_parts: list[str] = []
    with tempfile.TemporaryDirectory(prefix="sediment-upload-extract-") as temp_dir:
        temp_root = Path(temp_dir)
        for index, upload in enumerate(normalized_uploads, start=1):
            temp_path = temp_root / f"{index}_{sanitize_filename(Path(upload['filename']).name)}"
            temp_path.write_bytes(upload["file_bytes"])
            text = extract_upload_text(temp_path, upload["mime_type"]).strip()
            if not text:
                continue
            extracted_parts.append(f"## {upload['relative_path']}\n\n{text}")

    if not extracted_parts:
        raise ValueError("could not extract text from uploaded document")

    source_label = _bundle_label(filename=filename, uploads=normalized_uploads)
    notes_parts: list[str] = []
    if skipped:
        notes_parts.append("Skipped unsupported files: " + ", ".join(skipped[:8]))
        LOGGER.warning(
            "document.prepare.skipped_uploads",
            "Skipped unsupported uploads while preparing document submission.",
            details={"skipped": skipped[:8], "skipped_count": len(skipped)},
        )
    notes = " | ".join(notes_parts) if notes_parts else None

    original_is_archive = mime_type in ZIP_MIME_TYPES
    if uploads or len(normalized_uploads) > 1 or original_is_archive:
        if original_is_archive and not uploads:
            stored_bytes = file_bytes
            stored_filename = sanitize_filename(filename or "upload.zip")
        else:
            stored_bytes = _zip_document_bundle(normalized_uploads)
            stored_filename = sanitize_filename(source_label)
            if not stored_filename.lower().endswith(".zip"):
                stored_filename = f"{stored_filename}.zip"
        stored_mime_type = "application/zip"
    else:
        only = normalized_uploads[0]
        stored_bytes = only["file_bytes"]
        stored_filename = sanitize_filename(only["filename"])
        stored_mime_type = only["mime_type"]

    return {
        "title": source_label,
        "filename": stored_filename,
        "mime_type": stored_mime_type,
        "file_bytes": stored_bytes,
        "extracted_text": "\n\n".join(extracted_parts).strip(),
        "notes": notes,
        "file_count": len(normalized_uploads),
    }


def prepare_document_staging_upload(
    *,
    filename: str,
    mime_type: str,
    file_bytes: bytes,
    uploads: list[dict[str, Any]] | None,
    max_upload_bytes: int,
) -> dict[str, Any]:
    normalized_uploads, skipped = _expand_document_uploads(
        filename=filename,
        mime_type=mime_type,
        file_bytes=file_bytes,
        uploads=uploads,
        max_upload_bytes=max_upload_bytes,
    )
    if not normalized_uploads:
        raise ValueError("could not find a supported document to stage")

    source_label = _bundle_label(filename=filename, uploads=normalized_uploads)
    notes_parts: list[str] = []
    if skipped:
        notes_parts.append("Skipped unsupported files: " + ", ".join(skipped[:8]))
    original_is_archive = mime_type in ZIP_MIME_TYPES
    if uploads or len(normalized_uploads) > 1 or original_is_archive:
        if original_is_archive and not uploads:
            stored_bytes = file_bytes
            stored_filename = sanitize_filename(filename or "upload.zip")
        else:
            stored_bytes = _zip_document_bundle(normalized_uploads)
            stored_filename = sanitize_filename(source_label)
            if not stored_filename.lower().endswith(".zip"):
                stored_filename = f"{stored_filename}.zip"
        stored_mime_type = "application/zip"
    else:
        only = normalized_uploads[0]
        stored_bytes = only["file_bytes"]
        stored_filename = sanitize_filename(only["filename"])
        stored_mime_type = only["mime_type"]

    return {
        "title": source_label,
        "filename": stored_filename,
        "mime_type": stored_mime_type,
        "file_bytes": stored_bytes,
        "notes": " | ".join(notes_parts) if notes_parts else None,
        "file_count": len(normalized_uploads),
    }


def _build_submission_analysis_prompt(
    *,
    title: str,
    content: str,
    submission_type: str,
    shortlist: list[dict[str, Any]],
) -> str:
    shortlist_payload = [
        {
            "name": item["name"],
            "entry_type": item["entry_type"],
            "status": item["status"],
            "summary": item["summary"],
            "snippet": item["snippet"],
            "score": item["score"],
        }
        for item in shortlist
    ]
    return "\n\n".join(
        [
            "You are the Sediment submission triage assistant.",
            "Review this incoming user submission against the current KB shortlist.",
            "Be conservative. Suggest a better title or type only when clearly helpful.",
            "Return JSON only. No markdown fences and no prose before or after the JSON object.",
            "Focus on helping a human committer judge the submission quickly.",
            json.dumps(
                {
                    "submission": {
                        "title": title,
                        "type": submission_type,
                        "content": content,
                    },
                    "kb_shortlist": shortlist_payload,
                },
                ensure_ascii=False,
                indent=2,
            ),
        ]
    )


def _fallback_submission_analysis(
    *,
    title: str,
    submission_type: str,
    shortlist: list[dict[str, Any]],
) -> dict[str, Any]:
    related_entries = [
        {"name": item["name"], "reason": item["summary"] or item["snippet"] or "相关条目"}
        for item in shortlist[:4]
    ]
    highest_score = shortlist[0]["score"] if shortlist else 0
    duplicate_risk = "high" if highest_score >= 16 else "medium" if shortlist else "low"
    committer_action = "merge_or_link" if duplicate_risk == "high" else "ingest"
    if submission_type == "feedback":
        committer_action = "manual_review"
    summary = (
        "Agent analysis is unavailable, but Sediment found related KB entries for manual review."
        if shortlist
        else "Agent analysis is unavailable. The submission is stored and ready for manual review."
    )
    return {
        "status": "fallback",
        "summary": summary,
        "recommended_title": title,
        "recommended_type": (
            submission_type
            if submission_type in {"concept", "lesson", "feedback"}
            else "concept"
        ),
        "duplicate_risk": duplicate_risk,
        "committer_action": committer_action,
        "committer_note": (
            "Compare against the related entries before running ingest."
            if shortlist
            else "No obvious duplicate was found from the KB shortlist."
        ),
        "related_entries": related_entries,
        "warnings": [],
    }


def _normalize_submission_analysis(
    payload: dict[str, Any],
    *,
    title: str,
    submission_type: str,
    shortlist: list[dict[str, Any]],
) -> dict[str, Any]:
    fallback = _fallback_submission_analysis(
        title=title,
        submission_type=submission_type,
        shortlist=shortlist,
    )
    related_entries = payload.get("related_entries")
    normalized_entries: list[dict[str, str]] = []
    if isinstance(related_entries, list):
        for item in related_entries[:6]:
            if not isinstance(item, dict):
                continue
            name = normalize_name(str(item.get("name", "")), fallback="")
            reason = normalize_name(str(item.get("reason", "")), fallback="相关条目")
            if name:
                normalized_entries.append({"name": name, "reason": reason})
    return {
        "summary": normalize_name(str(payload.get("summary", "")), fallback=fallback["summary"]),
        "recommended_title": normalize_name(
            str(payload.get("recommended_title", "")),
            fallback=fallback["recommended_title"],
        ),
        "recommended_type": (
            str(payload.get("recommended_type", "")).strip()
            if str(payload.get("recommended_type", "")).strip() in {"concept", "lesson", "feedback"}
            else fallback["recommended_type"]
        ),
        "duplicate_risk": (
            str(payload.get("duplicate_risk", "")).strip()
            if str(payload.get("duplicate_risk", "")).strip() in {"low", "medium", "high"}
            else fallback["duplicate_risk"]
        ),
        "committer_action": (
            str(payload.get("committer_action", "")).strip()
            if str(payload.get("committer_action", "")).strip()
            in {"manual_review", "ingest", "merge_or_link"}
            else fallback["committer_action"]
        ),
        "committer_note": normalize_name(
            str(payload.get("committer_note", "")),
            fallback=fallback["committer_note"],
        ),
        "related_entries": normalized_entries or fallback["related_entries"],
        "warnings": [],
    }


def _expand_document_uploads(
    *,
    filename: str,
    mime_type: str,
    file_bytes: bytes,
    uploads: list[dict[str, Any]] | None,
    max_upload_bytes: int,
) -> tuple[list[dict[str, Any]], list[str]]:
    if uploads:
        return _normalize_supplied_uploads(uploads, max_upload_bytes=max_upload_bytes)
    normalized_mime = mime_type or infer_mime_type(filename or "") or ""
    if normalized_mime in ZIP_MIME_TYPES:
        return _extract_zip_upload(
            filename=filename or "upload.zip",
            file_bytes=file_bytes,
            max_upload_bytes=max_upload_bytes,
        )
    if normalized_mime not in ALLOWED_MIME_TYPES:
        raise ValueError("unsupported upload type")
    return (
        [
            {
                "filename": filename or "upload.bin",
                "relative_path": sanitize_relative_upload_path(filename or "upload.bin"),
                "mime_type": normalized_mime,
                "file_bytes": file_bytes,
            }
        ],
        [],
    )


def _normalize_supplied_uploads(
    uploads: list[dict[str, Any]],
    *,
    max_upload_bytes: int,
) -> tuple[list[dict[str, Any]], list[str]]:
    normalized: list[dict[str, Any]] = []
    skipped: list[str] = []
    total_bytes = 0
    for item in uploads:
        if not isinstance(item, dict):
            continue
        raw_bytes = item.get("file_bytes", b"")
        if not isinstance(raw_bytes, bytes) or not raw_bytes:
            continue
        filename = str(item.get("filename") or item.get("relative_path") or "upload.bin")
        relative_path = sanitize_relative_upload_path(
            str(item.get("relative_path") or filename),
            fallback=filename,
        )
        mime_type = str(item.get("mime_type") or infer_mime_type(filename) or "").strip()
        if _is_ignorable_upload(relative_path):
            skipped.append(relative_path)
            continue
        total_bytes += len(raw_bytes)
        if total_bytes > max(1, max_upload_bytes):
            raise ValueError("uploaded document bundle is too large")
        if mime_type in ZIP_MIME_TYPES:
            expanded, expanded_skipped = _extract_zip_upload(
                filename=filename,
                file_bytes=raw_bytes,
                max_upload_bytes=max_upload_bytes,
            )
            normalized.extend(expanded)
            skipped.extend(expanded_skipped)
            continue
        if mime_type not in ALLOWED_MIME_TYPES:
            skipped.append(relative_path)
            continue
        normalized.append(
            {
                "filename": filename,
                "relative_path": relative_path,
                "mime_type": mime_type,
                "file_bytes": raw_bytes,
            }
        )
    return normalized, skipped


def _extract_zip_upload(
    *,
    filename: str,
    file_bytes: bytes,
    max_upload_bytes: int,
) -> tuple[list[dict[str, Any]], list[str]]:
    supported: list[dict[str, Any]] = []
    skipped: list[str] = []
    total_unpacked = 0
    with zipfile.ZipFile(io.BytesIO(file_bytes)) as archive:
        infos = [item for item in archive.infolist() if not item.is_dir()]
        if len(infos) > 200:
            raise ValueError("archive contains too many files")
        for info in infos:
            relative_path = sanitize_relative_upload_path(info.filename, fallback=filename)
            if _is_ignorable_upload(relative_path):
                skipped.append(relative_path)
                continue
            mime_type = infer_mime_type(relative_path) or ""
            if mime_type not in ALLOWED_MIME_TYPES or mime_type in ZIP_MIME_TYPES:
                skipped.append(relative_path)
                continue
            total_unpacked += max(0, int(info.file_size))
            if total_unpacked > max(1, max_upload_bytes) * 4:
                raise ValueError("archive expands beyond the safe size limit")
            supported.append(
                {
                    "filename": Path(relative_path).name,
                    "relative_path": relative_path,
                    "mime_type": mime_type,
                    "file_bytes": archive.read(info),
                }
            )
    return supported, skipped


def _zip_document_bundle(uploads: list[dict[str, Any]]) -> bytes:
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for item in uploads:
            archive.writestr(item["relative_path"], item["file_bytes"])
    return buffer.getvalue()


def _bundle_label(*, filename: str, uploads: list[dict[str, Any]]) -> str:
    if filename and filename.strip():
        base = Path(filename).stem.strip()
        if base:
            return base
    roots = {
        item["relative_path"].split("/", 1)[0]
        for item in uploads
        if "/" in item["relative_path"]
    }
    if len(roots) == 1:
        return next(iter(roots))
    if len(uploads) == 1:
        return Path(uploads[0]["filename"]).stem or uploads[0]["filename"]
    return "document-bundle"


def sanitize_relative_upload_path(raw_value: str, *, fallback: str | None = None) -> str:
    value = (raw_value or fallback or "upload.bin").replace("\\", "/").strip().lstrip("/")
    parts: list[str] = []
    for part in value.split("/"):
        token = part.strip()
        if not token or token in {".", ".."}:
            continue
        parts.append(sanitize_filename(token))
    if not parts:
        return sanitize_filename(fallback or "upload.bin")
    return "/".join(parts)


def _is_ignorable_upload(relative_path: str) -> bool:
    lower = relative_path.lower()
    return (
        lower.startswith("__macosx/")
        or lower.endswith(".ds_store")
        or lower.endswith("thumbs.db")
    )


def submit_text(
    *,
    store: PlatformStore,
    kb_path: str | Path,
    title: str,
    content: str,
    submitter_name: str,
    submitter_ip: str,
    submission_type: str = "text",
    submitter_user_id: str | None = None,
    notes: str | None = None,
    rate_limit_count: int = 1,
    rate_limit_window_seconds: int = 60,
    max_text_chars: int = 20_000,
    dedupe_window_seconds: int = 86_400,
) -> dict[str, Any]:
    cleaned_title = normalize_name(title, fallback="未命名提交")
    cleaned_content = content.strip()
    cleaned_submitter = normalize_name(submitter_name, fallback="Anonymous")
    if not cleaned_content:
        raise ValueError("submission content must not be empty")
    if len(cleaned_content) > max(1, max_text_chars):
        raise ValueError("submission content is too large")
    dedupe_hash = build_submission_hash(cleaned_title, cleaned_content)

    analysis = analyze_text_submission(
        kb_path=kb_path,
        title=cleaned_title,
        content=cleaned_content,
        submission_type=submission_type,
    )

    record = store.create_submission_checked(
        submission_type=submission_type,
        title=cleaned_title,
        raw_text=cleaned_content,
        extracted_text=cleaned_content,
        stored_file_path=None,
        mime_type="text/plain",
        submitter_name=cleaned_submitter,
        submitter_ip=submitter_ip,
        submitter_user_id=submitter_user_id,
        dedupe_hash=dedupe_hash,
        rate_limit_count=max(1, rate_limit_count),
        rate_limit_window_seconds=max(1, rate_limit_window_seconds),
        dedupe_window_seconds=max(0, dedupe_window_seconds),
        analysis=analysis,
        notes=notes,
    )
    store.add_audit_log(
        actor_name=cleaned_submitter,
        actor_id=submitter_user_id,
        actor_role="contributor",
        action="submission.create_text",
        target_type="submission",
        target_id=record["id"],
        details={"submission_type": submission_type},
    )
    LOGGER.info(
        "submission.text.created",
        "Created text submission.",
        submission_id=record["id"],
        user_id=submitter_user_id,
        details={
            "submission_type": submission_type,
            "submitter_ip": submitter_ip,
            "title": cleaned_title,
        },
    )
    return record


def submit_document(
    *,
    store: PlatformStore,
    uploads_dir: str | Path,
    filename: str,
    mime_type: str,
    file_bytes: bytes,
    uploads: list[dict[str, Any]] | None = None,
    submitter_name: str,
    submitter_ip: str,
    submitter_user_id: str | None = None,
    notes: str | None = None,
    rate_limit_count: int = 1,
    rate_limit_window_seconds: int = 60,
    max_upload_bytes: int = 10 * 1024 * 1024,
    dedupe_window_seconds: int = 86_400,
) -> dict[str, Any]:
    if not uploads and len(file_bytes) > max(1, max_upload_bytes):
        raise ValueError("uploaded file is too large")

    cleaned_submitter = normalize_name(submitter_name, fallback="Anonymous")
    prepared = prepare_document_submission(
        filename=filename,
        mime_type=mime_type,
        file_bytes=file_bytes,
        uploads=uploads,
        max_upload_bytes=max_upload_bytes,
    )
    safe_filename = sanitize_filename(prepared["filename"] or "upload.bin")
    stored_path = Path(uploads_dir) / f"{uuid.uuid4().hex}_{safe_filename}"
    stored_path.parent.mkdir(parents=True, exist_ok=True)
    stored_path.write_bytes(prepared["file_bytes"])
    extracted_text = str(prepared["extracted_text"]).strip()
    dedupe_hash = build_submission_hash(safe_filename, prepared["mime_type"], extracted_text)
    try:
        record = store.create_submission_checked(
            submission_type="document",
            title=str(prepared["title"]),
            raw_text=base64.b64encode(prepared["file_bytes"]).decode("ascii"),
            extracted_text=extracted_text,
            stored_file_path=str(stored_path),
            mime_type=str(prepared["mime_type"]),
            submitter_name=cleaned_submitter,
            submitter_ip=submitter_ip,
            submitter_user_id=submitter_user_id,
            dedupe_hash=dedupe_hash,
            rate_limit_count=max(1, rate_limit_count),
            rate_limit_window_seconds=max(1, rate_limit_window_seconds),
            dedupe_window_seconds=max(0, dedupe_window_seconds),
            notes=prepared["notes"] or notes,
        )
    except Exception:
        stored_path.unlink(missing_ok=True)
        raise
    store.add_audit_log(
        actor_name=cleaned_submitter,
        actor_id=submitter_user_id,
        actor_role="contributor",
        action="submission.create_document",
        target_type="submission",
        target_id=record["id"],
        details={
            "filename": safe_filename,
            "mime_type": prepared["mime_type"],
            "file_count": prepared["file_count"],
        },
    )
    LOGGER.info(
        "submission.document.created",
        "Created document submission.",
        submission_id=record["id"],
        user_id=submitter_user_id,
        details={
            "filename": safe_filename,
            "mime_type": prepared["mime_type"],
            "file_count": prepared["file_count"],
            "submitter_ip": submitter_ip,
        },
    )
    return record


def submit_feedback_item(
    *,
    store: PlatformStore,
    title: str,
    content: str,
    submitter_name: str,
    submitter_ip: str,
    submitter_user_id: str | None = None,
    notes: str | None = None,
    rate_limit_count: int = 1,
    rate_limit_window_seconds: int = 60,
    max_text_chars: int = 20_000,
    dedupe_window_seconds: int = 86_400,
) -> dict[str, Any]:
    cleaned_title = normalize_name(title, fallback="Untitled feedback")
    cleaned_content = content.strip()
    cleaned_submitter = normalize_name(submitter_name, fallback="Anonymous")
    if not cleaned_content:
        raise ValueError("submission content must not be empty")
    if len(cleaned_content) > max(1, max_text_chars):
        raise ValueError("submission content is too large")
    item = store.create_inbox_item_checked(
        item_type="text_feedback",
        title=cleaned_title,
        body_text=cleaned_content,
        stored_file_path=None,
        original_filename=None,
        mime_type="text/plain",
        submitter_name=cleaned_submitter,
        submitter_ip=submitter_ip,
        submitter_user_id=submitter_user_id,
        dedupe_hash=build_submission_hash(cleaned_title, cleaned_content),
        rate_limit_count=max(1, rate_limit_count),
        rate_limit_window_seconds=max(1, rate_limit_window_seconds),
        dedupe_window_seconds=max(0, dedupe_window_seconds),
        status="open",
        notes=notes,
    )
    store.add_audit_log(
        actor_name=cleaned_submitter,
        actor_id=submitter_user_id,
        actor_role="contributor",
        action="inbox.create_text_feedback",
        target_type="inbox_item",
        target_id=item["id"],
        details={"item_type": "text_feedback"},
    )
    return item


def submit_uploaded_document_item(
    *,
    store: PlatformStore,
    uploads_dir: str | Path,
    filename: str,
    mime_type: str,
    file_bytes: bytes,
    uploads: list[dict[str, Any]] | None = None,
    submitter_name: str,
    submitter_ip: str,
    submitter_user_id: str | None = None,
    notes: str | None = None,
    rate_limit_count: int = 1,
    rate_limit_window_seconds: int = 60,
    max_upload_bytes: int = 10 * 1024 * 1024,
    dedupe_window_seconds: int = 86_400,
) -> dict[str, Any]:
    if not uploads and len(file_bytes) > max(1, max_upload_bytes):
        raise ValueError("uploaded file is too large")
    cleaned_submitter = normalize_name(submitter_name, fallback="Anonymous")
    prepared = prepare_document_staging_upload(
        filename=filename,
        mime_type=mime_type,
        file_bytes=file_bytes,
        uploads=uploads,
        max_upload_bytes=max_upload_bytes,
    )
    safe_filename = sanitize_filename(prepared["filename"] or "upload.bin")
    stored_path = Path(uploads_dir) / f"{uuid.uuid4().hex}_{safe_filename}"
    stored_path.parent.mkdir(parents=True, exist_ok=True)
    stored_path.write_bytes(prepared["file_bytes"])
    try:
        item = store.create_inbox_item_checked(
            item_type="uploaded_document",
            title=str(prepared["title"]),
            body_text="",
            stored_file_path=str(stored_path),
            original_filename=safe_filename,
            mime_type=str(prepared["mime_type"]),
            submitter_name=cleaned_submitter,
            submitter_ip=submitter_ip,
            submitter_user_id=submitter_user_id,
            dedupe_hash=build_submission_hash(
                safe_filename,
                prepared["mime_type"],
                base64.b64encode(prepared["file_bytes"]).decode("ascii"),
            ),
            rate_limit_count=max(1, rate_limit_count),
            rate_limit_window_seconds=max(1, rate_limit_window_seconds),
            dedupe_window_seconds=max(0, dedupe_window_seconds),
            status="staged",
            notes=prepared["notes"] or notes,
        )
    except Exception:
        stored_path.unlink(missing_ok=True)
        raise
    store.add_audit_log(
        actor_name=cleaned_submitter,
        actor_id=submitter_user_id,
        actor_role="contributor",
        action="inbox.create_uploaded_document",
        target_type="inbox_item",
        target_id=item["id"],
        details={
            "filename": safe_filename,
            "mime_type": prepared["mime_type"],
            "file_count": prepared["file_count"],
        },
    )
    return item


def sanitize_filename(filename: str) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9._-]+", "_", filename.strip())
    return cleaned or "upload.bin"


def get_portal_home(kb_path: str | Path, *, store: PlatformStore) -> dict[str, Any]:
    data = inventory(kb_path)
    report = audit_kb(kb_path)
    inbox_counts = store.inbox_status_counts() if hasattr(store, "inbox_status_counts") else {}
    docs = data["docs"]
    recent = sorted(
        (
            {
                "name": name,
                "entry_type": doc["entry_type"],
                "status": doc["status"],
                "updated_at": Path(doc["path"]).stat().st_mtime if doc["path"] else 0,
            }
            for name, doc in docs.items()
            if doc["path"]
        ),
        key=lambda item: item["updated_at"],
        reverse=True,
    )[:8]
    for item in recent:
        item["updated_at"] = int(item["updated_at"])
    return {
        "counts": {
            "formal_entries": len(data["entries"]),
            "placeholders": len(data["placeholders"]),
            "indexes": len(data["indexes"]),
            "pending_submissions": (
                inbox_counts.get("open", 0)
                + inbox_counts.get("staged", 0)
                + inbox_counts.get("ready", 0)
                + inbox_counts.get("ingesting", 0)
            ),
            "health_issues": len(build_health_issue_queue(kb_path)),
        },
        "recent_updates": recent,
        "health_summary": summarize_health_report(report),
    }


def _graph_ref(subject_kind: str, subject_id: str) -> str:
    normalized_kind = str(subject_kind or "").strip()
    normalized_id = str(subject_id or "").strip()
    if not normalized_id:
        return ""
    if normalized_kind == "canonical_entry":
        return f"entry::{normalized_id}"
    if normalized_kind == "insight_proposal":
        return f"insight::{normalized_id}"
    if normalized_kind == "query_cluster":
        return f"cluster::{normalized_id}"
    if normalized_kind == "index_segment":
        return f"index::{normalized_id}"
    if normalized_kind == "cluster_anchor":
        return f"anchor::{normalized_id}"
    return normalized_id


def _split_graph_ref(value: str) -> tuple[str, str]:
    normalized = str(value or "").strip()
    if "::" not in normalized:
        return "", normalized
    prefix, raw_id = normalized.split("::", 1)
    return prefix, raw_id


def _event_timestamp(value: str | None) -> datetime:
    if not value:
        return datetime.now(timezone.utc)
    try:
        parsed = datetime.fromisoformat(str(value).replace("Z", "+00:00"))
    except ValueError:
        return datetime.now(timezone.utc)
    if parsed.tzinfo is None:
        return parsed.replace(tzinfo=timezone.utc)
    return parsed.astimezone(timezone.utc)


def _graph_event_score(event: dict[str, Any]) -> float:
    event_type = str(event.get("event_type") or "").strip()
    details = event.get("details") or {}
    base_energy = float(details.get("energy") or GRAPH_EVENT_BASE_ENERGY.get(event_type, 0.45))
    age_seconds = max(
        (datetime.now(timezone.utc) - _event_timestamp(str(event.get("created_at") or ""))).total_seconds(),
        0.0,
    )
    half_life_days = float(details.get("half_life_days") or 18.0)
    half_life_seconds = max(half_life_days, 1.0) * 24 * 60 * 60
    decay = math.exp(-age_seconds / half_life_seconds)
    reinforcement = 1.0 + float(details.get("reinforcement") or 0.0)
    return round(base_energy * decay * reinforcement, 4)


def _event_priority(event_type: str) -> int:
    return GRAPH_EVENT_PRIORITY.get(str(event_type or "").strip(), 0)


def record_graph_event(
    store: PlatformStore | None,
    *,
    event_type: str,
    subject_id: str,
    subject_kind: str,
    subject_label: str,
    entry_target: str | None = None,
    related_ids: list[str] | None = None,
    details: dict[str, Any] | None = None,
    created_at: str | None = None,
) -> dict[str, Any] | None:
    if store is None or not hasattr(store, "record_graph_event"):
        return None
    return store.record_graph_event(
        event_type=event_type,
        subject_id=subject_id,
        subject_kind=subject_kind,
        subject_label=subject_label,
        entry_target=entry_target,
        related_ids=related_ids or [],
        details=details or {},
        created_at=created_at,
    )


def emit_explore_graph_events(
    *,
    store: PlatformStore | None,
    cluster: dict[str, Any] | None,
    proposal: dict[str, Any] | None,
) -> None:
    if store is None or cluster is None:
        return
    source_entries = [
        _graph_ref("canonical_entry", item)
        for item in cluster.get("source_entries") or []
        if str(item).strip()
    ]
    proposal_ref = (
        _graph_ref("insight_proposal", str(proposal.get("id") or ""))
        if proposal is not None
        else ""
    )
    related_ids = list(source_entries)
    if proposal_ref:
        related_ids.append(proposal_ref)
    record_graph_event(
        store,
        event_type="ask_reinforced",
        subject_id=str(cluster.get("id") or cluster.get("cluster_key") or ""),
        subject_kind="query_cluster",
        subject_label=str(
            cluster.get("display_query")
            or cluster.get("normalized_subject")
            or cluster.get("cluster_key")
            or "Query cluster"
        ),
        related_ids=related_ids,
        details={
            "energy": float(cluster.get("demand_score") or 0.0) + 0.35,
            "stability": float(cluster.get("maturity_score") or 0.0),
            "intent": str(cluster.get("intent") or ""),
            "cluster_key": str(cluster.get("cluster_key") or ""),
            "visual_role": "reinforced_query",
            "source_entries": [item for item in cluster.get("source_entries") or [] if str(item).strip()],
            "insight_id": str(proposal.get("id") or cluster.get("insight_id") or "")
            if proposal is not None or cluster.get("insight_id")
            else "",
        },
        created_at=str(cluster.get("updated_at") or cluster.get("last_seen_at") or utc_now()),
    )
    if proposal is None:
        return
    supporting_entries = [
        _graph_ref("canonical_entry", item)
        for item in proposal.get("supporting_entries") or []
        if str(item).strip()
    ]
    record_graph_event(
        store,
        event_type="proposal_materialized",
        subject_id=str(proposal.get("id") or ""),
        subject_kind="insight_proposal",
        subject_label=str(proposal.get("title") or proposal.get("id") or "Insight proposal"),
        related_ids=[_graph_ref("query_cluster", str(cluster.get("id") or "")), *supporting_entries],
        details={
            "energy": 0.94,
            "stability": float(cluster.get("maturity_score") or 0.0),
            "kind": str(proposal.get("kind") or "concept"),
            "review_state": str(proposal.get("review_state") or "proposed"),
            "hypothesis": str(proposal.get("hypothesis") or ""),
            "proposed_answer": str(proposal.get("proposed_answer") or ""),
            "supporting_entries": [
                item for item in proposal.get("supporting_entries") or [] if str(item).strip()
            ],
            "trigger_queries": [
                item for item in proposal.get("trigger_queries") or [] if str(item).strip()
            ],
            "visual_role": "forming_insight",
        },
    )


def emit_managed_graph_events(
    *,
    kb_path: str | Path,
    store: PlatformStore | None,
    operation: str,
    request_payload: dict[str, Any],
    payload: dict[str, Any],
    apply_result: dict[str, Any],
    commit_sha: str,
) -> None:
    if store is None:
        return
    kb_root = Path(kb_path).resolve()
    applied_by_path = {
        str(item.get("relative_path") or ""): item for item in apply_result.get("operations", [])
    }
    planned_by_path = {
        str(item.get("relative_path") or ""): item for item in payload.get("operations", [])
    }

    if operation == "ingest":
        for relative_path, applied in applied_by_path.items():
            if not relative_path.startswith("entries/"):
                continue
            planned = planned_by_path.get(relative_path) or {}
            content = str(planned.get("content") or "")
            related_links = [
                _graph_ref("canonical_entry", item)
                for item in extract_wikilinks(content)
                if str(item).strip()
            ]
            event_type = "ingest_created" if applied.get("change_type") == "create" else "ingest_updated"
            record_graph_event(
                store,
                event_type=event_type,
                subject_id=str(planned.get("name") or Path(relative_path).stem),
                subject_kind="canonical_entry",
                subject_label=str(planned.get("name") or Path(relative_path).stem),
                entry_target=str(planned.get("name") or Path(relative_path).stem),
                related_ids=related_links,
                details={
                    "energy": 1.0 if event_type == "ingest_created" else 0.72,
                    "stability": 0.82 if event_type == "ingest_created" else 0.9,
                    "relative_path": relative_path,
                    "change_type": str(applied.get("change_type") or ""),
                    "commit_sha": commit_sha,
                    "summary": str(payload.get("summary") or ""),
                    "visual_role": "fresh_ingest" if event_type == "ingest_created" else "refreshed_entry",
                },
            )
        return

    if operation != "insight":
        return

    action = str(request_payload.get("action") or "").strip()
    insight_id = str(request_payload.get("insight_id") or "").strip()
    if action not in {"promote", "merge"} or not insight_id:
        return
    detail = get_insight_detail(kb_root, insight_id, store=store)
    proposal = detail["proposal"]
    supporting_entries = [
        _graph_ref("canonical_entry", item)
        for item in proposal.get("supporting_entries") or []
        if str(item).strip()
    ]
    event_type = "insight_promoted" if action == "promote" else "insight_merged"
    target_name = ""
    for relative_path, applied in applied_by_path.items():
        if relative_path.startswith("entries/") and applied.get("change_type") in {"create", "update"}:
            target_name = str((planned_by_path.get(relative_path) or {}).get("name") or Path(relative_path).stem)
            break
    if not target_name:
        target_name = str(
            request_payload.get("new_title")
            or request_payload.get("target_name")
            or proposal.get("title")
            or insight_id
        ).strip()
    record_graph_event(
        store,
        event_type=event_type,
        subject_id=target_name,
        subject_kind="canonical_entry",
        subject_label=target_name,
        entry_target=target_name,
        related_ids=[_graph_ref("insight_proposal", insight_id), *supporting_entries],
        details={
            "energy": 1.04 if event_type == "insight_promoted" else 0.88,
            "stability": 0.98 if event_type == "insight_promoted" else 0.95,
            "proposal_id": insight_id,
            "proposal_title": str(proposal.get("title") or insight_id),
            "supporting_entries": [
                item for item in proposal.get("supporting_entries") or [] if str(item).strip()
            ],
            "visual_role": "recent_canonical",
            "commit_sha": commit_sha,
            "review_note": str(request_payload.get("note") or ""),
        },
    )


def _query_cluster_stats(kb_path: str | Path, store: PlatformStore | None) -> dict[str, Any]:
    if store is None or not hasattr(store, "list_signal_clusters"):
        return {
            "clusters": [],
            "emerging": [],
            "stress_points": [],
            "cluster_coverage": 0.0,
            "ready_clusters": [],
        }
    data = inventory(kb_path)
    clusters = list(store.list_signal_clusters(limit=200))
    for cluster in clusters:
        source_entries = [item for item in cluster.get("source_entries") or [] if item in data["docs"]]
        cluster["source_entries"] = source_entries
        cluster["source_entry_count"] = len(source_entries)
    total_formal = max(len(data["entries"]), 1)
    covered_entries = {
        entry
        for cluster in clusters
        for entry in cluster.get("source_entries") or []
        if entry in data["entries"]
    }
    cluster_coverage = round(len(covered_entries) / total_formal, 3) if data["entries"] else 0.0
    emerging = sorted(
        [
            cluster
            for cluster in clusters
            if cluster.get("status") in {"clustered", "ready", "materialized"}
        ],
        key=lambda item: (
            -float(item.get("demand_score") or 0.0),
            -float(item.get("maturity_score") or 0.0),
            str(item.get("last_seen_at") or ""),
        ),
    )[:8]
    stress_points: list[dict[str, Any]] = []
    by_entry: dict[str, dict[str, Any]] = defaultdict(
        lambda: {"entry": "", "cluster_count": 0, "signal_count": 0, "synthesized_count": 0, "query_examples": []}
    )
    for cluster in clusters:
        for entry in cluster.get("source_entries") or []:
            bucket = by_entry[entry]
            bucket["entry"] = entry
            bucket["cluster_count"] += 1
            bucket["signal_count"] += int(cluster.get("signal_count") or 0)
            if str(cluster.get("result_mode") or "") == "synthesized":
                bucket["synthesized_count"] += 1
            if cluster.get("display_query") and len(bucket["query_examples"]) < 3:
                bucket["query_examples"].append(cluster["display_query"])
    stress_points = sorted(
        by_entry.values(),
        key=lambda item: (-item["signal_count"], -item["synthesized_count"], item["entry"]),
    )[:8]
    ready_clusters = [
        cluster
        for cluster in clusters
        if cluster.get("status") == "ready" and not cluster.get("insight_id")
    ]
    return {
        "clusters": clusters,
        "emerging": emerging,
        "stress_points": stress_points,
        "cluster_coverage": cluster_coverage,
        "ready_clusters": ready_clusters,
    }


def _insights_dir(kb_path: str | Path) -> Path:
    return Path(kb_path).resolve() / "insights"


def list_insight_proposals(kb_path: str | Path) -> list[dict[str, Any]]:
    data = inventory(kb_path)
    return sorted(
        list(data.get("insight_docs", {}).values()),
        key=lambda item: (
            {"proposed": 0, "observing": 1, "promoted": 2, "merged": 3, "rejected": 4, "archived": 5}.get(
                str(item.get("review_state") or "proposed"),
                9,
            ),
            str(item.get("title") or item.get("id") or ""),
        ),
    )


def get_insight_detail(kb_path: str | Path, insight_id: str, *, store: PlatformStore | None = None) -> dict[str, Any]:
    data = inventory(kb_path)
    for proposal in data.get("insight_docs", {}).values():
        if proposal.get("id") == insight_id or proposal.get("name") == insight_id:
            cluster = None
            if store is not None and hasattr(store, "list_signal_clusters"):
                candidates = [
                    item
                    for item in store.list_signal_clusters(limit=200)
                    if item.get("insight_id") == proposal["id"]
                ]
                cluster = candidates[0] if candidates else None
            return {
                "proposal": proposal,
                "cluster": cluster,
                "recommended_action": infer_recommended_action(
                    kind=str(proposal.get("kind") or "concept"),
                    supporting_entries=list(proposal.get("supporting_entries") or []),
                ),
            }
    raise FileNotFoundError(insight_id)


def _render_frontmatter_with_body(frontmatter: dict[str, Any], body: str) -> str:
    yaml_block = yaml.safe_dump(frontmatter, allow_unicode=True, sort_keys=False).strip()
    return f"---\n{yaml_block}\n---\n\n{body.strip()}\n"


def _ensure_section(sections: dict[str, str], name: str, content: str) -> dict[str, str]:
    updated = dict(sections)
    existing = str(updated.get(name) or "").strip()
    addition = content.strip()
    if not addition:
        return updated
    if existing:
        if addition not in existing:
            updated[name] = f"{existing}\n\n{addition}".strip()
    else:
        updated[name] = addition
    return updated


def _compose_entry_body(*, title: str, summary: str, sections: dict[str, str], entry_type: str) -> str:
    order = ["Scope", "Related"] if entry_type == "concept" else ["Trigger", "Why", "Risks", "Related"]
    body_parts = [f"# {title}", "", summary.strip() or title]
    rendered_related = sections.get("Related", "")
    if rendered_related and "[[" not in rendered_related:
        rendered_related = "\n".join(f"- [[{item.strip()}]]" for item in rendered_related.splitlines() if item.strip())
    for name in order:
        content = sections.get(name, "").strip()
        if not content:
            continue
        body_parts.extend(["", f"## {name}", content])
    for name, content in sections.items():
        if name in order or not str(content).strip():
            continue
        body_parts.extend(["", f"## {name}", str(content).strip()])
    return "\n".join(body_parts).strip() + "\n"


def prepare_insight_review_payload(
    kb_path: str | Path,
    *,
    insight_id: str,
    action: str,
    target_name: str | None = None,
    note: str = "",
    entry_type: str | None = None,
    new_title: str | None = None,
) -> dict[str, Any]:
    root = Path(kb_path).resolve()
    detail = get_insight_detail(root, insight_id)
    proposal = detail["proposal"]
    insight_path = Path(proposal["path"])
    insight_content = insight_path.read_text(encoding="utf-8")
    insight_frontmatter, insight_body = split_frontmatter(insight_content)
    insight_sections, _ = split_sections(insight_body)
    summary = str(proposal.get("proposed_answer") or proposal.get("hypothesis") or proposal.get("title") or "").strip()
    supporting_entries = list(proposal.get("supporting_entries") or [])

    updated_frontmatter = dict(insight_frontmatter)
    updated_sections = dict(insight_sections)
    operations: list[dict[str, Any]] = []
    warnings: list[str] = []

    if note.strip():
        updated_sections = _ensure_section(updated_sections, "Review Notes", f"- {note.strip()}")

    if action == "observe":
        updated_frontmatter["review_state"] = "observing"
        operations.append(
            {
                "name": insight_path.stem,
                "relative_path": str(insight_path.relative_to(root)),
                "change_type": "update",
                "rationale": "Mark insight proposal as observing",
                "content": _render_frontmatter_with_body(
                    updated_frontmatter,
                    _compose_entry_body(
                        title=str(updated_frontmatter.get("title") or proposal["title"]),
                        summary=str(updated_frontmatter.get("hypothesis") or summary),
                        sections=updated_sections,
                        entry_type="concept",
                    ),
                ),
                "base_hash": content_hash(insight_content),
            }
        )
        return {"summary": "Insight kept in observing state.", "warnings": warnings, "operations": operations}

    if action == "reject":
        updated_frontmatter["review_state"] = "rejected"
        operations.append(
            {
                "name": insight_path.stem,
                "relative_path": str(insight_path.relative_to(root)),
                "change_type": "update",
                "rationale": "Reject insight proposal",
                "content": _render_frontmatter_with_body(
                    updated_frontmatter,
                    _compose_entry_body(
                        title=str(updated_frontmatter.get("title") or proposal["title"]),
                        summary=str(updated_frontmatter.get("hypothesis") or summary),
                        sections=updated_sections,
                        entry_type="concept",
                    ),
                ),
                "base_hash": content_hash(insight_content),
            }
        )
        return {"summary": "Insight rejected.", "warnings": warnings, "operations": operations}

    if action == "promote":
        canonical_name = str(new_title or proposal.get("title") or proposal["name"]).strip()
        canonical_entry_type = str(entry_type or ("lesson" if proposal.get("kind") == "lesson" else "concept")).strip()
        canonical_frontmatter = {
            "type": canonical_entry_type,
            "status": "inferred",
            "aliases": [],
            "sources": [f"insight:{proposal['id']}"],
        }
        canonical_sections = {
            "Scope" if canonical_entry_type == "concept" else "Trigger": "Derived from a reviewed insight proposal.",
            "Why" if canonical_entry_type == "lesson" else "Related": "\n".join(
                f"- [[{item}]]" for item in supporting_entries
            ).strip(),
        }
        if canonical_entry_type == "lesson":
            canonical_sections["Why"] = summary
            canonical_sections["Risks"] = "Validate and refine this inferred workflow with domain owners."
            canonical_sections["Related"] = "\n".join(f"- [[{item}]]" for item in supporting_entries).strip()
        else:
            canonical_sections["Scope"] = summary
            canonical_sections["Related"] = "\n".join(f"- [[{item}]]" for item in supporting_entries).strip()
        canonical_body = _compose_entry_body(
            title=canonical_name,
            summary=summary,
            sections=canonical_sections,
            entry_type=canonical_entry_type,
        )
        operations.append(
            {
                "name": canonical_name,
                "relative_path": f"entries/{canonical_name}.md",
                "change_type": "update" if (root / "entries" / f"{canonical_name}.md").exists() else "create",
                "rationale": "Promote reviewed insight proposal into canonical knowledge",
                "content": _render_frontmatter_with_body(canonical_frontmatter, canonical_body),
                "base_hash": content_hash((root / "entries" / f"{canonical_name}.md").read_text(encoding="utf-8"))
                if (root / "entries" / f"{canonical_name}.md").exists()
                else None,
            }
        )
        updated_frontmatter["review_state"] = "promoted"
        updated_sections = _ensure_section(updated_sections, "Review Notes", f"- Promoted to [[{canonical_name}]].")
        operations.append(
            {
                "name": insight_path.stem,
                "relative_path": str(insight_path.relative_to(root)),
                "change_type": "update",
                "rationale": "Mark insight as promoted",
                "content": _render_frontmatter_with_body(
                    updated_frontmatter,
                    _compose_entry_body(
                        title=str(updated_frontmatter.get("title") or proposal["title"]),
                        summary=str(updated_frontmatter.get("hypothesis") or summary),
                        sections=updated_sections,
                        entry_type="concept",
                    ),
                ),
                "base_hash": content_hash(insight_content),
            }
        )
        return {"summary": f"Promoted insight into `{canonical_name}`.", "warnings": warnings, "operations": operations}

    if action == "merge":
        if not target_name:
            raise ValueError("target_name is required for merge")
        target_path = resolve_kb_document_path(root, target_name)
        if target_path is None:
            raise FileNotFoundError(target_name)
        target_content = target_path.read_text(encoding="utf-8")
        target_frontmatter, target_body = split_frontmatter(target_content)
        target_sections, target_preamble = split_sections(target_body)
        target_entry_type = str(target_frontmatter.get("type") or "concept").strip() or "concept"
        target_summary = target_preamble.strip() or target_name
        if target_entry_type == "lesson":
            target_sections = _ensure_section(target_sections, "Why", summary)
        else:
            target_sections = _ensure_section(target_sections, "Scope", summary)
        target_sections = _ensure_section(
            target_sections,
            "Related",
            "\n".join(f"- [[{item}]]" for item in supporting_entries if item != target_name).strip(),
        )
        operations.append(
            {
                "name": target_name,
                "relative_path": str(target_path.relative_to(root)),
                "change_type": "update",
                "rationale": "Merge reviewed insight proposal into canonical entry",
                "content": _render_frontmatter_with_body(
                    target_frontmatter,
                    _compose_entry_body(
                        title=target_name,
                        summary=target_summary,
                        sections=target_sections,
                        entry_type=target_entry_type,
                    ),
                ),
                "base_hash": content_hash(target_content),
            }
        )
        updated_frontmatter["review_state"] = "merged"
        updated_sections = _ensure_section(updated_sections, "Review Notes", f"- Merged into [[{target_name}]].")
        operations.append(
            {
                "name": insight_path.stem,
                "relative_path": str(insight_path.relative_to(root)),
                "change_type": "update",
                "rationale": "Mark insight as merged",
                "content": _render_frontmatter_with_body(
                    updated_frontmatter,
                    _compose_entry_body(
                        title=str(updated_frontmatter.get("title") or proposal["title"]),
                        summary=str(updated_frontmatter.get("hypothesis") or summary),
                        sections=updated_sections,
                        entry_type="concept",
                    ),
                ),
                "base_hash": content_hash(insight_content),
            }
        )
        return {"summary": f"Merged insight into `{target_name}`.", "warnings": warnings, "operations": operations}

    raise ValueError(f"unsupported insight review action: {action}")


def materialize_insight_proposal(
    kb_path: str | Path,
    *,
    cluster: dict[str, Any],
) -> dict[str, Any]:
    root = Path(kb_path).resolve()
    insights_dir = _insights_dir(root)
    insights_dir.mkdir(parents=True, exist_ok=True)
    language = str(cluster.get("language") or "en")
    title = insight_title_from_cluster(
        normalized_subject=str(cluster.get("normalized_subject") or cluster.get("display_query") or ""),
        language=language,
        intent=str(cluster.get("intent") or "definition"),
    )
    cluster_key = str(cluster.get("cluster_key") or "")
    insight_id = f"insight-{slugify_filename(title)}"
    if cluster_key:
        fingerprint = hashlib.sha1(cluster_key.encode("utf-8")).hexdigest()[:8]
        insight_id = f"{insight_id}-{fingerprint}"
    frontmatter = insight_frontmatter(cluster, insight_id=insight_id, title=title)
    content = render_insight_markdown(frontmatter)
    filename = f"{slugify_filename(title)}.md"
    target = insights_dir / filename
    suffix = 2
    while target.exists():
        existing = parse_insight(target)
        if existing.get("id") == insight_id:
            break
        target = insights_dir / f"{slugify_filename(title)}-{suffix}.md"
        suffix += 1
    target.write_text(content, encoding="utf-8")
    proposal = parse_insight(target)
    proposal["recommended_action"] = infer_recommended_action(
        kind=str(proposal.get("kind") or "concept"),
        supporting_entries=list(proposal.get("supporting_entries") or []),
    )
    return proposal


def record_explore_signal(
    *,
    kb_path: str | Path,
    store: PlatformStore,
    question: str,
    entrypoint: str,
    strategy: str,
    result: dict[str, Any],
    actor_fingerprint_value: str = "",
    response_language: str | None = None,
) -> dict[str, Any]:
    kb_snapshot = inventory(kb_path)
    kb_language = str(kb_snapshot.get("default_language") or "en")
    query_language = detect_query_language(question, default_language=kb_language)
    normalized_query = normalize_query_for_kb(question, kb_language=kb_language)
    intent = detect_intent(normalized_query, language=query_language)
    subject = normalize_subject(normalized_query, language=query_language)
    cluster_key = build_cluster_key(
        language=kb_language,
        intent=intent,
        normalized_subject=subject,
    )
    sources = [str(item).strip() for item in result.get("sources") or [] if str(item).strip()]
    mode = str(result.get("mode") or "").strip()
    if not mode:
        if not sources:
            mode = "gap"
        elif len(sources) == 1 and not (result.get("gaps") or []):
            mode = "direct"
        else:
            mode = "synthesized"
    signal = store.record_question_signal(
        raw_query=question,
        normalized_query=normalized_query,
        query_language=query_language,
        kb_language=kb_language,
        response_language=response_language or query_language,
        entrypoint=entrypoint,
        strategy=strategy,
        result_mode=mode,
        confidence=str(result.get("confidence") or "low"),
        source_entries=sources,
        actor_fingerprint=actor_fingerprint_value,
        cluster_key=cluster_key,
        intent=intent,
        normalized_subject=subject,
        answer_excerpt=str(result.get("answer") or "")[:500],
    )
    cluster = store.get_signal_cluster_by_key(cluster_key) or {}
    proposal = None
    if cluster.get("status") == "ready" and not cluster.get("insight_id"):
        proposal = materialize_insight_proposal(kb_path, cluster=cluster)
        store.attach_insight_to_cluster(
            cluster_id=str(cluster["id"]),
            insight_id=str(proposal["id"]),
            status="materialized",
        )
        cluster = store.get_signal_cluster(str(cluster["id"])) or cluster
    emit_explore_graph_events(
        store=store,
        cluster=cluster,
        proposal=proposal,
    )
    envelope = {
        "query": question,
        "entrypoint": entrypoint,
        "strategy": strategy,
        "mode": mode,
        "answer": str(result.get("answer") or ""),
        "confidence": str(result.get("confidence") or "low"),
        "sources": sources,
        "proposal_state": "materialized" if proposal else str(cluster.get("status") or "none"),
        "query_language": query_language,
        "kb_language": kb_language,
        "response_language": response_language or query_language,
        "kb_normalized_query": normalized_query,
        "signal_id": signal.get("id"),
        "signal_cluster_id": cluster.get("id"),
        "insight_id": proposal.get("id") if proposal else cluster.get("insight_id"),
    }
    return {"signal": signal, "cluster": cluster, "proposal": proposal, "envelope": envelope}


def search_kb(kb_path: str | Path, query: str, *, limit: int = 20) -> list[dict[str, Any]]:
    raw_query = query.strip()
    if not raw_query:
        return []
    data = inventory(kb_path)
    terms = [term for term in re.split(r"\s+", raw_query) if term]
    results: list[dict[str, Any]] = []
    for name, doc in data["docs"].items():
        score = 0
        haystacks = {
            "title": doc["title"],
            "aliases": " ".join(doc["aliases"]),
            "summary": doc["summary"],
            "body": doc["body"],
        }
        for term in terms:
            lowered = term.casefold()
            if lowered in haystacks["title"].casefold():
                score += 10
            if lowered in haystacks["aliases"].casefold():
                score += 8
            if lowered in haystacks["summary"].casefold():
                score += 6
            if lowered in haystacks["body"].casefold():
                score += 3
        if score <= 0:
            continue
        score += min(doc["inbound_count"], 5)
        if doc["kind"] == "formal":
            score += 2
        snippet = build_search_snippet(doc["body"], terms)
        results.append(
            {
                "name": name,
                "title": doc["title"],
                "kind": doc["kind"],
                "entry_type": doc["entry_type"],
                "status": doc["status"],
                "summary": doc["summary"],
                "snippet": snippet,
                "score": score,
            }
        )
    return sorted(results, key=lambda item: (-item["score"], item["name"]))[:limit]


def search_kb_suggestions(
    kb_path: str | Path,
    query: str,
    *,
    limit: int = 8,
) -> list[dict[str, Any]]:
    raw_query = query.strip()
    if not raw_query:
        return []
    data = inventory(kb_path)
    terms = [term for term in re.split(r"\s+", raw_query) if term]
    suggestions: list[dict[str, Any]] = []
    for name, doc in data["docs"].items():
        score = 0
        matched_field = ""
        for term in terms:
            lowered = term.casefold()
            if doc["title"].casefold().startswith(lowered) or name.casefold().startswith(lowered):
                score += 16
                matched_field = matched_field or "title"
            elif lowered in doc["title"].casefold():
                score += 12
                matched_field = matched_field or "title"
            elif any(alias.casefold().startswith(lowered) for alias in doc["aliases"]):
                score += 10
                matched_field = matched_field or "aliases"
            elif lowered in " ".join(doc["aliases"]).casefold():
                score += 8
                matched_field = matched_field or "aliases"
            elif lowered in doc["summary"].casefold():
                score += 4
                matched_field = matched_field or "summary"
        if score <= 0:
            continue
        suggestions.append(
            {
                "name": name,
                "title": doc["title"],
                "kind": doc["kind"],
                "entry_type": doc["entry_type"],
                "status": doc["status"],
                "summary": doc["summary"],
                "matched_field": matched_field or "summary",
                "score": score + min(int(doc["inbound_count"]), 5),
            }
        )
    return sorted(suggestions, key=lambda item: (-item["score"], item["name"]))[:limit]


def build_search_snippet(body: str, terms: list[str]) -> str:
    compact = re.sub(r"\s+", " ", body).strip()
    if not compact:
        return ""
    lowered = compact.casefold()
    for term in terms:
        idx = lowered.find(term.casefold())
        if idx != -1:
            start = max(0, idx - 80)
            end = min(len(compact), idx + 140)
            snippet = compact[start:end].strip()
            if start > 0:
                snippet = "..." + snippet
            if end < len(compact):
                snippet = snippet + "..."
            return snippet
    return compact[:180] + ("..." if len(compact) > 180 else "")


def get_entry_detail(kb_path: str | Path, name: str) -> dict[str, Any]:
    data = inventory(kb_path)
    doc = data["docs"].get(name) or data["index_docs"].get(name)
    if doc is None:
        raise FileNotFoundError(name)
    path = resolve_kb_document_path(kb_path, name)
    content = path.read_text(encoding="utf-8") if path else ""
    validation = None
    if path and (path.parent.name == "indexes" or path.name == index_config()["root_file"]):
        validation = validate_index(path)
    elif path:
        kind = "placeholder" if "/placeholders/" in str(path) else "formal"
        validation = validate_entry(path=path, kind=kind)
    sections_map = dict(doc.get("sections_map") or {})
    canonical_section_order = ("Scope", "Trigger", "Why", "Risks", "Related")
    canonical_sections = [
        {"name": section_name, "content": str(sections_map.get(section_name, "")).strip()}
        for section_name in canonical_section_order
        if str(sections_map.get(section_name, "")).strip()
    ]
    residual_sections = [
        {"name": section_name, "content": str(section_body).strip()}
        for section_name, section_body in sections_map.items()
        if section_name not in canonical_section_order and str(section_body).strip()
    ]
    residual_parts: list[str] = []
    preamble = str(doc.get("preamble", "")).strip()
    if preamble:
        residual_parts.append(preamble)
    for section in residual_sections:
        residual_parts.append(f"## {section['name']}\n{section['content']}")
    residual_markdown = "\n\n".join(part for part in residual_parts if part).strip()
    related_links = extract_wikilinks(sections_map.get("Related", "")) or list(
        doc.get("graph_links") or []
    )
    return {
        "name": name,
        "path": str(path) if path else None,
        "content": content,
        "content_hash": content_hash(content),
        "metadata": doc,
        "validation": validation,
        "structured": {
            "title": doc.get("title") or name,
            "kind": doc.get("kind"),
            "entry_type": doc.get("entry_type"),
            "status": doc.get("status"),
            "summary": doc.get("summary"),
            "aliases": list(doc.get("aliases") or []),
            "sources": list(doc.get("sources") or []),
            "related_links": related_links,
            "canonical_sections": canonical_sections,
            "residual_markdown": residual_markdown,
            "validation_cues": {
                "valid": bool(validation.get("valid")) if isinstance(validation, dict) else None,
                "warnings": list(validation.get("warnings") or []) if isinstance(validation, dict) else [],
                "hard_failures": list(validation.get("hard_failures") or [])
                if isinstance(validation, dict)
                else [],
            },
        },
    }


def graph_payload(
    kb_path: str | Path,
    *,
    store: PlatformStore | None = None,
    graph_kind: str = "portal",
    focus: str | None = None,
    scene: str | None = None,
) -> dict[str, Any]:
    data = inventory(kb_path)
    cluster_stats = _query_cluster_stats(kb_path, store)
    proposals_by_id = {
        str(item["id"]): item for item in list_insight_proposals(kb_path)
    }
    clusters_by_id = {
        str(item["id"]): item for item in cluster_stats["clusters"]
    }
    scene_key = "admin" if graph_kind == "admin" else ("full" if str(scene or "").strip() == "full" else "home")
    budget = GRAPH_SCENE_BUDGETS[scene_key]
    anchor_specs = {
        "fresh": {"id": "anchor::fresh", "label": "Fresh knowledge basin", "x": -90, "y": 22, "z": 30},
        "forming": {"id": "anchor::forming", "label": "Forming knowledge basin", "x": 10, "y": 58, "z": -18},
        "stable": {"id": "anchor::stable", "label": "Stable knowledge basin", "x": 98, "y": -24, "z": 34},
    }
    graph_events = (
        store.list_graph_events(limit=600)
        if store is not None and hasattr(store, "list_graph_events")
        else []
    )
    selected_events: list[dict[str, Any]] = []
    for item in graph_events:
        score = _graph_event_score(item)
        if score <= 0.08:
            continue
        enriched = dict(item)
        details = enriched.get("details") or {}
        age_seconds = max(
            (datetime.now(timezone.utc) - _event_timestamp(str(enriched.get("created_at") or ""))).total_seconds(),
            0.0,
        )
        half_life_days = float(details.get("half_life_days") or 18.0)
        half_life_seconds = max(half_life_days, 1.0) * 24 * 60 * 60
        recentness = math.exp(-age_seconds / half_life_seconds)
        stability = float(details.get("stability") or 0.0)
        burst_level = max(
            0.08,
            min(
                1.0,
                float(details.get("burst_level") or 0.0)
                or score * 0.72
                + (0.24 if str(enriched.get("event_type") or "").startswith("ingest") else 0.0)
                + (0.16 if str(enriched.get("event_type") or "") == "proposal_materialized" else 0.0)
                - stability * 0.14,
            ),
        )
        formation_stage = str(details.get("formation_stage") or "").strip()
        if not formation_stage:
            event_type = str(enriched.get("event_type") or "")
            if event_type.startswith("ingest") and burst_level >= 0.7:
                formation_stage = "bursting"
            elif event_type in {"proposal_materialized", "ask_reinforced"}:
                formation_stage = "condensing"
            elif event_type in {"insight_promoted", "insight_merged"}:
                formation_stage = "stable"
            else:
                formation_stage = "stirring"
        enriched["_score"] = score
        enriched["_recentness"] = round(recentness, 4)
        enriched["_burst_level"] = round(burst_level, 4)
        enriched["_formation_stage"] = formation_stage
        selected_events.append(enriched)
    selected_events.sort(
        key=lambda item: (
            -float(item.get("_score") or 0.0),
            -_event_priority(str(item.get("event_type") or "")),
            str(item.get("created_at") or ""),
        )
    )
    selected_events = selected_events[: int(budget["events"])]

    nodes_by_id: dict[str, dict[str, Any]] = {}
    edges_by_key: dict[tuple[str, str, str], dict[str, Any]] = {}
    stage_priority = {
        "dormant": 0,
        "stable": 1,
        "stirring": 2,
        "condensing": 3,
        "bursting": 4,
    }

    def _stable_seed(value: str) -> int:
        return int(hashlib.sha1(value.encode("utf-8")).hexdigest()[:8], 16)

    def _position_near(anchor_id: str, node_id: str, spread: float = 42.0) -> tuple[float, float, float]:
        anchor = anchor_specs.get(anchor_id.replace("anchor::", ""), anchor_specs["stable"])
        seed = _stable_seed(node_id)
        dx = ((seed % 200) / 100 - 1) * spread
        dy = (((seed // 7) % 200) / 100 - 1) * spread * 0.55
        dz = (((seed // 13) % 200) / 100 - 1) * spread
        return (
            round(anchor["x"] + dx, 3),
            round(anchor["y"] + dy, 3),
            round(anchor["z"] + dz, 3),
        )

    def _anchor_for_event(event_type: str, visual_role: str = "") -> str:
        if event_type.startswith("ingest"):
            return "anchor::fresh"
        if event_type in {"proposal_materialized", "ask_reinforced"} or "forming" in visual_role:
            return "anchor::forming"
        return "anchor::stable"

    def _canonical_node(name: str) -> dict[str, Any]:
        doc = data["docs"].get(name) or {}
        return {
            "id": _graph_ref("canonical_entry", name),
            "label": doc.get("title") or name,
            "kind": "formal",
            "node_type": "canonical_entry",
            "entry_type": doc.get("entry_type") or "concept",
            "status": doc.get("status") or "fact",
            "state": "stable" if doc.get("status") in {"fact", "inferred"} else "soft",
            "summary": doc.get("summary") or "",
            "details": {
                "related_links": list(doc.get("graph_links") or []),
                "aliases": list(doc.get("aliases") or []),
                "sources": list(doc.get("sources") or []),
            },
            "entry_target": name,
        }

    def _proposal_node(proposal_id: str) -> dict[str, Any]:
        proposal = proposals_by_id.get(proposal_id) or {}
        return {
            "id": _graph_ref("insight_proposal", proposal_id),
            "label": proposal.get("title") or proposal_id,
            "kind": "insight",
            "node_type": "insight_proposal",
            "entry_type": proposal.get("kind") or "concept",
            "status": proposal.get("review_state") or "proposed",
            "state": proposal.get("review_state") or "proposed",
            "summary": proposal.get("proposed_answer") or proposal.get("hypothesis") or "",
            "details": {
                "hypothesis": proposal.get("hypothesis") or "",
                "proposed_answer": proposal.get("proposed_answer") or "",
                "supporting_entries": list(proposal.get("supporting_entries") or []),
                "trigger_queries": list(proposal.get("trigger_queries") or []),
            },
            "entry_target": None,
        }

    def _cluster_node(cluster_id: str) -> dict[str, Any]:
        cluster = clusters_by_id.get(cluster_id) or {}
        return {
            "id": _graph_ref("query_cluster", cluster_id),
            "label": cluster.get("display_query") or cluster.get("normalized_subject") or cluster_id,
            "kind": "signal_cluster",
            "node_type": "query_cluster",
            "entry_type": cluster.get("intent") or "definition",
            "status": cluster.get("status") or "captured",
            "state": cluster.get("status") or "captured",
            "summary": cluster.get("normalized_subject") or "",
            "details": {
                "intent": cluster.get("intent") or "",
                "source_entries": list(cluster.get("source_entries") or []),
                "demand_score": float(cluster.get("demand_score") or 0.0),
                "maturity_score": float(cluster.get("maturity_score") or 0.0),
            },
            "entry_target": None,
        }

    def _index_node(name: str, doc: dict[str, Any]) -> dict[str, Any]:
        return {
            "id": _graph_ref("index_segment", name),
            "label": doc.get("title") or name,
            "kind": "index",
            "node_type": "index_segment",
            "entry_type": "index",
            "status": "n/a",
            "state": "stable",
            "summary": doc.get("summary") or "",
            "details": {"segment": doc.get("segment") or "", "links": list(doc.get("links") or [])},
            "entry_target": None,
        }

    def ensure_anchor(anchor_id: str) -> None:
        anchor = anchor_specs.get(anchor_id.replace("anchor::", ""), anchor_specs["stable"])
        if anchor["id"] not in nodes_by_id:
            nodes_by_id[anchor["id"]] = {
                "id": anchor["id"],
                "label": anchor["label"],
                "kind": "anchor",
                "node_type": "cluster_anchor",
                "entry_type": "anchor",
                "status": "stable",
                "state": "stable",
                "summary": anchor["label"],
                "visual_role": "knowledge_basin",
                "event_type": "",
                "energy": 0.35,
                "stability": 1.0,
                "burst_level": 0.08,
                "formation_stage": "stable",
                "recentness": 0.12,
                "weight": 1.3,
                "x": anchor["x"],
                "y": anchor["y"],
                "z": anchor["z"],
                "fx": anchor["x"],
                "fy": anchor["y"],
                "fz": anchor["z"],
            }

    def ensure_node(
        node: dict[str, Any],
        *,
        event_type: str = "",
        visual_role: str = "",
        energy: float = 0.4,
        stability: float = 0.7,
        anchor_id: str = "anchor::stable",
        burst_level: float = 0.12,
        formation_stage: str = "stable",
        recentness: float = 0.2,
    ) -> None:
        node_id = str(node.get("id") or "")
        if not node_id:
            return
        ensure_anchor(anchor_id)
        existing = nodes_by_id.get(node_id)
        target = existing or dict(node)
        priority = _event_priority(event_type)
        current_priority = _event_priority(str(target.get("event_type") or ""))
        if not existing or priority >= current_priority:
            target["event_type"] = event_type
            target["visual_role"] = visual_role or target.get("visual_role") or "stable_canonical"
        target["energy"] = round(max(float(target.get("energy") or 0.0), energy), 4)
        target["stability"] = round(max(float(target.get("stability") or 0.0), stability), 4)
        target["burst_level"] = round(max(float(target.get("burst_level") or 0.0), burst_level), 4)
        target["recentness"] = round(max(float(target.get("recentness") or 0.0), recentness), 4)
        current_stage = str(target.get("formation_stage") or "stable")
        if stage_priority.get(formation_stage, 0) >= stage_priority.get(current_stage, 0):
            target["formation_stage"] = formation_stage
        target["weight"] = round(
            max(float(target.get("weight") or 1.0), 1.0 + target["energy"] * 1.9 + target["recentness"] * 0.7),
            4,
        )
        target["anchor_id"] = anchor_id
        if "x" not in target or existing is None:
            x, y, z = _position_near(anchor_id, node_id, spread=34.0 if target["node_type"] == "canonical_entry" else 26.0)
            target["x"], target["y"], target["z"] = x, y, z
        nodes_by_id[node_id] = target

    def add_edge(
        source: str,
        target: str,
        *,
        edge_type: str,
        strength: float,
        activation: float,
        formation_role: str,
        pulse_level: float = 0.0,
    ) -> None:
        if not source or not target or source == target:
            return
        key = tuple(sorted((source, target)) + [edge_type])  # type: ignore[list-item]
        existing = edges_by_key.get(key)
        if existing:
            existing["strength"] = round(max(float(existing.get("strength") or 0.0), strength), 4)
            existing["activation"] = round(max(float(existing.get("activation") or 0.0), activation), 4)
            existing["pulse_level"] = round(max(float(existing.get("pulse_level") or 0.0), pulse_level), 4)
            return
        edges_by_key[key] = {
            "source": source,
            "target": target,
            "kind": edge_type,
            "edge_type": edge_type,
            "strength": round(strength, 4),
            "activation": round(activation, 4),
            "formation_role": formation_role,
            "pulse_level": round(pulse_level, 4),
        }

    def ensure_reference(ref: str, *, energy: float, anchor_id: str, event_type: str, visual_role: str) -> str:
        kind, raw_id = _split_graph_ref(ref)
        if kind == "entry" and raw_id in data["docs"]:
            ensure_node(
                _canonical_node(raw_id),
                event_type=event_type,
                visual_role=visual_role,
                energy=energy,
                stability=0.88 if data["docs"][raw_id].get("status") in {"fact", "inferred"} else 0.66,
                anchor_id=anchor_id,
                burst_level=0.12,
                formation_stage="stable",
                recentness=0.24,
            )
            return _graph_ref("canonical_entry", raw_id)
        if kind == "insight" and raw_id in proposals_by_id:
            ensure_node(
                _proposal_node(raw_id),
                event_type=event_type,
                visual_role=visual_role,
                energy=energy,
                stability=0.48,
                anchor_id=anchor_id,
                burst_level=0.5,
                formation_stage="condensing",
                recentness=0.54,
            )
            return _graph_ref("insight_proposal", raw_id)
        if kind == "cluster" and raw_id in clusters_by_id:
            ensure_node(
                _cluster_node(raw_id),
                event_type=event_type,
                visual_role=visual_role,
                energy=energy,
                stability=float(clusters_by_id[raw_id].get("maturity_score") or 0.25),
                anchor_id=anchor_id,
                burst_level=0.42,
                formation_stage="stirring",
                recentness=0.46,
            )
            return _graph_ref("query_cluster", raw_id)
        if kind == "index" and raw_id in data["index_docs"]:
            ensure_node(
                _index_node(raw_id, data["index_docs"][raw_id]),
                event_type=event_type,
                visual_role="segment_context",
                energy=energy,
                stability=0.9,
                anchor_id="anchor::stable",
                burst_level=0.14,
                formation_stage="stable",
                recentness=0.22,
            )
            return _graph_ref("index_segment", raw_id)
        return ""

    for event in selected_events:
        event_type = str(event.get("event_type") or "")
        details = event.get("details") or {}
        visual_role = str(details.get("visual_role") or "").strip()
        anchor_id = _anchor_for_event(event_type, visual_role)
        subject_kind = str(event.get("subject_kind") or "")
        subject_id = str(event.get("subject_id") or "")
        subject_ref = _graph_ref(subject_kind, subject_id)
        score = float(event.get("_score") or 0.0)
        if subject_kind == "canonical_entry" and subject_id in data["docs"]:
            ensure_node(
                _canonical_node(subject_id),
                event_type=event_type,
                visual_role=visual_role or ("recent_canonical" if event_type.startswith("insight_") else "fresh_ingest"),
                energy=score,
                stability=float(details.get("stability") or 0.92),
                anchor_id=anchor_id,
                burst_level=float(event.get("_burst_level") or 0.2),
                formation_stage=str(event.get("_formation_stage") or "stable"),
                recentness=float(event.get("_recentness") or 0.2),
            )
        elif subject_kind == "insight_proposal" and subject_id in proposals_by_id:
            ensure_node(
                _proposal_node(subject_id),
                event_type=event_type,
                visual_role=visual_role or "forming_insight",
                energy=score,
                stability=float(details.get("stability") or 0.42),
                anchor_id=anchor_id,
                burst_level=float(event.get("_burst_level") or 0.52),
                formation_stage=str(event.get("_formation_stage") or "condensing"),
                recentness=float(event.get("_recentness") or 0.58),
            )
        elif subject_kind == "query_cluster" and subject_id in clusters_by_id:
            ensure_node(
                _cluster_node(subject_id),
                event_type=event_type,
                visual_role=visual_role or "reinforced_query",
                energy=score,
                stability=float(details.get("stability") or 0.3),
                anchor_id=anchor_id,
                burst_level=float(event.get("_burst_level") or 0.44),
                formation_stage=str(event.get("_formation_stage") or "stirring"),
                recentness=float(event.get("_recentness") or 0.48),
            )
        elif subject_kind == "index_segment" and subject_id in data["index_docs"]:
            ensure_node(
                _index_node(subject_id, data["index_docs"][subject_id]),
                event_type=event_type,
                visual_role="segment_context",
                energy=score,
                stability=0.9,
                anchor_id="anchor::stable",
                burst_level=0.12,
                formation_stage="stable",
                recentness=float(event.get("_recentness") or 0.22),
            )
        ensure_anchor(anchor_id)
        if subject_ref:
            add_edge(
                anchor_id,
                subject_ref,
                edge_type="belongs_to_cluster",
                strength=max(0.35, score),
                activation=max(0.25, score),
                formation_role="knowledge_basin",
                pulse_level=max(0.12, float(event.get("_burst_level") or 0.0) * 0.58),
            )

        related_refs = [str(item).strip() for item in event.get("related_ids") or [] if str(item).strip()]
        for ref in related_refs:
            related_ref = ensure_reference(
                ref,
                energy=max(0.22, score * 0.68),
                anchor_id=anchor_id,
                event_type=event_type,
                visual_role="supporting_entry" if ref.startswith("entry::") else "forming_context",
            )
            if not related_ref or not subject_ref:
                continue
            edge_type = "weak_affinity"
            formation_role = "story_context"
            pulse_level = max(0.08, float(event.get("_recentness") or 0.0) * 0.42)
            if event_type == "ask_reinforced":
                edge_type = "ask_reinforcement"
                formation_role = "reinforcement"
                pulse_level = max(0.42, float(event.get("_burst_level") or 0.0))
            elif event_type == "proposal_materialized":
                edge_type = "supports"
                formation_role = "formation_support"
                pulse_level = max(0.46, float(event.get("_burst_level") or 0.0))
            elif event_type in {"insight_promoted", "insight_merged"} and ref.startswith("insight::"):
                edge_type = "routes_to"
                formation_role = "canonicalization"
                pulse_level = max(0.52, float(event.get("_burst_level") or 0.0))
            elif event_type.startswith("ingest"):
                edge_type = "weak_affinity"
                formation_role = "ingest_context"
                pulse_level = max(0.28, float(event.get("_burst_level") or 0.0) * 0.7)
            add_edge(
                subject_ref,
                related_ref,
                edge_type=edge_type,
                strength=max(0.22, score * 0.9),
                activation=max(0.18, score),
                formation_role=formation_role,
                pulse_level=pulse_level,
            )

    def _fill_stable_context(minimum_visible: int = 9) -> None:
        ensure_anchor("anchor::stable")
        visible_entry_refs = {
            node_id
            for node_id, node in nodes_by_id.items()
            if node.get("node_type") == "canonical_entry"
        }
        top_entries = sorted(
            [
                (name, doc)
                for name, doc in data["docs"].items()
                if doc["kind"] == "formal"
            ],
            key=lambda item: (-int(item[1].get("inbound_count") or 0), item[0]),
        )[:10]
        for name, doc in top_entries:
            entry_ref = _graph_ref("canonical_entry", name)
            if entry_ref in visible_entry_refs and len(visible_entry_refs) >= minimum_visible:
                continue
            ensure_node(
                _canonical_node(name),
                event_type="",
                visual_role="stable_canonical",
                energy=max(0.24, 0.15 + float(doc.get("inbound_count") or 0.0) * 0.07),
                stability=0.94,
                anchor_id="anchor::stable",
                burst_level=0.1,
                formation_stage="stable",
                recentness=0.18,
            )
            visible_entry_refs.add(entry_ref)
            add_edge(
                "anchor::stable",
                entry_ref,
                edge_type="belongs_to_cluster",
                strength=0.3,
                activation=0.18,
                formation_role="knowledge_basin",
                pulse_level=0.08,
            )
        for name, doc in top_entries:
            for target in doc.get("graph_links") or []:
                source_ref = _graph_ref("canonical_entry", name)
                target_ref = _graph_ref("canonical_entry", target)
                if source_ref in visible_entry_refs and target_ref in visible_entry_refs:
                    add_edge(
                        source_ref,
                        target_ref,
                        edge_type="weak_affinity",
                        strength=0.26,
                        activation=0.16,
                        formation_role="stable_context",
                        pulse_level=0.06,
                    )

    if not selected_events:
        _fill_stable_context()
    elif len([node for node in nodes_by_id.values() if node.get("node_type") != "cluster_anchor"]) < 8:
        _fill_stable_context()

    visible_doc_names = {
        raw_id
        for node_id in nodes_by_id
        for kind, raw_id in [_split_graph_ref(node_id)]
        if kind == "entry" and raw_id in data["docs"]
    }
    for name in sorted(visible_doc_names):
        doc = data["docs"][name]
        for target in doc.get("graph_links") or []:
            if target in visible_doc_names:
                add_edge(
                    _graph_ref("canonical_entry", name),
                    _graph_ref("canonical_entry", target),
                    edge_type="weak_affinity",
                    strength=0.24,
                    activation=0.18,
                    formation_role="knowledge_context",
                    pulse_level=0.08,
                )

    if scene_key == "admin":
        index_candidates = [
            (name, doc)
            for name, doc in data["index_docs"].items()
            if not doc.get("is_root")
            and any(link in visible_doc_names for link in doc.get("links") or [])
        ][:6]
        for name, doc in index_candidates:
            index_ref = ensure_reference(
                _graph_ref("index_segment", name),
                energy=0.22,
                anchor_id="anchor::stable",
                event_type="",
                visual_role="segment_context",
            )
            for link in doc.get("links") or []:
                if link in visible_doc_names:
                    add_edge(
                        index_ref,
                        _graph_ref("canonical_entry", link),
                        edge_type="belongs_to_cluster",
                        strength=0.3,
                        activation=0.2,
                        formation_role="index_projection",
                        pulse_level=0.08,
                    )

    nodes = list(nodes_by_id.values())
    edges = list(edges_by_key.values())

    if focus:
        focus_id = focus
        for candidate in (
            focus,
            _graph_ref("canonical_entry", focus),
            _graph_ref("insight_proposal", focus),
            _graph_ref("query_cluster", focus),
        ):
            if any(node["id"] == candidate for node in nodes):
                focus_id = candidate
                break
        related_ids = {focus_id}
        frontier = {focus_id}
        for _ in range(2):
            next_frontier: set[str] = set()
            for edge in edges:
                if edge["source"] in frontier or edge["target"] in frontier:
                    related_ids.add(edge["source"])
                    related_ids.add(edge["target"])
                    next_frontier.add(edge["source"])
                    next_frontier.add(edge["target"])
            frontier = next_frontier
        nodes = [node for node in nodes if node["id"] in related_ids]
        edges = [
            edge for edge in edges if edge["source"] in related_ids and edge["target"] in related_ids
        ]
    else:
        nodes = sorted(
            nodes,
            key=lambda item: (
                0 if item["node_type"] == "cluster_anchor" else 1,
                -float(item.get("energy") or 0.0),
                -float(item.get("weight") or 0.0),
                item["id"],
            ),
        )[: int(budget["nodes"])]
        visible = {node["id"] for node in nodes}
        edges = [
            edge
            for edge in sorted(
                edges,
                key=lambda item: (
                    -float(item.get("activation") or 0.0),
                    -float(item.get("strength") or 0.0),
                    item["source"],
                    item["target"],
                ),
            )
            if edge["source"] in visible and edge["target"] in visible
        ][: int(budget["edges"])]

    focus_seed = next(
        (node["id"] for node in nodes if node["node_type"] != "cluster_anchor"),
        nodes[0]["id"] if nodes else "",
    )
    story_caption = ""
    if selected_events:
        lead = selected_events[0]
        story_caption = str(
            (lead.get("details") or {}).get("proposed_answer")
            or lead.get("subject_label")
            or ""
        ).strip()
    elif nodes:
        story_caption = str(
            next((node.get("label") for node in nodes if node["node_type"] != "cluster_anchor"), "")
            or nodes[0].get("label")
            or ""
        ).strip()

    return {
        "graph_version": "insights-v2",
        "graph_kind": graph_kind,
        "scene_mode": budget["scene_mode"],
        "focus_seed": focus_seed,
        "story_caption": story_caption,
        "ambient_seed": hashlib.sha1(
            f"{kb_path}:{scene_key}:{','.join(str(item.get('id') or '') for item in selected_events[:6])}".encode("utf-8")
        ).hexdigest()[:12],
        "playback_events": [
            {
                "id": str(item.get("id") or f"event-{index}"),
                "event_type": str(item.get("event_type") or ""),
                "subject_ref": _graph_ref(str(item.get("subject_kind") or ""), str(item.get("subject_id") or "")),
                "related_refs": [str(ref) for ref in item.get("related_ids") or [] if str(ref).strip()],
                "caption": str(
                    (item.get("details") or {}).get("proposed_answer")
                    or item.get("subject_label")
                    or ""
                ).strip(),
            }
            for index, item in enumerate(selected_events[: min(6, len(selected_events))])
        ],
        "kb_language": data.get("default_language") or "en",
        "generated_at": utc_now(),
        "stats": {
            "node_count": len(nodes),
            "edge_count": len(edges),
            "formal_entry_count": len(data["entries"]),
            "insight_count": len(data.get("insights") or []),
            "query_cluster_count": len(cluster_stats["clusters"]),
            "cluster_coverage": cluster_stats["cluster_coverage"],
            "event_count": len(selected_events),
        },
        "nodes": nodes,
        "edges": edges,
    }


def summarize_health_report(report: dict[str, Any], *, cluster_coverage: float | None = None) -> dict[str, Any]:
    return {
        "formal_entry_count": report["formal_entry_count"],
        "placeholder_count": report["placeholder_count"],
        "hard_fail_entry_count": report["hard_fail_entry_count"],
        "dangling_link_count": report["dangling_link_count"],
        "orphan_entry_count": report["orphan_entry_count"],
        "promotable_placeholder_count": report["promotable_placeholder_count"],
        "canonical_gap_count": report["canonical_gap_count"],
        "invalid_index_count": report["invalid_index_count"],
        "cluster_coverage": cluster_coverage if cluster_coverage is not None else 0.0,
    }


def build_health_issue_queue(
    kb_path: str | Path,
    *,
    store: PlatformStore | None = None,
) -> list[dict[str, Any]]:
    report = audit_kb(kb_path)
    data = inventory(kb_path)
    cluster_stats = _query_cluster_stats(kb_path, store)
    issues: list[dict[str, Any]] = []

    for item in report["entry_validation"]:
        if item["hard_failures"]:
            issues.append(
                issue(
                    issue_type="hard_failure",
                    severity="blocking",
                    target=item["name"],
                    summary="; ".join(item["hard_failures"]),
                    suggested_action="edit_entry",
                    evidence={"hard_failures": item["hard_failures"]},
                )
            )

    for link in report["dangling_links"]:
        issues.append(
            issue(
                issue_type="dangling_link",
                severity="high",
                target=link["source_file"],
                summary=f"链接 {link['link']} 没有目标",
                suggested_action="run_tidy",
                evidence=link,
            )
        )

    for orphan in report["orphan_entries"]:
        issues.append(
            issue(
                issue_type="orphan_entry",
                severity="medium",
                target=orphan,
                summary="正式条目缺少足够连接",
                suggested_action="run_tidy",
                evidence={"entry": orphan},
            )
        )

    for item in report["promotable_placeholders"]:
        issues.append(
            issue(
                issue_type="promotable_placeholder",
                severity="medium",
                target=item["name"],
                summary=f"placeholder 被引用 {item['ref_count']} 次，值得提升",
                suggested_action="promote_placeholder",
                evidence=item,
            )
        )

    for item in report["canonical_gaps"]:
        issues.append(
            issue(
                issue_type="canonical_gap",
                severity="high",
                target=item["name"],
                summary="正式条目依赖 placeholder 作为关键概念",
                suggested_action="run_tidy",
                evidence=item,
            )
        )

    for item in report["provenance_contamination"]:
        issues.append(
            issue(
                issue_type="provenance_contamination",
                severity="medium",
                target=item["name"],
                summary="来源型链接污染了知识图谱",
                suggested_action="edit_entry",
                evidence=item,
            )
        )

    for item in report["index_validation"]:
        if item["hard_failures"]:
            issues.append(
                issue(
                    issue_type="invalid_index",
                    severity="blocking",
                    target=item["name"],
                    summary="; ".join(item["hard_failures"]),
                    suggested_action="edit_entry",
                    evidence=item,
                )
            )

    for item in report["overloaded_indexes"]:
        issues.append(
            issue(
                issue_type="overloaded_index",
                severity="medium",
                target=item["name"],
                summary="索引规模超过阈值",
                suggested_action="run_tidy",
                evidence=item,
            )
        )

    for item in report["unknown_index_links"]:
        issues.append(
            issue(
                issue_type="unknown_index_link",
                severity="high",
                target=item["index"],
                summary=f"索引链接 {item['link']} 不存在或不应作为入口",
                suggested_action="edit_entry",
                evidence=item,
            )
        )

    for name in report["uncovered_formal_entries"]:
        issues.append(
            issue(
                issue_type="uncovered_formal_entry",
                severity="medium",
                target=name,
                summary="正式条目未被任何索引覆盖",
                suggested_action="run_tidy",
                evidence={"entry": name},
            )
        )

    for name, doc in data["docs"].items():
        if doc["kind"] != "formal":
            continue
        if doc["status"] in {"inferred", "disputed"}:
            issues.append(
                issue(
                    issue_type="low_confidence_entry",
                    severity="low" if doc["status"] == "inferred" else "medium",
                    target=name,
                    summary=f"条目状态为 {doc['status']}，需要人工关注",
                    suggested_action="review_entry",
                    evidence={"status": doc["status"]},
                )
            )

    root_doc = data["index_docs"].get("index.root")
    if root_doc and len(data["entries"]) >= 8 and len(root_doc.get("links") or []) <= 1:
        issues.append(
            issue(
                issue_type="root_index_skeleton",
                severity="medium",
                target="index.root",
                summary="根索引仍接近空骨架，无法承载稳定入口聚类",
                suggested_action="run_tidy",
                evidence={"root_links": list(root_doc.get("links") or []), "formal_entry_count": len(data["entries"])},
            )
        )

    non_root_indexes = [item for item in data["index_docs"].values() if not item.get("is_root")]
    if len(data["entries"]) >= 12 and not non_root_indexes:
        issues.append(
            issue(
                issue_type="missing_segment_index",
                severity="medium",
                target="indexes",
                summary="正式条目规模已增长，但缺少分段 index 作为稳定聚类入口",
                suggested_action="run_tidy",
                evidence={"formal_entry_count": len(data["entries"])},
            )
        )

    cluster_coverage = cluster_stats["cluster_coverage"]
    if len(data["entries"]) >= 6 and cluster_coverage < 0.5:
        issues.append(
            issue(
                issue_type="low_cluster_coverage",
                severity="medium",
                target="knowledge-base",
                summary="知识聚类覆盖率偏低，索引网络尚未形成足够稳定的知识盆地",
                suggested_action="run_tidy",
                evidence={"cluster_coverage": cluster_coverage},
            )
        )

    for cluster in cluster_stats["ready_clusters"][:10]:
        issues.append(
            issue(
                issue_type="latent_cluster_ready",
                severity="low",
                target=str(cluster.get("normalized_subject") or cluster.get("display_query") or cluster.get("cluster_key") or "cluster"),
                summary="隐性知识簇已经成熟，建议物化为 insight proposal 或进入整理队列",
                suggested_action="run_tidy",
                evidence=cluster,
            )
        )

    if non_root_indexes and cluster_coverage < 0.75:
        issues.append(
            issue(
                issue_type="index_cluster_drift",
                severity="low",
                target="indexes",
                summary="现有 index 入口与实际高频知识簇出现偏移，建议重构 segment index",
                suggested_action="run_tidy",
                evidence={"cluster_coverage": cluster_coverage, "index_count": len(non_root_indexes)},
            )
        )

    return sorted(issues, key=lambda item: (severity_rank(item["severity"]), item["target"]))


def issue(
    *,
    issue_type: str,
    severity: str,
    target: str,
    summary: str,
    suggested_action: str,
    evidence: dict[str, Any],
) -> dict[str, Any]:
    return {
        "type": issue_type,
        "severity": severity,
        "target": target,
        "summary": summary,
        "suggested_action": suggested_action,
        "evidence": evidence,
    }


def severity_rank(severity: str) -> int:
    return {"blocking": 0, "high": 1, "medium": 2, "low": 3}.get(severity, 4)


def get_health_payload(
    kb_path: str | Path,
    *,
    store: PlatformStore | None = None,
) -> dict[str, Any]:
    report = audit_kb(kb_path)
    cluster_stats = _query_cluster_stats(kb_path, store)
    issues = build_health_issue_queue(kb_path, store=store)
    counts: dict[str, int] = {}
    for item in issues:
        counts[item["severity"]] = counts.get(item["severity"], 0) + 1
    return {
        "summary": summarize_health_report(report, cluster_coverage=cluster_stats["cluster_coverage"]),
        "severity_counts": counts,
        "issues": issues,
        "emerging_clusters": cluster_stats["emerging"],
        "canonical_stress_points": cluster_stats["stress_points"],
    }


def _issue_document_name(
    issue: dict[str, Any],
    *,
    docs: dict[str, dict[str, Any]],
    index_docs: dict[str, dict[str, Any]],
) -> str | None:
    evidence = issue.get("evidence") if isinstance(issue.get("evidence"), dict) else {}
    candidates = [
        issue.get("target"),
        evidence.get("name"),
        evidence.get("entry"),
        evidence.get("index"),
        evidence.get("source_file"),
    ]
    for raw_value in candidates:
        value = str(raw_value or "").strip()
        if not value:
            continue
        normalized = Path(value).stem if "/" in value or value.endswith(".md") else value
        if normalized in docs or normalized in index_docs:
            return normalized
    return None


def _admin_document_collections(
    kb_path: str | Path,
) -> tuple[
    dict[str, Any],
    list[dict[str, Any]],
    dict[str, list[dict[str, Any]]],
    dict[str, dict[str, Any]],
]:
    root = Path(kb_path).resolve()
    data = inventory(root)
    docs = data["docs"]
    index_docs = data["index_docs"]
    issues = build_health_issue_queue(root)
    issues_by_document: dict[str, list[dict[str, Any]]] = defaultdict(list)
    normalized_issues: list[dict[str, Any]] = []
    for issue in issues:
        enriched = dict(issue)
        document_name = _issue_document_name(issue, docs=docs, index_docs=index_docs)
        enriched["document_name"] = document_name
        normalized_issues.append(enriched)
        if document_name:
            issues_by_document[document_name].append(enriched)

    def build_document(name: str, doc: dict[str, Any], *, group: str) -> dict[str, Any]:
        path = resolve_kb_document_path(root, name)
        document_issues = issues_by_document.get(name, [])
        relative_path = None
        updated_at = 0
        if path is not None and path.exists():
            relative_path = str(path.resolve().relative_to(root))
            updated_at = int(path.stat().st_mtime)
        return {
            "name": name,
            "title": str(doc.get("title") or name),
            "group": group,
            "kind": str(doc.get("kind") or group),
            "entry_type": str(doc.get("entry_type") or ("index" if group == "index" else "")),
            "status": str(doc.get("status") or ("n/a" if group == "index" else "")),
            "summary": str(doc.get("summary") or ""),
            "relative_path": relative_path,
            "updated_at": updated_at,
            "aliases": list(doc.get("aliases") or []),
            "links": list(doc.get("graph_links") or doc.get("links") or []),
            "issue_count": len(document_issues),
            "issue_summaries": [str(item.get("summary") or "") for item in document_issues[:3]],
            "indexes": [],
        }

    documents_by_name: dict[str, dict[str, Any]] = {}
    for name in data["entries"]:
        documents_by_name[name] = build_document(name, docs[name], group="formal")
    for name in data["placeholders"]:
        documents_by_name[name] = build_document(name, docs[name], group="placeholder")
    for name in data["indexes"]:
        documents_by_name[name] = build_document(name, index_docs[name], group="index")
    return data, normalized_issues, issues_by_document, documents_by_name


def kb_document_browser_payload(kb_path: str | Path) -> dict[str, Any]:
    data, normalized_issues, _issues_by_document, documents_by_name = _admin_document_collections(kb_path)
    groups = {
        "formal": [documents_by_name[name] for name in data["entries"]],
        "placeholder": [documents_by_name[name] for name in data["placeholders"]],
        "index": [documents_by_name[name] for name in data["indexes"]],
    }
    return {
        "counts": {key: len(value) for key, value in groups.items()},
        "groups": groups,
        "health_issues": normalized_issues,
    }


def kb_file_management_payload(kb_path: str | Path) -> dict[str, Any]:
    data, normalized_issues, _issues_by_document, documents_by_name = _admin_document_collections(kb_path)
    index_docs = data["index_docs"]
    index_names = set(data["indexes"])
    document_names = set(documents_by_name) - index_names
    parent_indexes: dict[str, set[str]] = defaultdict(set)
    for name, record in index_docs.items():
        for target in record.get("links") or []:
            if target in index_names and target != name:
                parent_indexes[target].add(name)

    indexed_documents: dict[str, set[str]] = defaultdict(set)

    def build_index_node(name: str, *, stack: tuple[str, ...] = ()) -> tuple[dict[str, Any], set[str]]:
        record = documents_by_name[name]
        direct_doc_targets: list[str] = []
        child_index_targets: list[str] = []
        for target in record.get("links") or []:
            if target in index_names and target != name:
                child_index_targets.append(target)
            elif target in document_names:
                direct_doc_targets.append(target)
        direct_doc_targets = list(dict.fromkeys(direct_doc_targets))
        child_index_targets = [
            target for target in dict.fromkeys(child_index_targets) if target not in stack
        ]

        child_nodes: list[dict[str, Any]] = []
        reachable_documents: set[str] = set(direct_doc_targets)
        for child_name in child_index_targets:
            child_node, child_documents = build_index_node(child_name, stack=(*stack, name))
            child_nodes.append(child_node)
            reachable_documents.update(child_documents)
        for document_name in reachable_documents:
            indexed_documents[document_name].add(name)
        node = {
            "name": record["name"],
            "title": record["title"],
            "summary": record["summary"],
            "relative_path": record["relative_path"],
            "issue_count": record["issue_count"],
            "segment": str(index_docs[name].get("segment") or record["name"]),
            "is_root": bool(index_docs[name].get("is_root")),
            "entry_count": int(index_docs[name].get("entry_count") or 0),
            "estimated_tokens": int(index_docs[name].get("estimated_tokens") or 0),
            "last_tidied_at": str(index_docs[name].get("last_tidied_at") or ""),
            "direct_documents": [documents_by_name[target] for target in direct_doc_targets],
            "child_indexes": child_nodes,
            "reachable_document_count": len(reachable_documents),
        }
        return node, reachable_documents

    root_indexes = [
        name
        for name in data["indexes"]
        if bool(index_docs[name].get("is_root")) or not parent_indexes.get(name)
    ]
    top_index_names = list(dict.fromkeys(root_indexes or data["indexes"]))
    top_index_names.sort(
        key=lambda name: (
            not bool(index_docs[name].get("is_root")),
            str(documents_by_name[name].get("title") or name).casefold(),
        )
    )
    top_indexes: list[dict[str, Any]] = []
    reachable_from_top: set[str] = set()
    for name in top_index_names:
        node, reachable_documents = build_index_node(name)
        top_indexes.append(node)
        reachable_from_top.update(reachable_documents)

    for document_name, index_set in indexed_documents.items():
        if document_name in documents_by_name:
            documents_by_name[document_name]["indexes"] = sorted(index_set)

    unindexed_documents = [
        documents_by_name[name]
        for name in sorted(document_names - reachable_from_top)
    ]
    return {
        "counts": {
            "formal": len(data["entries"]),
            "placeholder": len(data["placeholders"]),
            "index": len(data["indexes"]),
            "indexed": len(reachable_from_top),
            "unindexed": len(unindexed_documents),
        },
        "documents_by_name": documents_by_name,
        "top_indexes": top_indexes,
        "unindexed_documents": unindexed_documents,
        "health_issues": normalized_issues,
    }


def search_kb_file_suggestions(
    kb_path: str | Path,
    query: str,
    *,
    limit: int = 8,
) -> list[dict[str, Any]]:
    raw_query = query.strip()
    if not raw_query:
        return []
    payload = kb_file_management_payload(kb_path)
    terms = [term for term in re.split(r"\s+", raw_query) if term]
    suggestions: list[dict[str, Any]] = []
    for document in payload["documents_by_name"].values():
        score = 0
        matched_field = ""
        aliases = list(document.get("aliases") or [])
        relative_path = str(document.get("relative_path") or "")
        summary = str(document.get("summary") or "")
        title = str(document.get("title") or "")
        name = str(document.get("name") or "")
        for term in terms:
            lowered = term.casefold()
            if title.casefold().startswith(lowered) or name.casefold().startswith(lowered):
                score += 16
                matched_field = matched_field or "title"
            elif lowered in title.casefold():
                score += 12
                matched_field = matched_field or "title"
            elif any(alias.casefold().startswith(lowered) for alias in aliases):
                score += 10
                matched_field = matched_field or "aliases"
            elif lowered in " ".join(aliases).casefold():
                score += 8
                matched_field = matched_field or "aliases"
            elif lowered in relative_path.casefold():
                score += 6
                matched_field = matched_field or "path"
            elif lowered in summary.casefold():
                score += 4
                matched_field = matched_field or "summary"
        if score <= 0:
            continue
        suggestions.append(
            {
                "name": name,
                "title": title,
                "group": str(document.get("group") or ""),
                "kind": str(document.get("kind") or ""),
                "entry_type": str(document.get("entry_type") or ""),
                "status": str(document.get("status") or ""),
                "summary": summary,
                "relative_path": relative_path,
                "indexes": list(document.get("indexes") or []),
                "matched_field": matched_field or "summary",
                "score": score + min(int(document.get("issue_count") or 0), 5),
            }
        )
    return sorted(suggestions, key=lambda item: (-item["score"], item["name"]))[:limit]


def content_hash(content: str) -> str:
    return hashlib.sha256(content.encode("utf-8")).hexdigest()


def list_reviews_with_jobs(
    *,
    store: PlatformStore,
    decision: str | None = None,
) -> list[dict[str, Any]]:
    reviews = store.list_reviews(decision=decision, limit=200)
    enriched: list[dict[str, Any]] = []
    for review in reviews:
        payload = dict(review)
        payload["job"] = store.get_job(review["job_id"])
        if review.get("submission_id"):
            payload["submission"] = store.get_submission(review["submission_id"])
        enriched.append(payload)
    return enriched


def determine_target_path(
    kb_path: str | Path,
    *,
    name: str,
    content: str,
    relative_path: str | None = None,
) -> Path:
    root = Path(kb_path).resolve()
    if relative_path:
        target = (root / relative_path).resolve()
        if target == root or root not in target.parents:
            raise ValueError("relative path escapes knowledge base root")
        if target.exists() and target.is_dir():
            raise ValueError("relative path must point to a file inside the knowledge base")
        return target

    existing = resolve_kb_document_path(root, name)
    if existing is not None:
        return existing.resolve()

    frontmatter, _ = split_frontmatter(content)
    item_type = str(frontmatter.get("type", "")).strip()
    if item_type == "placeholder":
        return root / "placeholders" / f"{name}.md"
    if item_type in {"concept", "lesson"}:
        return root / "entries" / f"{name}.md"
    if str(frontmatter.get("kind", "")).strip() == "index" or name.startswith("index."):
        if name == "index.root":
            return root / index_config()["root_file"]
        return root / "indexes" / f"{name}.md"
    return root / "entries" / f"{name}.md"


def validate_target_content(path: Path, content: str) -> dict[str, Any] | None:
    if path.name == index_config()["root_file"] or path.parent.name == "indexes":
        temp_path = _write_temp_for_validation(path.name, content)
        try:
            return validate_index(temp_path)
        finally:
            temp_path.unlink(missing_ok=True)
    if path.parent.name == "insights":
        return validate_insight_content(path, content)
    kind = "placeholder" if path.parent.name == "placeholders" else "formal"
    return validate_entry(text=content, name=path.stem, kind=kind)


def _write_temp_for_validation(filename: str, content: str) -> Path:
    tmp = tempfile.NamedTemporaryFile(
        prefix="sediment-validate-",
        suffix=f"-{filename}",
        delete=False,
    )
    tmp.write(content.encode("utf-8"))
    tmp.flush()
    tmp.close()
    return Path(tmp.name)


def build_diff(path_label: str, old_content: str, new_content: str) -> str:
    diff = difflib.unified_diff(
        old_content.splitlines(keepends=True),
        new_content.splitlines(keepends=True),
        fromfile=f"a/{path_label}",
        tofile=f"b/{path_label}",
    )
    return "".join(diff)


def _atomic_write_text(path: Path, content: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    temp = tempfile.NamedTemporaryFile(
        prefix=f".{path.name}.",
        suffix=".tmp",
        dir=path.parent,
        delete=False,
    )
    temp_path = Path(temp.name)
    try:
        temp.write(content.encode("utf-8"))
        temp.flush()
        temp.close()
        temp_path.replace(path)
    except Exception:
        temp.close()
        temp_path.unlink(missing_ok=True)
        raise


def apply_operations(
    kb_path: str | Path,
    operations: list[dict[str, Any]],
    *,
    actor_name: str,
    actor_id: str | None = None,
    actor_role: str,
    store: PlatformStore,
) -> dict[str, Any]:
    root = Path(kb_path).resolve()
    prepared: list[dict[str, Any]] = []
    seen_paths: set[str] = set()
    for operation in operations:
        name = operation.get("name") or Path(operation.get("relative_path", "")).stem
        target_path = determine_target_path(
            root,
            name=name,
            content=operation.get("content", ""),
            relative_path=operation.get("relative_path"),
        )
        old_content = target_path.read_text(encoding="utf-8") if target_path.exists() else ""
        target_exists = target_path.exists()
        current_hash = content_hash(old_content)
        base_hash = operation.get("base_hash")
        if base_hash and target_path.exists() and current_hash != base_hash:
            raise RuntimeError(
                f"{target_path.name} changed since review was prepared; refresh the diff first."
            )
        rel_path = str(target_path.relative_to(root))
        if rel_path in seen_paths:
            raise ValueError(f"duplicate operation target detected: {rel_path}")
        seen_paths.add(rel_path)

        change_type = operation.get("change_type", "update")
        if change_type == "delete":
            prepared.append(
                {
                    "name": name,
                    "target_path": target_path,
                    "relative_path": rel_path,
                    "change_type": "delete",
                    "target_exists": target_exists,
                    "old_content": old_content,
                }
            )
            continue

        new_content = operation["content"].strip() + "\n"
        validation = validate_target_content(target_path, new_content)
        if validation and validation.get("hard_failures"):
            raise ValueError("; ".join(validation["hard_failures"]))
        prepared.append(
            {
                "name": name,
                "target_path": target_path,
                "relative_path": rel_path,
                "change_type": change_type,
                "target_exists": target_exists,
                "old_content": old_content,
                "new_content": new_content,
            }
        )

    applied: list[dict[str, Any]] = []
    mutated: list[dict[str, Any]] = []
    try:
        for item in prepared:
            target_path = item["target_path"]
            rel_path = item["relative_path"]
            mutated.append(item)
            if item["change_type"] == "delete":
                if target_path.exists():
                    target_path.unlink()
                applied.append(
                    {
                        "name": item["name"],
                        "relative_path": rel_path,
                        "change_type": "delete",
                        "diff": build_diff(rel_path, item["old_content"], ""),
                    }
                )
                continue

            new_content = item["new_content"]
            _atomic_write_text(target_path, new_content)
            applied.append(
                {
                    "name": item["name"],
                    "relative_path": rel_path,
                    "change_type": item["change_type"],
                    "diff": build_diff(rel_path, item["old_content"], new_content),
                    "new_hash": content_hash(new_content),
                }
            )
    except Exception as exc:
        rollback_errors: list[str] = []
        for item in reversed(mutated):
            target_path = item["target_path"]
            try:
                if item["target_exists"]:
                    _atomic_write_text(target_path, item["old_content"])
                else:
                    target_path.unlink(missing_ok=True)
            except Exception as rollback_exc:  # noqa: BLE001
                rollback_errors.append(
                    f"{target_path}: {rollback_exc}"
                )
        if rollback_errors:
            raise RuntimeError(
                f"{exc}; rollback failed for {', '.join(rollback_errors)}"
            ) from exc
        raise

    store.add_audit_log(
        actor_name=actor_name,
        actor_id=actor_id,
        actor_role=actor_role,
        action="kb.apply_operations",
        target_type="knowledge_base",
        target_id=str(root),
        details={"operations": applied},
    )
    return {"operations": applied, "health": summarize_health_report(audit_kb(root))}


def save_entry(
    kb_path: str | Path,
    *,
    name: str,
    content: str,
    expected_hash: str | None,
    actor_name: str,
    actor_id: str | None = None,
    actor_role: str = "committer",
    store: PlatformStore,
) -> dict[str, Any]:
    root = Path(kb_path).resolve()
    existing = resolve_kb_document_path(root, name)
    if existing is None:
        raise FileNotFoundError(name)
    target = existing.resolve()
    old_content = target.read_text(encoding="utf-8")
    if expected_hash and old_content and content_hash(old_content) != expected_hash:
        raise RuntimeError("entry changed since it was loaded; refresh before saving")
    result = apply_operations(
        root,
        [
            {
                "name": name,
                "relative_path": str(target.relative_to(root)),
                "change_type": "update",
                "content": content,
                "base_hash": content_hash(old_content),
            }
        ],
        actor_name=actor_name,
        actor_id=actor_id,
        actor_role=actor_role,
        store=store,
    )
    detail = get_entry_detail(root, name)
    detail["write_result"] = result
    detail["write_result"]["actor"] = {
        "id": actor_id,
        "name": actor_name,
        "role": actor_role,
    }
    return detail


def stage_workspace_copy(kb_path: str | Path, workspaces_dir: str | Path, job_id: str) -> Path:
    workspaces_root = Path(workspaces_dir)
    workspaces_root.mkdir(parents=True, exist_ok=True)
    workspace = workspaces_root / job_id
    if workspace.exists():
        shutil.rmtree(workspace)
    shutil.copytree(kb_path, workspace / "knowledge-base")
    return workspace
