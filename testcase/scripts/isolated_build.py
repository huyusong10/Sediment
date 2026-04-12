"""
isolated_build.py — Sediment 隔离构建脚本

封装项目隔离、KB 构建、MCP Server 启动的完整流程。
所有操作均在隔离目录内执行，确保不影响源码目录。

用法：
    # 完整构建并启动测试
    python isolated_build.py --build-type full
    python isolated_build.py --build-type batched

    # 仅创建隔离目录（不构建）
    python isolated_build.py --dry-run

    # 作为模块导入
    from isolated_build import IsolatedBuilder
    async with IsolatedBuilder() as builder:
        await builder.build_full()
        server = builder.start_mcp_server()
"""

import argparse
import asyncio
import json
import os
import random
import re
import shutil
import sys
import tempfile
import time
from pathlib import Path

from mcp_server.llm_cli import build_cli_command

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

PROJECT_ROOT = Path(__file__).resolve().parent.parent.parent
TESTCASE_DIR = PROJECT_ROOT / 'testcase'
MATERIAL_DIR = TESTCASE_DIR / 'material'
SKILLS_DIR = PROJECT_ROOT / 'skills'

MCP_HOST = '127.0.0.1'
FULL_BUILD_INGEST_BATCHES = 3

if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def log(msg: str):
    print(f"[isolated_build] {msg}", flush=True)


def load_prompt_template(path: Path, **format_kwargs) -> str:
    """Load a prompt template from a source-controlled file."""
    text = path.read_text(encoding='utf-8').strip()
    text = re.sub(r'^---\n.*?\n---\n?', '', text, flags=re.DOTALL)
    if format_kwargs:
        return text.format(**format_kwargs)
    return text


def configured_llm_cli() -> str:
    """Return the CLI contract used by both runtime and benchmark workflows."""
    return os.environ.get("SEDIMENT_CLI", "claude").strip()


def chunk_list(lst: list, n: int) -> list[list]:
    """Split list into n roughly equal chunks."""
    k, m = divmod(len(lst), n)
    return [lst[i * k + min(i, m):(i + 1) * k + min(i + 1, m)] for i in range(n)]


def should_include_material_file(path: Path) -> bool:
    """Filter out benchmark-maintenance files that are not part of the actual project material set."""
    if path.name.startswith('.'):
        return False
    if any(part in {'__pycache__'} for part in path.parts):
        return False
    if path.suffix.lower() == '.py' and re.match(r'^(extract|fix|search)_', path.name):
        return False
    return True


def get_material_files() -> list[Path]:
    """Get all material files, sorted for deterministic batching."""
    files = []
    for f in MATERIAL_DIR.rglob('*'):
        if f.is_file() and should_include_material_file(f):
            files.append(f)
    files.sort(key=lambda p: str(p))
    return files


def collect_kb_diagnostics(kb_dir: Path) -> dict:
    """Collect health-style diagnostics for a built KB."""
    entries_dir = kb_dir / 'entries'
    placeholders_dir = kb_dir / 'placeholders'

    entry_files = sorted(p for p in entries_dir.glob('*.md') if p.name != '.gitkeep') if entries_dir.exists() else []
    placeholder_files = sorted(
        p for p in placeholders_dir.glob('*.md') if p.name != '.gitkeep'
    ) if placeholders_dir.exists() else []

    entry_sizes = [len(p.read_text(encoding='utf-8')) for p in entry_files]
    avg_entry_size = round(sum(entry_sizes) / len(entry_sizes), 1) if entry_sizes else 0.0

    diagnostics = {
        'kb_path': str(kb_dir),
        'entry_count': len(entry_files),
        'placeholder_count': len(placeholder_files),
        'avg_entry_size': avg_entry_size,
        'sample_entries': [p.stem for p in entry_files[:10]],
        'sample_placeholders': [p.stem for p in placeholder_files[:10]],
    }

    try:
        from mcp_server.kb import audit_kb, count_placeholder_refs

        audit_report = audit_kb(kb_dir)
        placeholder_refs = count_placeholder_refs(str(kb_dir))

        diagnostics.update({
            'formal_entry_count': audit_report.get('formal_entry_count', 0),
            'invalid_placeholder_count': audit_report.get('invalid_placeholder_count', 0),
            'invalid_placeholder_entries': audit_report.get('invalid_placeholder_entries', []),
            'dangling_link_count': audit_report.get('dangling_link_count', 0),
            'dangling_links': audit_report.get('dangling_links', []),
            'orphan_entry_count': audit_report.get('orphan_entry_count', 0),
            'orphan_entries': audit_report.get('orphan_entries', []),
            'canonical_gap_count': audit_report.get('canonical_gap_count', 0),
            'canonical_gaps': audit_report.get('canonical_gaps', []),
            'provenance_contamination_count': audit_report.get('provenance_contamination_count', 0),
            'provenance_contamination': audit_report.get('provenance_contamination', []),
            'placeholder_ref_summary': {
                'high': sum(1 for item in placeholder_refs if item.get('ref_count', 0) >= 5),
                'medium': sum(
                    1 for item in placeholder_refs if 2 <= item.get('ref_count', 0) <= 4
                ),
                'low': sum(1 for item in placeholder_refs if item.get('ref_count', 0) == 1),
            },
            'top_placeholders': placeholder_refs[:20],
            'promotable_placeholders': audit_report.get('promotable_placeholders', []),
        })
    except Exception as exc:
        diagnostics['health_error'] = str(exc)

    return diagnostics


# ---------------------------------------------------------------------------
# Material Extraction
# ---------------------------------------------------------------------------

def extract_material_text(file_path: Path) -> str:
    """Extract text content from a material file based on its extension."""
    ext = file_path.suffix.lower()

    if ext == '.md':
        return file_path.read_text(encoding='utf-8')

    if ext == '.txt':
        return file_path.read_text(encoding='utf-8')

    if ext == '.py':
        return _extract_python_docs(file_path)

    if ext in ('.cpp', '.h'):
        return _extract_cpp_docs(file_path)

    if ext == '.xml':
        return file_path.read_text(encoding='utf-8')

    if ext == '.yaml':
        return file_path.read_text(encoding='utf-8')

    if ext == '.json':
        return _extract_json_docs(file_path)

    if ext == '.puml':
        return file_path.read_text(encoding='utf-8')

    if ext == '.docx':
        try:
            from docx import Document
            doc = Document(file_path)
            return '\n'.join(p.text for p in doc.paragraphs)
        except ImportError:
            return f"[DOCX file: {file_path.name} - python-docx not installed]"

    if ext == '.pptx':
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
            return '\n'.join(texts)
        except ImportError:
            return f"[PPTX file: {file_path.name} - python-pptx not installed]"

    return f"[Unsupported format: {ext}]"


def _extract_python_docs(file_path: Path) -> str:
    """Extract docstrings, comments, and signatures from Python files."""
    import ast

    source = file_path.read_text(encoding='utf-8')

    try:
        tree = ast.parse(source)
    except SyntaxError:
        return source

    parts = []

    if (tree.body and isinstance(tree.body[0], ast.Expr)
            and isinstance(tree.body[0].value, ast.Constant)):
        parts.append(f"# 文件: {file_path.name}")
        parts.append(tree.body[0].value.value)
        parts.append("")

    comment_lines = []
    for i, line in enumerate(source.split('\n'), 1):
        stripped = line.strip()
        if stripped.startswith('#') and len(stripped) > 2:
            comment_lines.append(stripped.lstrip('# ').strip())

    for node in ast.walk(tree):
        if isinstance(node, (ast.ClassDef, ast.FunctionDef, ast.AsyncFunctionDef)):
            kind = '类' if isinstance(node, ast.ClassDef) else '函数'
            sig = f"{kind}: {node.name}"
            docstring = ast.get_docstring(node) or ''

            if docstring or comment_lines:
                parts.append(f"## {sig}")
                if docstring:
                    parts.append(docstring)
                parts.append("")

    if comment_lines:
        parts.append("## 代码注释")
        parts.extend(comment_lines)

    result = '\n'.join(parts)

    if len(result) < 200:
        if len(source) > 8000:
            result = source[:4000] + "\n...[中间部分省略]...\n" + source[-4000:]
        else:
            result = source

    return result


def _extract_cpp_docs(file_path: Path) -> str:
    """Extract comments and signatures from C++ files."""
    source = file_path.read_text(encoding='utf-8')

    block_comments = re.findall(r'/\*(.*?)\*/', source, re.DOTALL)
    line_comments = re.findall(r'//\s*(.+)', source)
    signatures = re.findall(r'(?:class|struct|namespace|enum)\s+(\w+)', source)
    func_sigs = re.findall(r'(\w+)\s*\([^)]*\)\s*(?:const\s*)?(?:override\s*)?(?:\{|;)', source)

    parts = [f"# 文件: {file_path.name}"]

    if block_comments:
        parts.append("## 块注释")
        parts.extend(c.strip() for c in block_comments if c.strip())

    if line_comments:
        parts.append("## 行注释")
        parts.extend(line_comments)

    if signatures:
        parts.append("## 类型定义")
        parts.extend(f"- {s}" for s in signatures)

    if func_sigs:
        parts.append("## 函数")
        parts.extend(f"- {f}" for f in func_sigs)

    result = '\n'.join(parts)

    if len(result) < 200:
        if len(source) > 8000:
            result = source[:4000] + "\n...[中间部分省略]...\n" + source[-4000:]
        else:
            result = source

    return result


def _extract_json_docs(file_path: Path) -> str:
    """Extract meaningful domain knowledge from JSON files."""
    source = file_path.read_text(encoding='utf-8')

    try:
        data = json.loads(source)
    except json.JSONDecodeError:
        return source

    def extract_keys(obj, prefix=''):
        parts = []
        if isinstance(obj, dict):
            for key, value in obj.items():
                full_key = f"{prefix}.{key}" if prefix else key
                if isinstance(value, dict):
                    parts.append(f"{full_key}: (object)")
                    parts.extend(extract_keys(value, full_key))
                elif isinstance(value, list):
                    if value and isinstance(value[0], dict):
                        parts.append(f"{full_key}: (array of {len(value)} objects)")
                        parts.extend(extract_keys(value[0], f"{full_key}[0]"))
                    else:
                        sample = str(value[:5])
                        parts.append(f"{full_key}: {sample}")
                elif isinstance(value, (str, int, float, bool)):
                    parts.append(f"{full_key}: {value}")
                else:
                    parts.append(f"{full_key}: {type(value).__name__}")
        return parts

    lines = [f"# 文件: {file_path.name}"]

    if isinstance(data, dict):
        if 'metrics' in data and isinstance(data['metrics'], list):
            lines.append(f"\n## 概述")
            lines.append(f"文件描述: {data.get('description', '')}")
            lines.append(f"版本: {data.get('version', '')}")
            lines.append(f"维护者: {data.get('maintainer', '')}")
            lines.append(f"\n## 指标定义详情")
            lines.append(f"该文件定义了系统中所有关键监控指标的元数据。")
            for metric in data['metrics']:
                name = metric.get('name', '未知')
                code = metric.get('code', '')
                desc = metric.get('description', '')
                unit = metric.get('unit', '')
                lines.append(f"\n### {name} ({code})")
                lines.append(f"- 描述: {desc}")
                lines.append(f"- 单位: {unit}")
        elif 'nodes' in data and isinstance(data['nodes'], list):
            lines.append(f"\n## 部署拓扑")
            lines.append(f"文件描述: {data.get('description', '')}")
            summary = data.get('topology_summary', {})
            if isinstance(summary, dict):
                summary_text = summary.get('summary_text')
                if summary_text:
                    lines.append(f"- 概览: {summary_text}")
                for key, value in summary.items():
                    if key == 'resonator_roles':
                        lines.append(f"- 谐振腔角色分布: {value}")
                    elif key != 'summary_text':
                        lines.append(f"- {key}: {value}")
            ops = data.get('operations_profile', {})
            if isinstance(ops, dict):
                lines.append(f"\n## 运维画像")
                for key, value in ops.items():
                    lines.append(f"- {key}: {value}")
            connections = data.get('connections', [])
            if isinstance(connections, list):
                lines.append(f"- connection_count: {len(connections)}")
            for node in data['nodes']:
                nid = node.get('id', '')
                role = node.get('role', '')
                zone = node.get('zone', '')
                lines.append(f"- {nid}: 角色={role}, 区域={zone}")
        else:
            lines.append("## 数据结构")
            for key, value in data.items():
                if isinstance(value, dict):
                    lines.append(f"\n### {key}")
                    for line in extract_keys(value, key):
                        lines.append(f"- {line}")
                elif isinstance(value, list):
                    lines.append(f"\n### {key} ({len(value)} items)")
    elif isinstance(data, list):
        lines.append(f"## 数组 ({len(data)} items)")

    result = '\n'.join(lines)

    if len(result) < 200:
        if len(source) > 8000:
            result = source[:4000] + "\n...[中间部分省略]...\n" + source[-4000:]
        else:
            result = source

    return result


# ---------------------------------------------------------------------------
# Ingest Prompt Builder
# ---------------------------------------------------------------------------

def build_ingest_prompt(materials: list[Path], isolated_dir: Path) -> str:
    """Build the ingest prompt for Claude."""
    kb_dir = isolated_dir / 'knowledge-base'
    entry_names = sorted(p.stem for p in (kb_dir / 'entries').glob('*.md'))
    placeholder_names = sorted(p.stem for p in (kb_dir / 'placeholders').glob('*.md'))
    has_existing_kb = bool(entry_names or placeholder_names)

    parts = [
        load_prompt_template(SKILLS_DIR / 'ingest' / 'SKILL.md'),
        """
## Benchmark Harness Instructions

- This is a non-interactive benchmark run. Do not ask the user for confirmation.
- Use `SEDIMENT_KB_PATH` as the only knowledge base root.
- Write formal entries to `entries/` and placeholders to `placeholders/`.
- Keep source document names in frontmatter `sources`, using plain text only.
- Never create KB entries or placeholders for source document titles, filenames, report names,
  manual names, plan names, or slide deck titles unless the title itself is independently reused
  as a first-class domain concept.
- Preserve exact source term spellings for first-class concepts. Do not invent near-synonym titles.
- Benchmark priority: produce a KB that answers direct `什么是X` queries from canonical bare-term entries.
        """.strip(),
    ]

    if has_existing_kb:
        inventory_preview = ', '.join(entry_names[:20])
        placeholder_preview = ', '.join(placeholder_names[:12])
        parts.append(
            f"""
## Incremental KB Instructions

- This KB already has {len(entry_names)} formal entries and {len(placeholder_names)} placeholders.
- Before creating new entries, inspect the existing KB to avoid near-duplicates and to preserve graph continuity.
- Prefer enriching or linking an existing concept over creating a parallel title for the same idea.
- If a first-class concept currently exists only as a placeholder and the new materials define it, promote it into the canonical formal entry.
- If you create a sentence-style lesson entry for a core concept, ensure the canonical bare-term concept entry also exists.
- Only edit existing entries when the new material clearly strengthens or corrects them; do not rewrite unrelated entries.
- Existing entry preview: {inventory_preview or '(none)'}
- Existing placeholder preview: {placeholder_preview or '(none)'}
            """.strip()
        )
    else:
        parts.append(
            """
## Initial KB Instructions

- The KB is currently empty. Focus on creating the first coherent graph from these materials.
            """.strip()
        )

    for f in materials:
        content = extract_material_text(f)
        if len(content) > 12000:
            content = content[:6000] + "\n...[中间部分省略]...\n" + content[-6000:]
        parts.append(f"\n{'='*60}\n文件: {f.relative_to(MATERIAL_DIR)}\n{'='*60}\n{content}\n")

    return '\n'.join(parts)


def _sample_noncanonical_titles(kb_dir: Path, limit: int = 12) -> list[str]:
    markers = ('必须', '需要', '应该', '不应', '不得', '导致', '引发', '可', '会', '前', '后', '时')
    titles = []
    for path in sorted((kb_dir / 'entries').glob('*.md')):
        title = path.stem
        if len(title) >= 7 and any(marker in title for marker in markers):
            titles.append(title)
    return titles[:limit]


def build_tidy_prompt(kb_dir: Path, focus: str = 'general') -> str:
    """Build the tidy prompt for Claude."""
    entry_count = len(list((kb_dir / 'entries').glob('*.md')))
    placeholder_count = len(list((kb_dir / 'placeholders').glob('*.md')))
    placeholder_preview = ', '.join(p.stem for p in sorted((kb_dir / 'placeholders').glob('*.md'))[:12])
    noncanonical_preview = ', '.join(_sample_noncanonical_titles(kb_dir))

    focus_block = ""
    if focus == 'canonicalization':
        focus_block = f"""
## Final Convergence Instructions

- This is the final KB convergence pass, not a light cleanup pass.
- Prioritize canonical bare-term concept entries for roles, tools, states, metrics, protocols, and operations.
- Promote placeholders into formal entries when the KB already contains enough evidence.
- Merge or fold shallow sentence-style definitional entries into canonical concept entries when they are not truly distinct lessons.
- Keep sentence-style entries only when they encode a distinct rule, failure pattern, anti-pattern, or causal lesson.
- Reduce duplicate concept surfaces. Prefer one strong canonical entry with aliases over several weak definitional siblings.
- Remove provenance-only placeholders or entries created from source-document titles. Keep those titles as plain provenance in `Source`, not as KB concepts.
- Normalize invented near-synonyms and typo-like variants back to the exact source term.
- Do not create MOC / index notes in this pass unless absolutely necessary.
- Placeholder preview: {placeholder_preview or '(none)'}
- Sentence-style entry preview: {noncanonical_preview or '(none)'}
        """.strip()

    return '\n\n'.join([
        load_prompt_template(SKILLS_DIR / 'tidy' / 'SKILL.md'),
        f"""
## Benchmark Harness Instructions

- This is a non-interactive benchmark run. Do not wait for user confirmation.
- Knowledge base path: `{kb_dir}`
- Current entries: {entry_count}
- Current placeholders: {placeholder_count}
- Execute the tidy workflow directly in this isolated knowledge base.
- You may create placeholders, repair links, and improve clearly broken structure directly.
- Be conservative about destructive edits: preserve formal entries unless a duplicate merge is obviously correct and the information is retained.
- Provenance is not graph structure: do not create or preserve KB nodes that only represent source documents.
        """.strip(),
        focus_block,
    ])


# ---------------------------------------------------------------------------
# Ingest & Tidy Execution
# ---------------------------------------------------------------------------

async def _stream_subprocess(proc: asyncio.subprocess.Process, label: str):
    """Stream subprocess stdout/stderr to terminal in real time."""
    async def _read_stream(stream, label_tag):
        while True:
            line = await stream.readline()
            if not line:
                break
            text = line.decode('utf-8', errors='replace').rstrip()
            if text:
                print(f"  [{label_tag}] {text}", flush=True)

    tasks = []
    if proc.stdout:
        tasks.append(asyncio.create_task(_read_stream(proc.stdout, f"{label}-out")))
    if proc.stderr:
        tasks.append(asyncio.create_task(_read_stream(proc.stderr, f"{label}-err")))
    if tasks:
        await asyncio.gather(*tasks)


async def _spawn_agent_process(
    *,
    prompt: str,
    cwd: Path,
    env: dict[str, str],
    max_budget_usd: str,
) -> asyncio.subprocess.Process:
    cli_value = configured_llm_cli()
    command, stdin_data = build_cli_command(
        cli_value,
        prompt,
        extra_args=[
            "--permission-mode",
            "auto",
            "--allowed-tools",
            "Write",
            "Edit",
            "Bash",
            "Read",
            "Glob",
            "--max-budget-usd",
            max_budget_usd,
            "--no-session-persistence",
        ],
    )
    proc = await asyncio.create_subprocess_exec(
        *command,
        cwd=str(cwd),
        env=env,
        stdin=asyncio.subprocess.PIPE if stdin_data is not None else asyncio.subprocess.DEVNULL,
        stdout=asyncio.subprocess.PIPE,
        stderr=asyncio.subprocess.PIPE,
    )
    if stdin_data is not None and proc.stdin is not None:
        proc.stdin.write(stdin_data.encode("utf-8"))
        await proc.stdin.drain()
        proc.stdin.close()
    return proc


async def run_ingest(isolated_dir: Path, materials: list[Path]) -> bool:
    """Run ingest through the shared Sediment CLI contract."""
    kb_entries = isolated_dir / 'knowledge-base' / 'entries'
    kb_entries.mkdir(parents=True, exist_ok=True)

    MAX_RETRIES = 2

    for attempt in range(1, MAX_RETRIES + 1):
        prompt = build_ingest_prompt(materials, isolated_dir)

        env = os.environ.copy()
        env['SEDIMENT_KB_PATH'] = str(isolated_dir / 'knowledge-base')
        env['CLAUDE_CODE'] = '1'

        log(f"Running ingest for {len(materials)} files (attempt {attempt}/{MAX_RETRIES})...")
        start = time.time()

        try:
            proc = await _spawn_agent_process(
                prompt=prompt,
                cwd=isolated_dir,
                env=env,
                max_budget_usd='15',
            )

            await _stream_subprocess(proc, 'ingest')
            try:
                await asyncio.wait_for(proc.wait(), timeout=3600)
            except asyncio.TimeoutError:
                log(f"  Ingest timed out after 60 minutes, killing process...")
                proc.kill()
                await proc.wait()
                if attempt < MAX_RETRIES:
                    log("Retrying in 5 seconds...")
                    await asyncio.sleep(5)
                    continue
                return False

            elapsed = time.time() - start
            entry_count = len(list(kb_entries.glob('*.md')))
            log(f"  Ingest complete in {elapsed:.1f}s ({elapsed/60:.1f}min). Created {entry_count} entries.")

            if proc.returncode != 0:
                log(f"Warning: ingest agent exited with code {proc.returncode}")
                if attempt < MAX_RETRIES:
                    log("Retrying in 5 seconds...")
                    await asyncio.sleep(5)
                    continue

            return entry_count > 0

        except Exception as e:
            log(f"Ingest attempt {attempt} failed: {e}")
            if attempt < MAX_RETRIES:
                log("Retrying in 5 seconds...")
                await asyncio.sleep(5)
                continue
            return False

    return False


async def run_tidy(isolated_dir: Path, focus: str = 'general') -> bool:
    """Run tidy through the shared Sediment CLI contract."""
    kb_dir = isolated_dir / 'knowledge-base'
    if not kb_dir.exists():
        return False

    entry_count = len(list((kb_dir / 'entries').glob('*.md')))
    placeholder_count = len(list((kb_dir / 'placeholders').glob('*.md')))
    log(f"Running tidy ({focus}). Current: {entry_count} entries, {placeholder_count} placeholders.")

    MAX_RETRIES = 2

    for attempt in range(1, MAX_RETRIES + 1):
        prompt = build_tidy_prompt(kb_dir, focus=focus)

        env = os.environ.copy()
        env['SEDIMENT_KB_PATH'] = str(kb_dir)

        log(f"Tidy ({focus}) attempt {attempt}/{MAX_RETRIES}...")
        start = time.time()

        try:
            proc = await _spawn_agent_process(
                prompt=prompt,
                cwd=isolated_dir,
                env=env,
                max_budget_usd='5',
            )

            await _stream_subprocess(proc, 'tidy')
            try:
                await asyncio.wait_for(proc.wait(), timeout=600)
            except asyncio.TimeoutError:
                log(f"  Tidy timed out after 10 minutes, killing process...")
                proc.kill()
                await proc.wait()
                if attempt < MAX_RETRIES:
                    log("Retrying in 5 seconds...")
                    await asyncio.sleep(5)
                    continue
                return False

            elapsed = time.time() - start
            new_entry_count = len(list((kb_dir / 'entries').glob('*.md')))
            new_placeholder_count = len(list((kb_dir / 'placeholders').glob('*.md')))
            log(
                f"  Tidy ({focus}) complete in {elapsed:.1f}s. "
                f"Now: {new_entry_count} entries, {new_placeholder_count} placeholders."
            )

            if proc.returncode != 0:
                log(f"Warning: tidy ({focus}) exited with code {proc.returncode}")
                if attempt < MAX_RETRIES:
                    log("Retrying in 5 seconds...")
                    await asyncio.sleep(5)
                    continue

            return True

        except Exception as e:
            log(f"Tidy ({focus}) attempt {attempt} failed: {e}")
            if attempt < MAX_RETRIES:
                log("Retrying in 5 seconds...")
                await asyncio.sleep(5)
                continue
            return False

    return False


# ---------------------------------------------------------------------------
# MCP Server
# ---------------------------------------------------------------------------

class MCPServer:
    """Manages the MCP server process, running inside the isolated directory."""

    def __init__(self, kb_path: Path, port: int, isolated_dir: Path):
        self.kb_path = kb_path
        self.port = port
        self.isolated_dir = isolated_dir
        self.process = None
        self.base_url = f'http://{MCP_HOST}:{port}'
        self.sse_url = f'{self.base_url}/sediment/'

    async def start(self):
        """Start the MCP server."""
        # Kill any existing process on the target port
        try:
            proc = await asyncio.create_subprocess_exec(
                'lsof', '-ti', f':{self.port}',
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
            )
            stdout, _ = await proc.communicate()
            if stdout:
                for pid in stdout.decode().strip().split('\n'):
                    pid = pid.strip()
                    if pid:
                        log(f"  Killing existing process on port {self.port}: PID {pid}")
                        await asyncio.create_subprocess_exec('kill', '-9', pid)
                        await asyncio.sleep(0.5)
        except Exception:
            pass

        env = os.environ.copy()
        env['SEDIMENT_KB_PATH'] = str(self.kb_path)
        env['SEDIMENT_PORT'] = str(self.port)
        env['SEDIMENT_HOST'] = MCP_HOST
        env['SEDIMENT_RUNTIME_ALLOW_MATERIAL_FALLBACK'] = '0'
        # Ensure Python can find mcp_server module from the isolated directory
        env['PYTHONPATH'] = str(self.isolated_dir)

        venv_python = self.isolated_dir / '.venv' / 'bin' / 'python'
        if venv_python.exists():
            cmd = [
                str(venv_python), '-m', 'mcp_server.server',
            ]
            log(f"  Using venv python: {venv_python}")
        else:
            # Use 'uv run' directly with the uv binary
            cmd = [
                'uv', 'run', 'python', '-m', 'mcp_server.server',
            ]
            log(f"  Using uv run")

        log(f"  Launching MCP server: {' '.join(cmd)}")
        self.process = await asyncio.create_subprocess_exec(
            *cmd,
            cwd=str(self.isolated_dir),
            env=env,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )

        # Wait for server to be ready
        import httpx
        log(f"  Waiting for MCP server on {self.base_url}...")
        for i in range(120):
            await asyncio.sleep(0.5)
            elapsed = (i + 1) * 0.5
            if elapsed.is_integer() and int(elapsed) % 20 == 0:
                log(f"  Still waiting... ({int(elapsed)}s)")

            if self.process.returncode is not None:
                stdout_tail = ''
                stderr_tail = ''
                if self.process.stdout:
                    stdout_tail = (await self.process.stdout.read()).decode(errors='ignore')[-1200:]
                if self.process.stderr:
                    stderr_tail = (await self.process.stderr.read()).decode(errors='ignore')[-1200:]
                log(f"  MCP server exited early with code {self.process.returncode}")
                if stdout_tail.strip():
                    log(f"  server stdout tail:\n{stdout_tail.rstrip()}")
                if stderr_tail.strip():
                    log(f"  server stderr tail:\n{stderr_tail.rstrip()}")
                return False
            try:
                async with httpx.AsyncClient(timeout=2) as client:
                    resp = await client.post(
                        self.sse_url,
                        json={
                            'jsonrpc': '2.0',
                            'id': 1,
                            'method': 'tools/list',
                        },
                        headers={'Content-Type': 'application/json'},
                    )
                    if resp.status_code == 200:
                        log(f"  MCP server ready on port {self.port} (after {elapsed:.0f}s)")
                        return True
            except Exception:
                continue

        log("  MCP server failed to start (timeout after 60s)")
        return False

    async def call_tool(self, tool_name: str, arguments: dict) -> str:
        """Call an MCP tool via direct JSON-RPC."""
        import httpx
        try:
            async with httpx.AsyncClient(timeout=300) as client:
                resp = await client.post(
                    self.sse_url,
                    json={
                        'jsonrpc': '2.0',
                        'id': random.randint(1, 10000),
                        'method': 'tools/call',
                        'params': {
                            'name': tool_name,
                            'arguments': arguments,
                        },
                    },
                    headers={'Content-Type': 'application/json'},
                )
                if resp.status_code == 200:
                    data = resp.json()
                    result = data.get('result', {})
                    content = result.get('content', [])
                    if content:
                        return content[0].get('text', '')
                    return json.dumps(result, ensure_ascii=False)
                return f"HTTP {resp.status_code}"
        except Exception as e:
            return f"ERROR: {e}"

    async def stop(self):
        """Stop the MCP server."""
        if self.process:
            try:
                self.process.terminate()
                try:
                    await asyncio.wait_for(self.process.wait(), timeout=10)
                except asyncio.TimeoutError:
                    self.process.kill()
            except ProcessLookupError:
                pass  # Process already exited
            log(f"MCP server on port {self.port} stopped")


# ---------------------------------------------------------------------------
# IsolatedBuilder: Main orchestrator
# ---------------------------------------------------------------------------

class IsolatedBuilder:
    """
    Manages the full lifecycle: create isolated copy → build KB → start MCP → cleanup.

    Usage as context manager:
        async with IsolatedBuilder(build_type='full') as builder:
            server = builder.mcp_server
            # use server for testing...
    """

    def __init__(self, build_type: str = 'full', port: int = 18800):
        self.build_type = build_type
        self.port = port
        self.isolated_dir: Path | None = None
        self.kb_dir: Path | None = None
        self.mcp_server: MCPServer | None = None

    async def __aenter__(self):
        await self.create_isolated_copy()
        await self.build()
        return self

    async def __aexit__(self, *args):
        await self.cleanup()

    async def create_isolated_copy(self) -> Path:
        """Create an isolated copy of the project."""
        log(f"Creating isolated copy...")
        self.isolated_dir = Path(tempfile.mkdtemp(prefix=f'sediment-{self.build_type}-'))
        log(f"  Isolated dir: {self.isolated_dir}")

        log(f"  Copying project files (excluding .git, __pycache__, etc.)...")
        shutil.copytree(
            PROJECT_ROOT, self.isolated_dir, dirs_exist_ok=True,
            ignore=shutil.ignore_patterns(
                '.git', '__pycache__', '*.pyc',
                '.venv', 'results', 'results_bak', '.claude', 'node_modules',
            ),
        )
        log(f"  Copy complete")

        # Ensure KB directories exist
        self.kb_dir = self.isolated_dir / 'knowledge-base'
        self.kb_dir.mkdir(exist_ok=True)
        (self.kb_dir / 'entries').mkdir(exist_ok=True)
        (self.kb_dir / 'placeholders').mkdir(exist_ok=True)
        log(f"  KB directories created at {self.kb_dir}")

        return self.isolated_dir

    async def build(self):
        """Execute the build pipeline based on build_type."""
        if self.build_type == 'full':
            await self.build_full()
        elif self.build_type == 'batched':
            await self.build_batched()
        else:
            raise ValueError(f"Unknown build_type: {self.build_type}")

    async def build_full(self):
        """Full build: ingest all files in a few large chunks, tidy once at end."""
        materials = get_material_files()
        batches = [batch for batch in chunk_list(materials, FULL_BUILD_INGEST_BATCHES) if batch]
        total_size = sum(f.stat().st_size for f in materials) / 1024  # KB

        log(f"\n{'─'*60}")
        log(f"Phase 1: Full ingest in {len(batches)} chunks ({len(materials)} files, {total_size:.0f}KB)")
        log(f"{'─'*60}")

        total_ingest_start = time.time()
        for i, batch in enumerate(batches, 1):
            batch_size = sum(f.stat().st_size for f in batch) / 1024
            log(f"\n  [full {i}/{len(batches)}] ingest {len(batch)} files ({batch_size:.0f}KB)")
            batch_start = time.time()
            success = await run_ingest(self.isolated_dir, batch)
            batch_elapsed = time.time() - batch_start
            if not success:
                log(f"  Full ingest chunk {i} failed after {batch_elapsed:.1f}s")
            else:
                log(f"  Full ingest chunk {i} complete in {batch_elapsed:.1f}s")

            entry_count = len(list((self.kb_dir / 'entries').glob('*.md')))
            placeholder_count = len(list((self.kb_dir / 'placeholders').glob('*.md')))
            log(f"    KB now: {entry_count} entries, {placeholder_count} placeholders")

        log(f"Total ingest phase complete in {time.time() - total_ingest_start:.1f}s")

        log(f"\n{'─'*60}")
        log(f"Phase 2: Tidy")
        log(f"{'─'*60}")
        tidy_start = time.time()
        await run_tidy(self.isolated_dir, focus='general')
        tidy_elapsed = time.time() - tidy_start
        log(f"Tidy phase complete in {tidy_elapsed:.1f}s")

        log(f"\n{'─'*60}")
        log("Phase 3: Canonical convergence tidy")
        log(f"{'─'*60}")
        convergence_start = time.time()
        await run_tidy(self.isolated_dir, focus='canonicalization')
        convergence_elapsed = time.time() - convergence_start
        log(f"Canonical convergence complete in {convergence_elapsed:.1f}s")

        final_entry_count = len(list((self.kb_dir / 'entries').glob('*.md')))
        final_placeholder_count = len(list((self.kb_dir / 'placeholders').glob('*.md')))
        diagnostics = collect_kb_diagnostics(self.kb_dir)
        log(f"  After tidy: {final_entry_count} entries, {final_placeholder_count} placeholders")
        log(
            "  Health summary: "
            f"{diagnostics.get('orphan_entry_count', 0)} orphans, "
            f"{diagnostics.get('dangling_link_count', 0)} dangling links, "
            f"{diagnostics.get('placeholder_ref_summary', {}).get('high', 0)} high-ref placeholders"
        )
        log(f"{'─'*60}")

    async def build_batched(self):
        """Batched build: ingest 1/5 at a time, tidy after each."""
        materials = get_material_files()
        batches = [batch for batch in chunk_list(materials, 5) if batch]
        for i, batch in enumerate(batches):
            batch_size = sum(f.stat().st_size for f in batch) / 1024  # KB
            log(f"\n{'─'*60}")
            log(f"Phase {i+1}/5: Batch ingest ({len(batch)} files, {batch_size:.0f}KB)")
            log(f"{'─'*60}")
            batch_start = time.time()
            success = await run_ingest(self.isolated_dir, batch)
            if success:
                log(f"  Batch ingest complete in {time.time() - batch_start:.1f}s")
                log(f"  Running tidy...")
                tidy_start = time.time()
                await run_tidy(self.isolated_dir, focus='general')
                log(f"  Tidy complete in {time.time() - tidy_start:.1f}s")
            else:
                log(f"  Batch ingest failed after {time.time() - batch_start:.1f}s")
            await asyncio.sleep(1)

        log(f"\n{'─'*60}")
        log("Final global tidy across all batched ingests")
        log(f"{'─'*60}")
        final_tidy_start = time.time()
        await run_tidy(self.isolated_dir, focus='general')
        log(f"  Final global tidy complete in {time.time() - final_tidy_start:.1f}s")

        log(f"\n{'─'*60}")
        log("Final canonical convergence tidy")
        log(f"{'─'*60}")
        final_convergence_start = time.time()
        await run_tidy(self.isolated_dir, focus='canonicalization')
        log(f"  Final canonical convergence complete in {time.time() - final_convergence_start:.1f}s")

        final_entry_count = len(list((self.kb_dir / 'entries').glob('*.md')))
        final_placeholder_count = len(list((self.kb_dir / 'placeholders').glob('*.md')))
        log(f"\n{'─'*60}")
        log(f"Batched build complete: {final_entry_count} entries, {final_placeholder_count} placeholders")
        log(f"{'─'*60}")

    def start_mcp_server(self, port: int | None = None) -> MCPServer:
        """Start the MCP server. Returns the server instance."""
        if port is None:
            port = self.port
        self.mcp_server = MCPServer(self.kb_dir, port, self.isolated_dir)
        return self.mcp_server

    async def cleanup(self, remove_dir: bool = True):
        """Stop MCP server, clean up background processes, optionally remove isolated directory."""
        if self.isolated_dir:
            log(f"Cleaning up isolated directory: {self.isolated_dir}")
        # Kill any lingering LLM CLI subprocesses spawned in this isolated dir
        try:
            pkill_proc = await asyncio.create_subprocess_exec(
                'pkill', '-f', str(self.isolated_dir),
                stdout=asyncio.subprocess.DEVNULL,
                stderr=asyncio.subprocess.DEVNULL,
            )
            await pkill_proc.wait()
        except Exception:
            pass

        if self.mcp_server:
            await self.mcp_server.stop()
            self.mcp_server = None
        if not remove_dir:
            if self.isolated_dir and self.isolated_dir.exists():
                log(f"  Preserved isolated directory for diagnosis: {self.isolated_dir}")
            else:
                log("  Preserve requested, but isolated dir is already missing")
            return
        if self.isolated_dir and self.isolated_dir.exists():
            try:
                shutil.rmtree(self.isolated_dir, ignore_errors=True)
                log(f"  Cleanup complete: {self.isolated_dir}")
            except Exception as e:
                log(f"  Cleanup warning: {e}")
        else:
            log(f"  Cleanup skipped (no isolated dir)")


# ---------------------------------------------------------------------------
# CLI Entry Point
# ---------------------------------------------------------------------------

async def main():
    parser = argparse.ArgumentParser(description='Sediment Isolated Build')
    parser.add_argument(
        '--build-type', choices=['full', 'batched'], default='full',
        help='Build type: full (single ingest + tidy) or batched (5-batch with tidy each)',
    )
    parser.add_argument('--port', type=int, default=18800, help='MCP server port')
    parser.add_argument('--dry-run', action='store_true', help='Only create isolated copy, skip build')
    parser.add_argument('--no-cleanup', action='store_true', help='Keep isolated directory after build (for debugging)')
    args = parser.parse_args()

    log(f"Build type: {args.build_type}")
    log(f"Project root: {PROJECT_ROOT}")
    log(f"Material files: {len(get_material_files())}")

    builder = IsolatedBuilder(build_type=args.build_type, port=args.port)

    try:
        await builder.create_isolated_copy()

        if not args.dry_run:
            await builder.build()

            # Report KB status
            entry_count = len(list((builder.kb_dir / 'entries').glob('*.md')))
            placeholder_count = len(list((builder.kb_dir / 'placeholders').glob('*.md')))
            log(f"\nBuild complete: {entry_count} entries, {placeholder_count} placeholders")
            log(f"KB path: {builder.kb_dir}")

            # Optionally start MCP server for quick testing
            log("\nTo test with the MCP server, run:")
            log(f"  export SEDIMENT_KB_PATH={builder.kb_dir}")
            log(f"  cd {builder.isolated_dir}")
            log(f"  python -m mcp_server.server --port {args.port}")
        else:
            log(f"\nDry run complete. Isolated dir: {builder.isolated_dir}")

        if args.no_cleanup:
            log(f"\nIsolated directory preserved: {builder.isolated_dir}")
            log("Remember to clean up manually: rm -rf " + str(builder.isolated_dir))
            # Don't cleanup
            builder.isolated_dir = None
    finally:
        if not args.no_cleanup:
            await builder.cleanup()


if __name__ == '__main__':
    asyncio.run(main())
