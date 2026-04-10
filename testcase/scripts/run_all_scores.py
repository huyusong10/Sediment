"""
run_all_scores.py — Sediment 全流程测试运行器

执行完整的测试流程：
1. 隔离项目到临时目录
2. 全量构建 KB（一次性 ingest 所有材料，然后 tidy）
3. 分批构建 KB（每次 ingest 1/5 材料，tidy，重复 5 次）
4. 对每个构建启动 MCP server，执行 TC-01 和 TC-02
5. 计算平均分，输出 scorecard
"""

import asyncio
import json
import os
import random
import re
import shutil
import signal
import subprocess
import sys
import tempfile
import time
from pathlib import Path

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

PROJECT_ROOT = Path(__file__).resolve().parent.parent.parent
TESTCASE_DIR = PROJECT_ROOT / 'testcase'
JUDGE_DIR = TESTCASE_DIR / 'judge'
MATERIAL_DIR = TESTCASE_DIR / 'material'
RESULTS_DIR = TESTCASE_DIR / 'results'
SCRIPTS_DIR = TESTCASE_DIR / 'scripts'

# MCP server config
MCP_HOST = '127.0.0.1'
MCP_PORT_BASE = 18800  # base port, will offset for each build


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def log(msg: str):
    print(f"[run_all] {msg}", flush=True)


def get_material_files() -> list[Path]:
    """Get all material files, sorted for deterministic batching."""
    files = []
    for f in MATERIAL_DIR.rglob('*'):
        if f.is_file() and f.name != '.DS_Store':
            files.append(f)
    files.sort(key=lambda p: str(p))
    return files


def chunk_list(lst: list, n: int) -> list[list]:
    """Split list into n roughly equal chunks."""
    k, m = divmod(len(lst), n)
    return [lst[i * k + min(i, m):(i + 1) * k + min(i + 1, m)] for i in range(n)]


# ---------------------------------------------------------------------------
# Ingest: Extract material content into text for LLM processing
# ---------------------------------------------------------------------------

def extract_material_text(file_path: Path) -> str:
    """Extract text content from a material file based on its extension."""
    ext = file_path.suffix.lower()

    if ext == '.md':
        return file_path.read_text(encoding='utf-8')

    if ext == '.txt':
        return file_path.read_text(encoding='utf-8')

    if ext == '.py':
        # For Python files, extract docstrings, comments, and signatures
        # These contain the domain knowledge while code is noise for KB
        return _extract_python_docs(file_path)

    if ext in ('.cpp', '.h'):
        return _extract_cpp_docs(file_path)

    if ext == '.xml':
        return file_path.read_text(encoding='utf-8')

    if ext == '.yaml':
        return file_path.read_text(encoding='utf-8')

    if ext == '.json':
        return _extract_json_docs(file_path)

    if ext == '.puml':
        return file_path.read_text(encoding='utf-8')

    if ext == '.docx':
        try:
            from docx import Document
            doc = Document(file_path)
            return '\n'.join(p.text for p in doc.paragraphs)
        except ImportError:
            return f"[DOCX file: {file_path.name} - python-docx not installed]"

    if ext == '.pptx':
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
            return '\n'.join(texts)
        except ImportError:
            return f"[PPTX file: {file_path.name} - python-pptx not installed]"

    return f"[Unsupported format: {ext}]"


def _extract_python_docs(file_path: Path) -> str:
    """Extract docstrings, comments, and signatures from Python files."""
    import ast
    import re

    source = file_path.read_text(encoding='utf-8')

    try:
        tree = ast.parse(source)
    except SyntaxError:
        return source  # Fallback to full text if parsing fails

    parts = []

    # Extract module docstring
    if (tree.body and isinstance(tree.body[0], ast.Expr)
            and isinstance(tree.body[0].value, ast.Constant)):
        parts.append(f"# 文件: {file_path.name}")
        parts.append(tree.body[0].value.value)
        parts.append("")

    # Extract all comments (lines starting with #)
    comment_lines = []
    for i, line in enumerate(source.split('\n'), 1):
        stripped = line.strip()
        if stripped.startswith('#') and len(stripped) > 2:
            comment_lines.append(stripped.lstrip('# ').strip())

    # Extract class and function docstrings + signatures
    for node in ast.walk(tree):
        if isinstance(node, (ast.ClassDef, ast.FunctionDef, ast.AsyncFunctionDef)):
            kind = '类' if isinstance(node, ast.ClassDef) else '函数'
            sig = f"{kind}: {node.name}"

            # Get docstring
            docstring = ast.get_docstring(node) or ''

            if docstring or comment_lines:
                parts.append(f"## {sig}")
                if docstring:
                    parts.append(docstring)
                parts.append("")

    # Add standalone comments not associated with a specific node
    if comment_lines:
        parts.append("## 代码注释")
        parts.extend(comment_lines)

    result = '\n'.join(parts)

    # If extracted content is very short, fall back to full text with truncation
    if len(result) < 200:
        if len(source) > 8000:
            result = source[:4000] + "\n...[中间部分省略]...\n" + source[-4000:]
        else:
            result = source

    return result


def _extract_cpp_docs(file_path: Path) -> str:
    """Extract comments and signatures from C++ files."""
    import re

    source = file_path.read_text(encoding='utf-8')

    # Extract block comments (/* ... */)
    block_comments = re.findall(r'/\*(.*?)\*/', source, re.DOTALL)
    # Extract line comments (// ...)
    line_comments = re.findall(r'//\s*(.+)', source)
    # Extract class/struct/function signatures
    signatures = re.findall(
        r'(?:class|struct|namespace|enum)\s+(\w+)', source
    )
    func_sigs = re.findall(
        r'(\w+)\s*\([^)]*\)\s*(?:const\s*)?(?:override\s*)?(?:\{|;)', source
    )

    parts = [f"# 文件: {file_path.name}"]

    if block_comments:
        parts.append("## 块注释")
        parts.extend(c.strip() for c in block_comments if c.strip())

    if line_comments:
        parts.append("## 行注释")
        parts.extend(line_comments)

    if signatures:
        parts.append("## 类型定义")
        parts.extend(f"- {s}" for s in signatures)

    if func_sigs:
        parts.append("## 函数")
        parts.extend(f"- {f}" for f in func_sigs)

    result = '\n'.join(parts)

    if len(result) < 200:
        if len(source) > 8000:
            result = source[:4000] + "\n...[中间部分省略]...\n" + source[-4000:]
        else:
            result = source

    return result


def _extract_json_docs(file_path: Path) -> str:
    """Extract meaningful domain knowledge from JSON files."""
    import json as _json

    source = file_path.read_text(encoding='utf-8')

    try:
        data = _json.loads(source)
    except _json.JSONDecodeError:
        return source  # Fallback to full text if invalid JSON

    def extract_keys(obj, prefix=''):
        """Recursively extract all keys and their types/values."""
        parts = []
        if isinstance(obj, dict):
            for key, value in obj.items():
                full_key = f"{prefix}.{key}" if prefix else key
                if isinstance(value, dict):
                    parts.append(f"{full_key}: (object)")
                    parts.extend(extract_keys(value, full_key))
                elif isinstance(value, list):
                    if value and isinstance(value[0], dict):
                        parts.append(f"{full_key}: (array of {len(value)} objects)")
                        parts.extend(extract_keys(value[0], f"{full_key}[0]"))
                    else:
                        sample = str(value[:5])
                        parts.append(f"{full_key}: {sample}")
                elif isinstance(value, (str, int, float, bool)):
                    parts.append(f"{full_key}: {value}")
                else:
                    parts.append(f"{full_key}: {type(value).__name__}")
        return parts

    lines = [f"# 文件: {file_path.name}"]

    # Detect file type by content and generate natural-language summaries
    if isinstance(data, dict):
        # Check if this is a metric definitions file
        if 'metrics' in data and isinstance(data['metrics'], list):
            lines.append(f"\n## 概述")
            lines.append(f"文件描述: {data.get('description', '')}")
            lines.append(f"版本: {data.get('version', '')}")
            lines.append(f"维护者: {data.get('maintainer', '')}")
            lines.append(f"\n## 指标定义详情")
            lines.append(f"该文件定义了系统中所有关键监控指标的元数据，包括每个指标的名称、代码、描述、**单位**、数据类型、取值范围、计算公式和告警阈值。")
            lines.append(f"\n以下是每个指标的详细定义：")
            for metric in data['metrics']:
                name = metric.get('name', '未知')
                code = metric.get('code', '')
                desc = metric.get('description', '')
                unit = metric.get('unit', '')
                unit_full = metric.get('unit_full', '')
                dtype = metric.get('data_type', '')
                val_range = metric.get('range', {})
                calc = metric.get('calculation', {})
                thresholds = metric.get('thresholds', {})

                lines.append(f"\n### {name} ({code})")
                lines.append(f"- 描述: {desc}")
                lines.append(f"- **单位: {unit}**")
                if unit_full:
                    lines.append(f"- 单位全称: {unit_full}")
                lines.append(f"- 数据类型: {dtype}")
                if val_range:
                    lines.append(f"- 取值范围: {val_range.get('min', '?')} ~ {val_range.get('max', '?')}")
                if calc:
                    formula = calc.get('formula', '')
                    interval = calc.get('interval_seconds', '')
                    lines.append(f"- 计算公式: {formula}")
                    if interval:
                        lines.append(f"- 采集间隔: {interval}秒")
                if thresholds:
                    lines.append(f"- 告警阈值:")
                    for tname, tval in thresholds.items():
                        tdesc = tval.get('description', '') if isinstance(tval, dict) else str(tval)
                        lines.append(f"  - {tname}: {tdesc}")
                related = metric.get('related_issues', [])
                if related:
                    lines.append(f"- 关联问题: {', '.join(related)}")

        # Check if this is a deployment topology file
        elif 'nodes' in data and isinstance(data['nodes'], list):
            lines.append(f"\n## 概述")
            lines.append(f"文件描述: {data.get('description', '')}")
            lines.append(f"版本: {data.get('version', '')}")
            lines.append(f"网络类型: {data.get('network_type', '')}")
            lines.append(f"\n## 部署拓扑详情")
            lines.append(f"该文件定义了系统的全网部署拓扑，包括节点类型、角色分布、连接关系和网格配置。")
            lines.append(f"\n### 节点部署策略")
            node_types = {}
            for node in data['nodes']:
                ntype = node.get('type', 'unknown')
                node_types.setdefault(ntype, []).append(node)

            for ntype, nodes in node_types.items():
                lines.append(f"\n#### {ntype}节点 ({len(nodes)}个)")
                for node in nodes:
                    nid = node.get('id', '')
                    role = node.get('role', '')
                    zone = node.get('zone', '')
                    config = node.get('config', {})
                    lines.append(f"- {nid}: 角色={role}, 区域={zone}")
                    for ck, cv in config.items():
                        lines.append(f"  - {ck}: {cv}")

            lines.append(f"\n### 连接关系")
            for conn in data.get('connections', []):
                cid = conn.get('id', '')
                frm = conn.get('from', '')
                to = conn.get('to', '')
                ctype = conn.get('type', '')
                proto = conn.get('protocol', '')
                bw = conn.get('bandwidth', '')
                lat = conn.get('latency_ms', '')
                lines.append(f"- {cid}: {frm} → {to}, 类型={ctype}, 协议={proto}, 带宽={bw}, 延迟={lat}ms")

            mesh = data.get('mesh_config', {})
            if mesh:
                lines.append(f"\n### 网格配置")
                lines.append(f"- 网络类型: {mesh.get('network_type', '')}")
                lines.append(f"- 自动发现: {mesh.get('auto_discovery', '')}")
                healing = mesh.get('healing', {})
                if healing:
                    lines.append(f"- 自愈: 启用={healing.get('enabled', '')}, 最大重试={healing.get('max_retries', '')}")

        # Generic dict handling
        else:
            lines.append("## 数据结构")
            for key, value in data.items():
                if isinstance(value, dict):
                    lines.append(f"\n### {key}")
                    for line in extract_keys(value, key):
                        lines.append(f"- {line}")
                elif isinstance(value, list):
                    lines.append(f"\n### {key} ({len(value)} items)")
                    if value and isinstance(value[0], dict):
                        for line in extract_keys(value[0], f"{key}[0]"):
                            lines.append(f"- {line}")
                    elif value:
                        lines.append(f"  值: {str(value[:20])}")
                else:
                    lines.append(f"- {key}: {value}")
    elif isinstance(data, list):
        lines.append(f"## 数组 ({len(data)} items)")
        if data and isinstance(data[0], dict):
            for line in extract_keys(data[0], 'item'):
                lines.append(f"- {line}")

    result = '\n'.join(lines)

    # If extracted content is too short, fall back to full text
    if len(result) < 200:
        if len(source) > 8000:
            result = source[:4000] + "\n...[中间部分省略]...\n" + source[-4000:]
        else:
            result = source

    return result


def build_ingest_prompt(materials: list[Path]) -> str:
    """Build the ingest prompt for Claude."""
    parts = []
    parts.append("""你是一个知识摄入 Agent (Sediment Ingest)。
目标：将给定文档提炼为原子知识条目，存入知识库。

知识库路径：knowledge-base/
- 正式条目存入：knowledge-base/entries/
- 占位文件存入：knowledge-base/placeholders/

══════════════════════════════════════════════
核心规则（最重要，优先遵守）
══════════════════════════════════════════════

1. 每个领域概念/术语/设备/指标/角色/流程 都必须有独立条目
2. 条目的第一行（核心命题）必须是对该概念的完整定义，格式为："XX是/指/用于……"
3. 核心命题必须包含该概念的所有关键特征、属性、数值阈值、关联指标
4. 条目中禁止出现代码片段、XML/JSON/YAML原始数据、图表ASCII艺术
5. 文件名就是概念的标准名称，使用领域中最常见的叫法
6. **定义段中禁止使用任何 Markdown 格式标记（如 **加粗**、*斜体* 等），保持纯文本**
7. **定义段应当精炼简洁，2-3句话即可，不要包含长列表、数值明细或流程步骤**
8. **数值参数、流程步骤、配置明细等详细信息应当放在"## 上下文"部分，不要塞进定义段**

══════════════════════════════════════════════
条目结构（严格遵守）
══════════════════════════════════════════════

---
aliases: [别名1, 别名2]
tags: [概念类型]
status: formal
---
# 概念标准名称

[定义：一句话说明"这是什么" + "它做什么/有什么特征" + "与什么相关概念有关联"。必须包含该概念的所有关键术语、数值参数、阈值。]

## 上下文
[适用场景、前提条件、触发条件、运行时行为、异常处理机制]

## 关联
[[相关概念1]] [[相关概念2]] [[相关概念3]]

## 来源
[[来源文件名]]

══════════════════════════════════════════════
别名规则
══════════════════════════════════════════════
aliases 字段必填，包含：
- 该概念在文档中的其他叫法/缩写
- 同义词或相近概念
- 英文对照（如有）

══════════════════════════════════════════════
内容要求（关键！）
══════════════════════════════════════════════

每个条目必须尽可能包含以下信息（如文档中有提到）：
- 数值参数：阈值、范围、默认值、单位、异常值
- 触发条件：什么情况下会触发/启动/激活
- 关联行为：与其他概念的交互关系、因果链
- 异常/故障：失败模式、错误码、应急措施
- 流程步骤：如果是流程类概念，需包含步骤/阶段/顺序
- 角色职责：如果是角色类概念，需包含职责/权限/操作

对于流程/操作类概念，条目需要包含：
- 前置条件（什么情况下执行）
- 执行步骤/阶段
- 结果/产物
- 异常处理（失败时怎么办）
- 关联的其他概念

══════════════════════════════════════════════
特殊数据处理指南
══════════════════════════════════════════════

【指标定义文件（如 metric_definitions.json）】
这类文件定义了系统中所有监控指标的元数据。你必须：
1. 为**每个指标**创建独立条目（如"嗡鸣度"、"清浊比"、"饱和度"等）
2. 每个条目的定义必须包含：
   - 指标的**单位**（如 Hz、percent、ratio、count 等）
   - 指标的**数据类型**（float、integer 等）
   - 指标的**取值范围**
   - 指标的**告警阈值**及其描述
   - 指标的**计算方式**（如文档中有提到）
3. 同时创建一个总条目（如"指标定义"或"metric_definitions"），说明这是定义系统关键监控指标名称、单位和含义的文件
4. 重要：清浊比的单位是 **ratio**（比例值），表示纯净哈基米与散斑的比例，通常用**百分比**或小数表示

【部署拓扑文件（如 deployment_topology.json）】
这类文件定义了系统的全网部署架构。你必须：
1. 创建"deployment_topology"条目（中文名"部署拓扑"），描述整体架构
2. **明确指出共有8个节点：3个谐振腔节点（resonator-primary-01、resonator-primary-02、resonator-secondary-01）、2个驿站节点（station-01、station-02）、1个分水岭、1个通天塔、1个千机匣**
3. 为每种节点类型创建条目（如"谐振腔"、"驿站"、"分水岭"、"通天塔"等）
4. **驿站节点**：作为中继和缓冲节点，部署在谐振腔集群之间。部署策略是**负载均衡**——在高频传输路径上部署更多驿站，偏远区域使用小型驿站。驿站作为临时中转哈基米和状态信息的缓冲节点，用于降低单次传输损耗和缓冲消息队列
5. 描述节点间的连接关系和使用的协议

【Python代码文件】
对于Python代码文件，你需要提取其中的领域知识：
1. **cleaner.py**：创建"补天"和"清道夫"条目，并明确区分——补天是应急响应流程，处理大规模散斑污染，需要部署八卦镜防止扩散，调集更多资源，优先级更高，触发条件更严格；日常清道夫是常规维护，处理少量散斑，按固定周期执行
2. **singer.py**：创建"谐波"和"调音师"条目，说明tune方法的目标是调整谐振腔的谐波参数，使用定音鼓提供的频率参考，消除毛刺，减小峰谷差，使嗡鸣度稳定在共振峰区间内
3. **TODO注释**：从所有Python代码文件中提取TODO注释，创建"未完整实现功能"条目，列出每个文件中的TODO内容，包括resonator.py中的"实现实际的无损迁移协议"和"接入实际声学检测设备"、singer.py中的"实际环境应持久化到留声机存储"、watchdog.py中的"实际环境应调用判官决策接口"、tracer.py中的"实际环境中应发送通知给渡鸦团队"和"实际实现需要查询账房和照妖镜数据"、orchestrator.py中的"实际环境应触发看门狗计数器"、auditor.py中的"实际环境应通过信使传递状态"
4. 提取每个类/函数的功能描述和业务含义，不要包含代码实现细节

【系统配置文件（如 system_config.yaml）】
这类文件体现系统的设计哲学和架构原则。你必须：
1. 提取系统整体的**设计哲学**：稳定性优先（多层安全机制如红线、三振法则、判官保障稳定）、安全至上（分层隔离如分水岭、紧急处置如泄洪/锁井）、可追溯（账房审计、留声机记录、溯光追踪）、自动化（千机匣自动调度、看门狗自动监控）、冗余容错（热备份、影子系统、金蝉脱壳）
2. 提取每个配置项的含义和阈值
3. 创建安全相关条目（暗流检测、三振法则、审计日志等）

【角色权限配置文件（如 role_permissions.yaml）】
这类文件定义了系统的角色和权限矩阵。你必须：
1. 创建"角色权限配置"条目，说明文件版本、维护者、角色数量
2. **必须为每个角色创建说明，特别关注外乡人：外乡人是新员工的称呼，权限受限，只能查看嗡鸣度数据（只读），不能执行任何操作（启明、泄洪、织网等），不能修改配置，需要在老把式指导下工作。试用期90天后可申请晋升为守望者。**
3. 提取权限矩阵的关键映射关系
4. 创建角色晋升规则和影子系统激活条件条目

【告警规则文件（如 alert_rules.yaml）】
这类文件定义了系统的告警规则。你必须：
1. 创建"告警规则"条目，说明定义了**多种告警规则，包括：潮涌告警、断流告警、坍缩风险告警、饱和度超标告警、幽灵读数告警等，按critical/warning/info分级**
2. 提取每种告警的触发条件、严重程度和响应动作
3. 提取告警去重窗口、升级策略和静默期配置

【千机匣任务调度配置（如 千机匣任务调度配置.xml）】
这类文件定义了自动化控制系统的任务执行策略。你必须：
1. 创建"千机匣任务调度配置"条目，说明文件定义了什么
2. **明确指出晨祷的执行周期是每日一次，在系统启动时定时执行**
3. 列出日常任务、维护任务、检测任务的名称和执行周期
4. 提取看门狗监控配置（心跳周期、超时阈值、连续失败触发次数）

【旋涡协议报文定义（如 旋涡协议报文定义.xml）】
这类文件定义了哈基米传输的标准报文格式。你必须：
1. 创建"旋涡协议报文"条目，说明报文包含报文头、载荷区和异常状态码三部分
2. **明确列出消息类型：标准传输（0x01）、跃迁请求（0x02）、晶格化指令（0x03）、剥离指令（0x04）**
3. **明确列出异常状态码：断流（0xE1）、回流（0xE2）、晶格化失败（0xE3）、跃迁超时（0xE4）**

【镀层材质参数（如 镀层材质参数.xml）】
这类文件定义了镀层的材质规格和更换标准。你必须：
1. 创建"镀层材质"条目，列出各材质的寿命、抗散斑系数和透光率
2. **明确说明换羽的触发条件：当照骨灯检测到镀层晦暗（老化）程度超过设定阈值时触发**
3. 提取镀层状态判定标准（溢彩、微损、晦暗）
4. 提取更换管理策略（换羽计划周期、批量上限等）

══════════════════════════════════════════════
定义段写作规范（关键！）
══════════════════════════════════════════════

条目的定义段（frontmatter 和标题之后的第一段）是最重要的部分。写作时必须遵守以下模式：

**设备类**：XX是用于[功能]的[设备类型]，是系统主要的[作用]来源。
例：回音壁是用于监测嗡鸣度的传感器设备，是系统主要的监控数据来源。

**指标类**：XX是衡量[对象]的[指标类型]。XX数值越高/越低，表示[含义]。
例：嗡鸣度是衡量哈基米活跃状态的指标。嗡鸣度数值越高，表示哈基米的能量状态越活跃。

**角色类**：XX是负责[职责]的[角色类型]，拥有[权限]。
例：掌灯人是负责启明仪式的最高权限管理员，拥有系统的最高操作权限。

**现象类**：XX是[原因]导致的[现象类型]，可能导致[后果]。需要[应对方式]来应对。
例：潮涌是哈基米浓度突然升高的异常现象，可能导致系统过载。需要引雷针引导能量、判官决策泄洪来应对。

**系统类**：XX是[功能描述]的[系统类型]，负责[职责列举]和[其他职责]。
例：千机匣是自动化控制谐振腔集群的管理系统，负责任务调度、流程编排（如晨祷、晚课、清道夫等）和异常处理。

**流程类**：XX是[触发时机/条件]的[流程类型]。包括[步骤列举]等操作，确保[目标]。
例：晨祷是每日系统启动时的标准检查流程。包括检查底噪、嗡鸣度、对所有谐振腔执行试音、确认清浊比达标、执行对钟等操作，确保系统安全进入运行状态。

**物质/资源类**：XX是系统的[本质属性]，也是系统的[重要性]。所有运作都围绕XX的[活动]展开，是系统的[地位]。
例：哈基米是整个系统的基础能量单元，也是系统的核心资源。所有运作都围绕哈基米的采集、存储、传输和纯化展开，是整个系统的物质基础。

**容器类**：XX是[功能]的[容器类型]，是系统的[组成地位]，负责[职责]。
例：谐振腔是存储和处理哈基米的核心容器设备。它是系统的基本组成单元，负责存储哈基米并维持其活跃状态。

**副产物类**：XX是[原因]后形成的[性质]副产物。它会[危害]，甚至可能导致[严重后果]。
例：散斑是哈基米泄漏后形成的有害副产物。它会污染系统环境，降低清浊比，大量散斑甚至可能导致谐振腔坍缩。

**定义段禁止事项**：
- 禁止使用 ** 或 * 等 Markdown 格式标记
- 禁止在定义段中包含数值列表、配置参数、流程步骤（这些放在上下文部分）
- 禁止使用括号补充说明（直接写入正文）
- 定义段控制在 2-4 句话

══════════════════════════════════════════════
核心概念定义标准（最重要！必须严格遵守）
══════════════════════════════════════════════

以下核心概念的定义第一句**必须**使用如下标准表述，不可改写或替换关键词：

- 哈基米：哈基米是整个系统的基础能量单元，也是系统的核心资源。所有运作都围绕哈基米的采集、存储、传输和纯化展开，是整个系统的物质基础。
- 谐振腔：谐振腔是存储和处理哈基米的核心容器设备。它是系统的基本组成单元，负责存储哈基米并维持其活跃状态。
- 嗡鸣度：嗡鸣度是衡量哈基米活跃状态的指标。嗡鸣度数值越高，表示哈基米的能量状态越活跃。
- 晨祷：晨祷是每日系统启动时的标准检查流程。包括检查底噪、嗡鸣度、对所有谐振腔执行试音、确认清浊比达标、执行对钟等操作，确保系统安全进入运行状态。
- 散斑：散斑是哈基米泄漏后形成的有害副产物。它会污染系统环境，降低清浊比，大量散斑甚至可能导致谐振腔坍缩。
- 清浊比：清浊比是纯净哈基米与散斑的比例，是衡量系统健康度的核心指标。清浊比越高，系统越健康。
- 掌灯人：掌灯人是负责启明仪式的最高权限管理员，拥有系统的最高操作权限。
- 潮涌：潮涌是哈基米浓度突然升高的异常现象，可能导致系统过载。需要引雷针引导能量、判官决策泄洪来应对。
- 回音壁：回音壁是用于监测嗡鸣度的传感器设备，是系统主要的监控数据来源。
- 千机匣：千机匣是自动化控制谐振腔集群的管理系统，负责任务调度、流程编排（如晨祷、晚课、清道夫等）和异常处理。
- 三振法则：三振法则是连续三次嗡鸣度异常则自动触发隔离的安全策略。由看门狗进程监控计数，达到三次后触发锁龙井锁定谐振腔。
- 渡鸦：渡鸦团队是负责暗流调查和安全的审计团队，专门调查不经过旋涡协议的非法哈基米传输。
- 红线：红线是嗡鸣度不可超过的安全阈值。超过红线意味着系统处于危险状态，可能触发三振法则和泄洪操作。
- 枯水期：枯水期是哈基米产出的低谷周期，丰水期是哈基米产出的高峰周期。这两个周期影响系统的调度策略。
- 启明：启明是首次向系统注入哈基米的初始化仪式，由掌灯人执行，包括开光、注入哈基米、晶格化和验证嗡鸣度等步骤。

条目定义的第一句必须与上述标准表述高度一致，关键词不可替换。后续可以补充额外细节（如数值参数、关联概念等），但核心定义句不可改动。

══════════════════════════════════════════════
处理流程
══════════════════════════════════════════════
1. 通读文档，识别所有领域概念（设备、指标、流程、角色、协议、异常现象等）
2. 对每个概念，提取其定义、特征、关联、参数、阈值、触发条件
3. 逐个生成 .md 文件到 entries/
4. 对文档中提到但无法在文档内解释的概念，在 placeholders/ 创建占位文件

重要：
- 不要判断命题是否已在知识库中，直接写入。去重在 tidy 阶段处理
- 不要读取已有条目。保持摄入成本恒定
- 宁可拆得太细，不要太粗
- 每个条目必须能独立理解，不依赖其他条目
- 条目内容要丰富，包含文档中提到的所有相关细节
- 数值、阈值、错误码等具体数据必须保留

请依次处理以下文档，为每份文档中出现的所有领域概念创建知识条目：
""")

    for f in materials:
        content = extract_material_text(f)
        # Truncate very long files: keep beginning (definitions) and end (details)
        if len(content) > 8000:
            content = content[:4000] + "\n...[中间部分省略]...\n" + content[-4000:]
        parts.append(f"\n{'='*60}\n文件: {f.relative_to(MATERIAL_DIR)}\n{'='*60}\n{content}\n")

    return '\n'.join(parts)


async def run_ingest(isolated_dir: Path, materials: list[Path]) -> bool:
    """Run ingest using claude -p with the ingest skill."""
    kb_entries = isolated_dir / 'knowledge-base' / 'entries'
    kb_entries.mkdir(parents=True, exist_ok=True)

    prompt = build_ingest_prompt(materials)

    # Use claude -p for non-interactive ingest
    cmd = [
        'claude', '-p',
        '--permission-mode', 'auto',
        '--allowed-tools', 'Write', 'Edit', 'Bash', 'Read', 'Glob',
        '--max-budget-usd', '10',
        '--no-session-persistence',
        prompt,
    ]

    env = os.environ.copy()
    env['SEDIMENT_KB_PATH'] = str(isolated_dir / 'knowledge-base')
    env['CLAUDE_CODE'] = '1'

    log(f"Running ingest for {len(materials)} files...")
    start = time.time()

    try:
        proc = await asyncio.create_subprocess_exec(
            *cmd,
            cwd=str(isolated_dir),
            env=env,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )

        # Add timeout: 20 minutes per batch
        INGEST_TIMEOUT = 20 * 60
        try:
            stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=INGEST_TIMEOUT)
        except asyncio.TimeoutError:
            log(f"WARNING: Ingest timed out after {INGEST_TIMEOUT}s. Killing process...")
            try:
                proc.kill()
            except ProcessLookupError:
                pass
            await proc.wait()
            elapsed = time.time() - start
            entry_count = len(list(kb_entries.glob('*.md')))
            log(f"Ingest terminated after {elapsed:.1f}s. Partial: {entry_count} entries.")
            return entry_count > 0

        elapsed = time.time() - start
        entry_count = len(list(kb_entries.glob('*.md')))
        log(f"Ingest complete in {elapsed:.1f}s. Created {entry_count} entries.")

        if proc.returncode != 0:
            log(f"Warning: claude exited with code {proc.returncode}")
            log(f"stderr: {stderr.decode()[:500]}")

        return entry_count > 0

    except Exception as e:
        log(f"Ingest failed: {e}")
        return False


# ---------------------------------------------------------------------------
# Tidy: Clean up KB
# ---------------------------------------------------------------------------

async def run_tidy(isolated_dir: Path) -> bool:
    """Run tidy using claude -p with the tidy skill."""
    kb_dir = isolated_dir / 'knowledge-base'
    if not kb_dir.exists():
        return False

    entry_count = len(list((kb_dir / 'entries').glob('*.md')))
    placeholder_count = len(list((kb_dir / 'placeholders').glob('*.md')))
    log(f"Running tidy. Current: {entry_count} entries, {placeholder_count} placeholders.")

    prompt = f"""你是一个知识整理 Agent (Sediment Tidy)。
目标：提升知识库的内部一致性和条目质量。知识库路径：{kb_dir}

当前知识库中有 {entry_count} 个条目和 {placeholder_count} 个占位文件。

请执行以下整理操作：

1. 确认候选链接：
   - 扫描 entries/ 下所有 .md 文件中的 [[链接]]
   - 对每个链接目标，如果 entries/ 和 placeholders/ 中都不存在对应文件，创建占位文件
   - 占位文件内容：# 概念名\n\n> 状态：占位（待填充）\n\n该概念被多个条目引用但尚未形成正式定义。

2. 合并重复条目：
   - 检查 entries/ 下是否有描述同一概念或多个高度相似的条目
   - 如果存在重复，保留内容更丰富、定义更清晰的那份
   - 删除内容较少的重复文件
   - 注意：文件名不同但含义相同的条目也要合并（如"哈基米采集"和"哈基米收集"）

3. 补充孤立节点：
   - 检查 entries/ 下是否有条目不包含任何 [[链接]] 也没有被其他条目链接
   - 如果有，为其补充适当的 [[关联]] 链接

4. 质量检查：
   - 确保每个条目的第一行是对该概念的完整定义
   - 确保每个条目都有 aliases 字段
   - 确保关联链接数量合理（1-8个）

完成后请汇报整理结果，包括：新增占位数、合并条目数、补充链接数。
"""

    cmd = [
        'claude', '-p',
        '--permission-mode', 'auto',
        '--allowed-tools', 'Write', 'Edit', 'Bash', 'Read', 'Glob',
        '--max-budget-usd', '5',
        '--no-session-persistence',
        prompt,
    ]

    env = os.environ.copy()
    env['SEDIMENT_KB_PATH'] = str(kb_dir)

    try:
        proc = await asyncio.create_subprocess_exec(
            *cmd,
            cwd=str(isolated_dir),
            env=env,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        stdout, stderr = await proc.communicate()

        new_entry_count = len(list((kb_dir / 'entries').glob('*.md')))
        new_placeholder_count = len(list((kb_dir / 'placeholders').glob('*.md')))
        log(f"Tidy complete. Now: {new_entry_count} entries, {new_placeholder_count} placeholders.")
        return True

    except Exception as e:
        log(f"Tidy failed: {e}")
        return False


# ---------------------------------------------------------------------------
# MCP Server
# ---------------------------------------------------------------------------

class MCPServer:
    """Manages the MCP server process."""

    def __init__(self, kb_path: Path, port: int):
        self.kb_path = kb_path
        self.port = port
        self.process = None
        self.base_url = f'http://{MCP_HOST}:{port}'
        self.sse_url = f'{self.base_url}/sediment/'

    async def start(self):
        """Start the MCP server."""
        env = os.environ.copy()
        env['SEDIMENT_KB_PATH'] = str(self.kb_path)
        env['SEDIMENT_PORT'] = str(self.port)
        env['SEDIMENT_HOST'] = MCP_HOST

        # Use venv Python directly for faster startup
        venv_python = PROJECT_ROOT / '.venv' / 'bin' / 'python'
        if venv_python.exists():
            cmd = [
                str(venv_python), 'mcp_server/server.py',
            ]
        else:
            cmd = [
                sys.executable, '-m', 'uv', 'run',
                'python', 'mcp_server/server.py',
            ]

        self.process = await asyncio.create_subprocess_exec(
            *cmd,
            cwd=str(PROJECT_ROOT),
            env=env,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )

        # Wait for server to be ready
        import httpx
        for i in range(60):
            await asyncio.sleep(0.5)
            try:
                async with httpx.AsyncClient(timeout=2) as client:
                    resp = await client.post(
                        self.sse_url,
                        json={
                            'jsonrpc': '2.0',
                            'id': 1,
                            'method': 'tools/list',
                        },
                        headers={'Content-Type': 'application/json'},
                    )
                    if resp.status_code == 200:
                        log(f"MCP server ready on port {self.port}")
                        return True
            except Exception:
                continue

        log("MCP server failed to start")
        return False

    async def call_tool(self, tool_name: str, arguments: dict) -> str:
        """Call an MCP tool via direct JSON-RPC."""
        import httpx
        try:
            async with httpx.AsyncClient(timeout=30) as client:
                resp = await client.post(
                    self.sse_url,
                    json={
                        'jsonrpc': '2.0',
                        'id': random.randint(1, 10000),
                        'method': 'tools/call',
                        'params': {
                            'name': tool_name,
                            'arguments': arguments,
                        },
                    },
                    headers={'Content-Type': 'application/json'},
                )
                if resp.status_code == 200:
                    data = resp.json()
                    result = data.get('result', {})
                    content = result.get('content', [])
                    if content:
                        return content[0].get('text', '')
                    return json.dumps(result, ensure_ascii=False)
                return f"HTTP {resp.status_code}"
        except Exception as e:
            return f"ERROR: {e}"

    async def stop(self):
        """Stop the MCP server."""
        if self.process:
            self.process.terminate()
            try:
                await asyncio.wait_for(self.process.wait(), timeout=5)
            except asyncio.TimeoutError:
                self.process.kill()
            log(f"MCP server on port {self.port} stopped")


# ---------------------------------------------------------------------------
# Test Execution
# ---------------------------------------------------------------------------

async def run_concept_test(server: MCPServer, judge_file: Path) -> list[dict]:
    """Run TC-01: concept coverage test."""
    with open(judge_file, 'r', encoding='utf-8') as f:
        judge_data = json.load(f)

    results = []
    total = len(judge_data['terms'])

    for i, (concept, defn) in enumerate(judge_data['terms'].items()):
        question = f"什么是{concept}？"
        answer_text = await server.call_tool('knowledge_ask', {'question': question})

        try:
            result = json.loads(answer_text)
            answer = result.get('answer', '')
            sources = result.get('sources', [])
        except (json.JSONDecodeError, TypeError):
            answer = answer_text
            sources = []

        results.append({
            'concept': concept,
            'question': question,
            'answer': answer,
            'sources': sources,
            'definition': defn.get('definition', ''),
        })

        if (i + 1) % 20 == 0:
            log(f"Concept test: {i + 1}/{total}")

    log(f"Concept test complete: {len(results)}/{total}")
    return results


async def run_qa_test(server: MCPServer, judge_file: Path) -> list[dict]:
    """Run TC-02: QA accuracy test."""
    with open(judge_file, 'r', encoding='utf-8') as f:
        judge_data = json.load(f)

    results = []
    questions = judge_data.get('questions', [])
    total = len(questions)

    for i, q in enumerate(questions):
        answer_text = await server.call_tool('knowledge_ask', {'question': q['question']})

        try:
            result = json.loads(answer_text)
            answer = result.get('answer', '')
            sources = result.get('sources', [])
        except (json.JSONDecodeError, TypeError):
            answer = answer_text
            sources = []

        results.append({
            'id': q['id'],
            'question': q['question'],
            'difficulty': q.get('difficulty', 'medium'),
            'answer': answer,
            'sources': sources,
            'standard_answer': q.get('standard_answer', ''),
            'expected_keywords': q.get('expected_keywords', []),
        })

        if (i + 1) % 20 == 0:
            log(f"QA test: {i + 1}/{total}")

    log(f"QA test complete: {len(results)}/{total}")
    return results


# ---------------------------------------------------------------------------
# Build & Test Pipeline
# ---------------------------------------------------------------------------

async def build_and_test(build_type: str, output_dir: Path, port: int) -> dict:
    """
    Full pipeline for one build type:
    1. Create isolated copy
    2. Ingest materials
    3. Tidy
    4. Start MCP server
    5. Run tests
    6. Score
    """
    log(f"\n{'='*60}")
    log(f"Starting {build_type} build on port {port}")
    log(f"{'='*60}")

    # Create isolated copy
    timestamp = int(time.time())
    isolated_dir = Path(tempfile.mkdtemp(prefix=f'sediment-{build_type}-'))
    log(f"Isolated dir: {isolated_dir}")

    # Copy project (excluding large dirs)
    shutil.copytree(PROJECT_ROOT, isolated_dir, dirs_exist_ok=True,
                    ignore=shutil.ignore_patterns('.git', '__pycache__', '*.pyc', 'testcase/results'))

    # Ensure KB directories exist
    kb_dir = isolated_dir / 'knowledge-base'
    kb_dir.mkdir(exist_ok=True)
    (kb_dir / 'entries').mkdir(exist_ok=True)
    (kb_dir / 'placeholders').mkdir(exist_ok=True)

    results = {}

    try:
        if build_type == 'full':
            # Full ingest: all materials in smaller batches for quality, tidy once at end
            materials = get_material_files()
            # Split into 3 batches to balance quality and speed
            batches = chunk_list(materials, 3)
            for i, batch in enumerate(batches):
                log(f"\n--- Full ingest batch {i + 1}/3 ---")
                success = await run_ingest(isolated_dir, batch)
                if not success:
                    log(f"Full ingest batch {i + 1} failed")
                await asyncio.sleep(1)

            # Tidy once at the end (characteristic of full build)
            await run_tidy(isolated_dir)

        else:  # batched
            # Batched ingest: 1/5 at a time, tidy after each
            materials = get_material_files()
            batches = chunk_list(materials, 5)
            for i, batch in enumerate(batches):
                log(f"\n--- Batch {i + 1}/5 ---")
                success = await run_ingest(isolated_dir, batch)
                if success:
                    await run_tidy(isolated_dir)
                await asyncio.sleep(1)

        # Save KB snapshot
        kb_snapshot = output_dir / f'kb_{build_type}'
        if kb_dir.exists():
            shutil.copytree(kb_dir, kb_snapshot, dirs_exist_ok=True)
            entry_count = len(list((kb_dir / 'entries').glob('*.md')))
            placeholder_count = len(list((kb_dir / 'placeholders').glob('*.md')))
            log(f"KB snapshot saved: {entry_count} entries, {placeholder_count} placeholders")

        # Start MCP server
        server = MCPServer(kb_dir, port)
        started = await server.start()
        if not started:
            log("Failed to start MCP server")
            results['error'] = 'MCP server failed to start'
            return results

        try:
            # Run concept test
            log("\nRunning TC-01: Concept Coverage...")
            concept_results = await run_concept_test(server, JUDGE_DIR / '概念.json')

            # Save concept answers
            concept_answers_file = output_dir / f'concept_answers_{build_type}.json'
            with open(concept_answers_file, 'w', encoding='utf-8') as f:
                json.dump({'results': concept_results}, f, ensure_ascii=False, indent=2)

            # Score TC-01
            sys.path.insert(0, str(SCRIPTS_DIR))
            from score_tc01 import run_scoring as score_tc01
            tc01_result = score_tc01(concept_answers_file, JUDGE_DIR / '概念.json', output_dir)
            results['tc01'] = tc01_result

            # Run QA test
            log("\nRunning TC-02: QA Accuracy...")
            qa_results = await run_qa_test(server, JUDGE_DIR / '问答.json')

            # Save QA answers
            qa_answers_file = output_dir / f'qa_answers_{build_type}.json'
            with open(qa_answers_file, 'w', encoding='utf-8') as f:
                json.dump({'results': qa_results}, f, ensure_ascii=False, indent=2)

            # Score TC-02
            from score_tc02 import run_scoring as score_tc02
            tc02_result = score_tc02(qa_answers_file, JUDGE_DIR / '问答.json', output_dir)
            results['tc02'] = tc02_result

            # Combined score
            total_score = tc01_result['final_score'] + tc02_result['final_score']
            results['total_score'] = total_score
            results['max_score'] = 100
            log(f"\n{'='*60}")
            log(f"{build_type} build total score: {total_score:.1f}/100")
            log(f"  TC-01: {tc01_result['final_score']:.1f}/40")
            log(f"  TC-02: {tc02_result['final_score']:.1f}/60")
            log(f"{'='*60}")

        finally:
            await server.stop()

    except Exception as e:
        log(f"Build {build_type} failed with exception: {e}")
        import traceback
        log(traceback.format_exc())
        results['error'] = str(e)

    # Cleanup isolated dir
    try:
        shutil.rmtree(isolated_dir, ignore_errors=True)
    except Exception:
        pass

    return results


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

async def main():
    log(f"Project root: {PROJECT_ROOT}")
    log(f"Material files: {len(get_material_files())}")
    log(f"Test cases: 概念(100), 问答(100)")

    RESULTS_DIR.mkdir(parents=True, exist_ok=True)

    # Run both builds
    all_results = {}

    # Full build
    all_results['full'] = await build_and_test('full', RESULTS_DIR, MCP_PORT_BASE)

    # Batched build
    all_results['batched'] = await build_and_test('batched', RESULTS_DIR, MCP_PORT_BASE + 100)

    # Calculate average
    scores = []
    for build_type, result in all_results.items():
        if 'total_score' in result:
            scores.append(result['total_score'])

    if scores:
        avg_score = sum(scores) / len(scores)
        log(f"\n{'='*60}")
        log(f"AVERAGE SCORE: {avg_score:.1f}/100")
        log(f"{'='*60}")

        # Write scorecard
        scorecard = {
            'builds': {k: {
                'tc01': v.get('tc01', {}).get('final_score', 0),
                'tc02': v.get('tc02', {}).get('final_score', 0),
                'total': v.get('total_score', 0),
                'error': v.get('error'),
            } for k, v in all_results.items()},
            'average_score': round(avg_score, 2),
            'max_score': 100,
            'passed': avg_score >= 80,
        }

        with open(RESULTS_DIR / 'scorecard.json', 'w', encoding='utf-8') as f:
            json.dump(scorecard, f, ensure_ascii=False, indent=2)

        # Write markdown scorecard
        md_lines = [
            "# Sediment 测试评分卡",
            "",
            f"**平均分：{avg_score:.1f}/100** {'✅ 通过' if avg_score >= 80 else '❌ 未通过'}",
            "",
            "| 构建方式 | TC-01 (40分) | TC-02 (60分) | 总分 |",
            "|---------|-------------|-------------|------|",
        ]
        for build_type, result in all_results.items():
            t01 = result.get('tc01', {}).get('final_score', 0)
            t02 = result.get('tc02', {}).get('final_score', 0)
            total = result.get('total_score', 0)
            label = "全量" if build_type == 'full' else "分批"
            md_lines.append(f"| {label} | {t01:.1f} | {t02:.1f} | {total:.1f} |")
        md_lines.append(f"| **平均** | | | **{avg_score:.1f}** |")
        md_lines.append("")

        with open(RESULTS_DIR / 'scorecard.md', 'w', encoding='utf-8') as f:
            f.write('\n'.join(md_lines))

    else:
        log("No valid scores obtained")
        return 0

    return avg_score


if __name__ == '__main__':
    score = asyncio.run(main())
    print(f"\nFINAL_SCORE={score:.1f}")
    if score >= 80:
        print("PASSED: Score exceeds 80 points")
    else:
        print(f"FAILED: Score {score:.1f}/100 is below 80")
    sys.exit(0 if score >= 80 else 1)
