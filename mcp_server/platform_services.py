from __future__ import annotations

import base64
import difflib
import hashlib
import ipaddress
import re
import shutil
import tempfile
from collections.abc import Iterable
from pathlib import Path
from typing import Any

from docx import Document
from pptx import Presentation

from mcp_server.kb import (
    audit_kb,
    index_config,
    inventory,
    resolve_kb_document_path,
    split_frontmatter,
    validate_entry,
    validate_index,
)
from mcp_server.platform_store import PlatformStore

FORMAL_STATUSES = {"fact", "inferred", "disputed"}
ALLOWED_MIME_TYPES = {
    "text/plain",
    "text/markdown",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


def ensure_platform_state(
    *,
    store: PlatformStore,
    state_dir: str | Path,
    uploads_dir: str | Path,
    workspaces_dir: str | Path,
) -> None:
    Path(state_dir).mkdir(parents=True, exist_ok=True)
    Path(uploads_dir).mkdir(parents=True, exist_ok=True)
    Path(workspaces_dir).mkdir(parents=True, exist_ok=True)
    store.init()


def build_submission_hash(*parts: str) -> str:
    digest = hashlib.sha256()
    for part in parts:
        digest.update(part.encode("utf-8"))
        digest.update(b"\n---\n")
    return digest.hexdigest()


def normalize_name(raw_value: str, *, fallback: str) -> str:
    cleaned = re.sub(r"\s+", " ", raw_value).strip()
    return cleaned or fallback


def parse_trusted_proxy_cidrs(raw_value: str | None) -> tuple[ipaddress._BaseNetwork, ...]:
    if not raw_value:
        return ()
    networks: list[ipaddress._BaseNetwork] = []
    for item in raw_value.split(","):
        token = item.strip()
        if not token:
            continue
        networks.append(ipaddress.ip_network(token, strict=False))
    return tuple(networks)


def detect_submitter_ip(
    headers: dict[str, str] | None,
    client_host: str | None,
    *,
    trust_proxy_headers: bool = False,
    trusted_proxy_cidrs: Iterable[ipaddress._BaseNetwork] = (),
) -> str:
    headers = headers or {}
    normalized_client = normalize_ip(client_host)
    if not trust_proxy_headers:
        return normalized_client

    trusted_networks = tuple(trusted_proxy_cidrs)
    if trusted_networks and not ip_in_networks(normalized_client, trusted_networks):
        return normalized_client

    forwarded = normalize_ip(headers.get("x-forwarded-for", "").split(",")[0].strip())
    real_ip = normalize_ip(headers.get("x-real-ip", "").strip())
    return forwarded or real_ip or normalized_client


def normalize_ip(raw_value: str | None) -> str:
    value = (raw_value or "").strip()
    if not value:
        return "unknown"
    candidate = value
    if value.startswith("[") and "]" in value:
        candidate = value[1 : value.index("]")]
    elif value.count(":") == 1 and "." in value:
        host, _, port = value.rpartition(":")
        if port.isdigit():
            candidate = host
    try:
        return str(ipaddress.ip_address(candidate))
    except ValueError:
        return value


def ip_in_networks(
    client_ip: str,
    networks: Iterable[ipaddress._BaseNetwork],
) -> bool:
    try:
        address = ipaddress.ip_address(client_ip)
    except ValueError:
        return False
    return any(address in network for network in networks)


def extract_upload_text(file_path: str | Path, mime_type: str) -> str:
    path = Path(file_path)
    if mime_type == "text/plain":
        return path.read_text(encoding="utf-8", errors="replace")
    if mime_type == "text/markdown":
        return path.read_text(encoding="utf-8", errors="replace")
    if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        doc = Document(path)
        paragraphs = [item.text.strip() for item in doc.paragraphs if item.text.strip()]
        return "\n\n".join(paragraphs).strip()
    if mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        presentation = Presentation(path)
        chunks: list[str] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_lines: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = str(shape.text).strip()
                    if text:
                        slide_lines.append(text)
            if slide_lines:
                chunks.append(f"Slide {slide_index}\n" + "\n".join(slide_lines))
        return "\n\n".join(chunks).strip()
    raise ValueError(f"Unsupported mime type: {mime_type}")


def submit_text(
    *,
    store: PlatformStore,
    title: str,
    content: str,
    submitter_name: str,
    submitter_ip: str,
    submission_type: str = "text",
    submitter_user_id: str | None = None,
    notes: str | None = None,
    rate_limit_count: int = 1,
    rate_limit_window_seconds: int = 60,
    max_text_chars: int = 20_000,
    dedupe_window_seconds: int = 86_400,
) -> dict[str, Any]:
    cleaned_title = normalize_name(title, fallback="未命名提交")
    cleaned_content = content.strip()
    cleaned_submitter = normalize_name(submitter_name, fallback="Anonymous")
    if not cleaned_content:
        raise ValueError("submission content must not be empty")
    if len(cleaned_content) > max(1, max_text_chars):
        raise ValueError("submission content is too large")
    if store.submission_count_for_ip(
        submitter_ip,
        within_seconds=max(1, rate_limit_window_seconds),
    ) >= max(1, rate_limit_count):
        raise PermissionError("submission rate limit exceeded for this IP")
    dedupe_hash = build_submission_hash(cleaned_title, cleaned_content)
    if dedupe_window_seconds > 0 and store.find_recent_submission_by_hash(
        dedupe_hash,
        within_seconds=dedupe_window_seconds,
    ):
        raise FileExistsError("duplicate submission already exists in the review buffer")

    record = store.create_submission(
        submission_type=submission_type,
        title=cleaned_title,
        raw_text=cleaned_content,
        extracted_text=cleaned_content,
        stored_file_path=None,
        mime_type="text/plain",
        submitter_name=cleaned_submitter,
        submitter_ip=submitter_ip,
        submitter_user_id=submitter_user_id,
        dedupe_hash=dedupe_hash,
        notes=notes,
    )
    store.add_audit_log(
        actor_name=cleaned_submitter,
        actor_id=submitter_user_id,
        actor_role="contributor",
        action="submission.create_text",
        target_type="submission",
        target_id=record["id"],
        details={"submission_type": submission_type},
    )
    return record


def submit_document(
    *,
    store: PlatformStore,
    uploads_dir: str | Path,
    filename: str,
    mime_type: str,
    file_bytes: bytes,
    submitter_name: str,
    submitter_ip: str,
    submitter_user_id: str | None = None,
    notes: str | None = None,
    rate_limit_count: int = 1,
    rate_limit_window_seconds: int = 60,
    max_upload_bytes: int = 10 * 1024 * 1024,
    dedupe_window_seconds: int = 86_400,
) -> dict[str, Any]:
    if mime_type not in ALLOWED_MIME_TYPES:
        raise ValueError("unsupported upload type")
    if store.submission_count_for_ip(
        submitter_ip,
        within_seconds=max(1, rate_limit_window_seconds),
    ) >= max(1, rate_limit_count):
        raise PermissionError("submission rate limit exceeded for this IP")
    if len(file_bytes) > max(1, max_upload_bytes):
        raise ValueError("uploaded file is too large")

    cleaned_submitter = normalize_name(submitter_name, fallback="Anonymous")
    safe_filename = sanitize_filename(filename or "upload.bin")
    temp_submission_id = hashlib.sha256(file_bytes).hexdigest()[:16]
    stored_path = Path(uploads_dir) / f"{temp_submission_id}_{safe_filename}"
    stored_path.parent.mkdir(parents=True, exist_ok=True)
    stored_path.write_bytes(file_bytes)
    extracted_text = extract_upload_text(stored_path, mime_type)
    if not extracted_text:
        raise ValueError("could not extract text from uploaded document")
    dedupe_hash = build_submission_hash(safe_filename, mime_type, extracted_text)
    if dedupe_window_seconds > 0 and store.find_recent_submission_by_hash(
        dedupe_hash,
        within_seconds=dedupe_window_seconds,
    ):
        stored_path.unlink(missing_ok=True)
        raise FileExistsError("duplicate submission already exists in the review buffer")

    record = store.create_submission(
        submission_type="document",
        title=safe_filename,
        raw_text=base64.b64encode(file_bytes).decode("ascii"),
        extracted_text=extracted_text,
        stored_file_path=str(stored_path),
        mime_type=mime_type,
        submitter_name=cleaned_submitter,
        submitter_ip=submitter_ip,
        submitter_user_id=submitter_user_id,
        dedupe_hash=dedupe_hash,
        notes=notes,
    )
    store.add_audit_log(
        actor_name=cleaned_submitter,
        actor_id=submitter_user_id,
        actor_role="contributor",
        action="submission.create_document",
        target_type="submission",
        target_id=record["id"],
        details={"filename": safe_filename, "mime_type": mime_type},
    )
    return record


def sanitize_filename(filename: str) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9._-]+", "_", filename.strip())
    return cleaned or "upload.bin"


def get_portal_home(kb_path: str | Path, *, store: PlatformStore) -> dict[str, Any]:
    data = inventory(kb_path)
    report = audit_kb(kb_path)
    docs = data["docs"]
    recent = sorted(
        (
            {
                "name": name,
                "entry_type": doc["entry_type"],
                "status": doc["status"],
                "updated_at": Path(doc["path"]).stat().st_mtime if doc["path"] else 0,
            }
            for name, doc in docs.items()
            if doc["path"]
        ),
        key=lambda item: item["updated_at"],
        reverse=True,
    )[:8]
    for item in recent:
        item["updated_at"] = int(item["updated_at"])
    hottest = sorted(
        (
            {
                "name": name,
                "entry_type": doc["entry_type"],
                "status": doc["status"],
                "inbound_count": doc["inbound_count"],
                "summary": doc["summary"],
            }
            for name, doc in docs.items()
            if doc["kind"] == "formal"
        ),
        key=lambda item: (item["inbound_count"], len(item["summary"])),
        reverse=True,
    )[:8]
    return {
        "counts": {
            "formal_entries": len(data["entries"]),
            "placeholders": len(data["placeholders"]),
            "indexes": len(data["indexes"]),
            "pending_submissions": store.submission_status_counts().get("pending", 0),
            "health_issues": len(build_health_issue_queue(kb_path)),
        },
        "recent_updates": recent,
        "popular_entries": hottest,
        "health_summary": summarize_health_report(report),
    }


def search_kb(kb_path: str | Path, query: str, *, limit: int = 20) -> list[dict[str, Any]]:
    raw_query = query.strip()
    if not raw_query:
        return []
    data = inventory(kb_path)
    terms = [term for term in re.split(r"\s+", raw_query) if term]
    results: list[dict[str, Any]] = []
    for name, doc in data["docs"].items():
        score = 0
        haystacks = {
            "title": doc["title"],
            "aliases": " ".join(doc["aliases"]),
            "summary": doc["summary"],
            "body": doc["body"],
        }
        for term in terms:
            lowered = term.casefold()
            if lowered in haystacks["title"].casefold():
                score += 10
            if lowered in haystacks["aliases"].casefold():
                score += 8
            if lowered in haystacks["summary"].casefold():
                score += 6
            if lowered in haystacks["body"].casefold():
                score += 3
        if score <= 0:
            continue
        score += min(doc["inbound_count"], 5)
        if doc["kind"] == "formal":
            score += 2
        snippet = build_search_snippet(doc["body"], terms)
        results.append(
            {
                "name": name,
                "kind": doc["kind"],
                "entry_type": doc["entry_type"],
                "status": doc["status"],
                "summary": doc["summary"],
                "snippet": snippet,
                "score": score,
            }
        )
    return sorted(results, key=lambda item: (-item["score"], item["name"]))[:limit]


def build_search_snippet(body: str, terms: list[str]) -> str:
    compact = re.sub(r"\s+", " ", body).strip()
    if not compact:
        return ""
    lowered = compact.casefold()
    for term in terms:
        idx = lowered.find(term.casefold())
        if idx != -1:
            start = max(0, idx - 80)
            end = min(len(compact), idx + 140)
            snippet = compact[start:end].strip()
            if start > 0:
                snippet = "..." + snippet
            if end < len(compact):
                snippet = snippet + "..."
            return snippet
    return compact[:180] + ("..." if len(compact) > 180 else "")


def get_entry_detail(kb_path: str | Path, name: str) -> dict[str, Any]:
    data = inventory(kb_path)
    doc = data["docs"].get(name) or data["index_docs"].get(name)
    if doc is None:
        raise FileNotFoundError(name)
    path = resolve_kb_document_path(kb_path, name)
    content = path.read_text(encoding="utf-8") if path else ""
    validation = None
    if path and (path.parent.name == "indexes" or path.name == index_config()["root_file"]):
        validation = validate_index(path)
    elif path:
        kind = "placeholder" if "/placeholders/" in str(path) else "formal"
        validation = validate_entry(path=path, kind=kind)
    return {
        "name": name,
        "path": str(path) if path else None,
        "content": content,
        "content_hash": content_hash(content),
        "metadata": doc,
        "validation": validation,
    }


def graph_payload(kb_path: str | Path) -> dict[str, Any]:
    data = inventory(kb_path)
    nodes: list[dict[str, Any]] = []
    edges: list[dict[str, Any]] = []
    for name, doc in data["docs"].items():
        nodes.append(
            {
                "id": name,
                "label": name,
                "kind": doc["kind"],
                "entry_type": doc["entry_type"],
                "status": doc["status"],
                "inbound_count": doc["inbound_count"],
            }
        )
        for target in doc["graph_links"]:
            edges.append({"source": name, "target": target, "kind": "related"})
    for name, doc in data["index_docs"].items():
        nodes.append(
            {
                "id": name,
                "label": name,
                "kind": "index",
                "entry_type": "index",
                "status": "n/a",
                "inbound_count": 0,
            }
        )
        for target in doc["links"]:
            edges.append({"source": name, "target": target, "kind": "index_link"})
    return {"nodes": nodes, "edges": edges}


def summarize_health_report(report: dict[str, Any]) -> dict[str, Any]:
    return {
        "formal_entry_count": report["formal_entry_count"],
        "placeholder_count": report["placeholder_count"],
        "hard_fail_entry_count": report["hard_fail_entry_count"],
        "dangling_link_count": report["dangling_link_count"],
        "orphan_entry_count": report["orphan_entry_count"],
        "promotable_placeholder_count": report["promotable_placeholder_count"],
        "canonical_gap_count": report["canonical_gap_count"],
        "invalid_index_count": report["invalid_index_count"],
    }


def build_health_issue_queue(kb_path: str | Path) -> list[dict[str, Any]]:
    report = audit_kb(kb_path)
    data = inventory(kb_path)
    issues: list[dict[str, Any]] = []

    for item in report["entry_validation"]:
        if item["hard_failures"]:
            issues.append(
                issue(
                    issue_type="hard_failure",
                    severity="blocking",
                    target=item["name"],
                    summary="; ".join(item["hard_failures"]),
                    suggested_action="edit_entry",
                    evidence={"hard_failures": item["hard_failures"]},
                )
            )

    for link in report["dangling_links"]:
        issues.append(
            issue(
                issue_type="dangling_link",
                severity="high",
                target=link["entry"],
                summary=f"链接 {link['target']} 没有目标",
                suggested_action="run_tidy",
                evidence=link,
            )
        )

    for orphan in report["orphan_entries"]:
        issues.append(
            issue(
                issue_type="orphan_entry",
                severity="medium",
                target=orphan,
                summary="正式条目缺少足够连接",
                suggested_action="run_tidy",
                evidence={"entry": orphan},
            )
        )

    for item in report["promotable_placeholders"]:
        issues.append(
            issue(
                issue_type="promotable_placeholder",
                severity="medium",
                target=item["name"],
                summary=f"placeholder 被引用 {item['ref_count']} 次，值得提升",
                suggested_action="promote_placeholder",
                evidence=item,
            )
        )

    for item in report["canonical_gaps"]:
        issues.append(
            issue(
                issue_type="canonical_gap",
                severity="high",
                target=item["name"],
                summary="正式条目依赖 placeholder 作为关键概念",
                suggested_action="run_tidy",
                evidence=item,
            )
        )

    for item in report["provenance_contamination"]:
        issues.append(
            issue(
                issue_type="provenance_contamination",
                severity="medium",
                target=item["name"],
                summary="来源型链接污染了知识图谱",
                suggested_action="edit_entry",
                evidence=item,
            )
        )

    for item in report["index_validation"]:
        if item["hard_failures"]:
            issues.append(
                issue(
                    issue_type="invalid_index",
                    severity="blocking",
                    target=item["name"],
                    summary="; ".join(item["hard_failures"]),
                    suggested_action="edit_entry",
                    evidence=item,
                )
            )

    for item in report["overloaded_indexes"]:
        issues.append(
            issue(
                issue_type="overloaded_index",
                severity="medium",
                target=item["name"],
                summary="索引规模超过阈值",
                suggested_action="run_tidy",
                evidence=item,
            )
        )

    for item in report["unknown_index_links"]:
        issues.append(
            issue(
                issue_type="unknown_index_link",
                severity="high",
                target=item["index"],
                summary=f"索引链接 {item['link']} 不存在或不应作为入口",
                suggested_action="edit_entry",
                evidence=item,
            )
        )

    for name in report["uncovered_formal_entries"]:
        issues.append(
            issue(
                issue_type="uncovered_formal_entry",
                severity="medium",
                target=name,
                summary="正式条目未被任何索引覆盖",
                suggested_action="run_tidy",
                evidence={"entry": name},
            )
        )

    for name, doc in data["docs"].items():
        if doc["kind"] != "formal":
            continue
        if doc["status"] in {"inferred", "disputed"}:
            issues.append(
                issue(
                    issue_type="low_confidence_entry",
                    severity="low" if doc["status"] == "inferred" else "medium",
                    target=name,
                    summary=f"条目状态为 {doc['status']}，需要人工关注",
                    suggested_action="review_entry",
                    evidence={"status": doc["status"]},
                )
            )

    return sorted(issues, key=lambda item: (severity_rank(item["severity"]), item["target"]))


def issue(
    *,
    issue_type: str,
    severity: str,
    target: str,
    summary: str,
    suggested_action: str,
    evidence: dict[str, Any],
) -> dict[str, Any]:
    return {
        "type": issue_type,
        "severity": severity,
        "target": target,
        "summary": summary,
        "suggested_action": suggested_action,
        "evidence": evidence,
    }


def severity_rank(severity: str) -> int:
    return {"blocking": 0, "high": 1, "medium": 2, "low": 3}.get(severity, 4)


def get_health_payload(kb_path: str | Path) -> dict[str, Any]:
    report = audit_kb(kb_path)
    issues = build_health_issue_queue(kb_path)
    counts: dict[str, int] = {}
    for item in issues:
        counts[item["severity"]] = counts.get(item["severity"], 0) + 1
    return {
        "summary": summarize_health_report(report),
        "severity_counts": counts,
        "issues": issues,
    }


def content_hash(content: str) -> str:
    return hashlib.sha256(content.encode("utf-8")).hexdigest()


def list_reviews_with_jobs(
    *,
    store: PlatformStore,
    decision: str | None = None,
) -> list[dict[str, Any]]:
    reviews = store.list_reviews(decision=decision, limit=200)
    enriched: list[dict[str, Any]] = []
    for review in reviews:
        payload = dict(review)
        payload["job"] = store.get_job(review["job_id"])
        if review.get("submission_id"):
            payload["submission"] = store.get_submission(review["submission_id"])
        enriched.append(payload)
    return enriched


def determine_target_path(
    kb_path: str | Path,
    *,
    name: str,
    content: str,
    relative_path: str | None = None,
) -> Path:
    root = Path(kb_path).resolve()
    if relative_path:
        target = (root / relative_path).resolve()
        if root not in target.parents and target != root:
            raise ValueError("relative path escapes knowledge base root")
        return target

    existing = resolve_kb_document_path(root, name)
    if existing is not None:
        return existing.resolve()

    frontmatter, _ = split_frontmatter(content)
    item_type = str(frontmatter.get("type", "")).strip()
    if item_type == "placeholder":
        return root / "placeholders" / f"{name}.md"
    if item_type in {"concept", "lesson"}:
        return root / "entries" / f"{name}.md"
    if str(frontmatter.get("kind", "")).strip() == "index" or name.startswith("index."):
        if name == "index.root":
            return root / index_config()["root_file"]
        return root / "indexes" / f"{name}.md"
    return root / "entries" / f"{name}.md"


def validate_target_content(path: Path, content: str) -> dict[str, Any] | None:
    if path.name == index_config()["root_file"] or path.parent.name == "indexes":
        temp_path = _write_temp_for_validation(path.name, content)
        try:
            return validate_index(temp_path)
        finally:
            temp_path.unlink(missing_ok=True)
    kind = "placeholder" if path.parent.name == "placeholders" else "formal"
    return validate_entry(text=content, name=path.stem, kind=kind)


def _write_temp_for_validation(filename: str, content: str) -> Path:
    tmp = tempfile.NamedTemporaryFile(
        prefix="sediment-validate-",
        suffix=f"-{filename}",
        delete=False,
    )
    tmp.write(content.encode("utf-8"))
    tmp.flush()
    tmp.close()
    return Path(tmp.name)


def build_diff(path_label: str, old_content: str, new_content: str) -> str:
    diff = difflib.unified_diff(
        old_content.splitlines(keepends=True),
        new_content.splitlines(keepends=True),
        fromfile=f"a/{path_label}",
        tofile=f"b/{path_label}",
    )
    return "".join(diff)


def apply_operations(
    kb_path: str | Path,
    operations: list[dict[str, Any]],
    *,
    actor_name: str,
    actor_role: str,
    store: PlatformStore,
) -> dict[str, Any]:
    root = Path(kb_path).resolve()
    prepared: list[dict[str, Any]] = []
    seen_paths: set[str] = set()
    for operation in operations:
        name = operation.get("name") or Path(operation.get("relative_path", "")).stem
        target_path = determine_target_path(
            root,
            name=name,
            content=operation.get("content", ""),
            relative_path=operation.get("relative_path"),
        )
        old_content = target_path.read_text(encoding="utf-8") if target_path.exists() else ""
        current_hash = content_hash(old_content)
        base_hash = operation.get("base_hash")
        if base_hash and target_path.exists() and current_hash != base_hash:
            raise RuntimeError(
                f"{target_path.name} changed since review was prepared; refresh the diff first."
            )
        rel_path = str(target_path.relative_to(root))
        if rel_path in seen_paths:
            raise ValueError(f"duplicate operation target detected: {rel_path}")
        seen_paths.add(rel_path)

        change_type = operation.get("change_type", "update")
        if change_type == "delete":
            prepared.append(
                {
                    "name": name,
                    "target_path": target_path,
                    "relative_path": rel_path,
                    "change_type": "delete",
                    "old_content": old_content,
                }
            )
            continue

        new_content = operation["content"].strip() + "\n"
        validation = validate_target_content(target_path, new_content)
        if validation and validation.get("hard_failures"):
            raise ValueError("; ".join(validation["hard_failures"]))
        prepared.append(
            {
                "name": name,
                "target_path": target_path,
                "relative_path": rel_path,
                "change_type": change_type,
                "old_content": old_content,
                "new_content": new_content,
            }
        )

    applied: list[dict[str, Any]] = []
    for item in prepared:
        target_path = item["target_path"]
        rel_path = item["relative_path"]
        if item["change_type"] == "delete":
            if target_path.exists():
                target_path.unlink()
            applied.append(
                {
                    "name": item["name"],
                    "relative_path": rel_path,
                    "change_type": "delete",
                    "diff": build_diff(rel_path, item["old_content"], ""),
                }
            )
            continue

        new_content = item["new_content"]
        target_path.parent.mkdir(parents=True, exist_ok=True)
        target_path.write_text(new_content, encoding="utf-8")
        applied.append(
            {
                "name": item["name"],
                "relative_path": rel_path,
                "change_type": item["change_type"],
                "diff": build_diff(rel_path, item["old_content"], new_content),
                "new_hash": content_hash(new_content),
            }
        )

    store.add_audit_log(
        actor_name=actor_name,
        actor_role=actor_role,
        action="kb.apply_operations",
        target_type="knowledge_base",
        target_id=str(root),
        details={"operations": applied},
    )
    return {"operations": applied, "health": summarize_health_report(audit_kb(root))}


def save_entry(
    kb_path: str | Path,
    *,
    name: str,
    content: str,
    expected_hash: str | None,
    actor_name: str,
    store: PlatformStore,
) -> dict[str, Any]:
    root = Path(kb_path).resolve()
    target = determine_target_path(kb_path, name=name, content=content)
    old_content = target.read_text(encoding="utf-8") if target.exists() else ""
    if expected_hash and old_content and content_hash(old_content) != expected_hash:
        raise RuntimeError("entry changed since it was loaded; refresh before saving")
    result = apply_operations(
        root,
        [
            {
                "name": name,
                "relative_path": str(target.relative_to(root)),
                "change_type": "create" if not old_content else "update",
                "content": content,
                "base_hash": content_hash(old_content) if old_content else None,
            }
        ],
        actor_name=actor_name,
        actor_role="committer",
        store=store,
    )
    detail = get_entry_detail(root, name)
    detail["write_result"] = result
    return detail


def stage_workspace_copy(kb_path: str | Path, workspaces_dir: str | Path, job_id: str) -> Path:
    workspaces_root = Path(workspaces_dir)
    workspaces_root.mkdir(parents=True, exist_ok=True)
    workspace = workspaces_root / job_id
    if workspace.exists():
        shutil.rmtree(workspace)
    shutil.copytree(kb_path, workspace / "knowledge-base")
    return workspace
